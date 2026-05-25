import base64
import uuid

import pytest
from pptx import Presentation

from slidegen.services.presentation.page_classifier import PageType, PageTypeClassifier

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _blank_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


# --- Task 1: Shape summary extraction ---


def test_summarize_slide_extracts_text_picture_and_placeholder(tmp_path):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly Strategy"
    title_slide.placeholders[1].text = "2026 Planning"

    image_path = tmp_path / "pixel.png"
    image_path.write_bytes(PNG_BYTES)
    title_slide.shapes.add_picture(str(image_path), 100000, 100000, 200000, 200000)

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(title_slide)

    assert any(summary.is_placeholder and "Quarterly Strategy" in summary.text for summary in summaries)
    assert any(summary.is_picture for summary in summaries)
    assert all(summary.text == summary.text[: classifier.max_text_chars] for summary in summaries)
    assert all(summary.width >= 0 and summary.height >= 0 for summary in summaries)


# --- Task 2: Local rule classification ---


def test_rule_classifies_obvious_end_page():
    _, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 4000000, 600000).text = "Thank you"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    result = classifier.classify_by_rules(page_index=4, slide_count=5, summaries=summaries)

    assert result is not None
    assert result.page_type == PageType.END
    assert result.method == "rule"


def test_rule_classifies_empty_slide_as_unknown():
    _, slide = _blank_slide()

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    result = classifier.classify_by_rules(page_index=2, slide_count=5, summaries=summaries)

    assert result is not None
    assert result.page_type == PageType.UNKNOWN
    assert result.method == "rule"


def test_rule_classifies_simple_first_slide_as_cover():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Strategy"
    slide.placeholders[1].text = "Planning deck"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    result = classifier.classify_by_rules(page_index=0, slide_count=5, summaries=summaries)

    assert result is not None
    assert result.page_type == PageType.COVER
    assert result.method == "rule"


# --- Task 3: Agent classification and fallback ---


class FakeAgent:
    def __init__(self, content=None, exc: Exception | None = None):
        self.content = content
        self.exc = exc
        self.prompts = []

    async def arun(self, prompt):
        self.prompts.append(prompt)
        if self.exc is not None:
            raise self.exc
        return type("RunOutput", (), {"content": self.content})()


@pytest.mark.anyio
async def test_agent_classification_uses_structured_output():
    _, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Market drivers"
    slide.shapes.add_textbox(100000, 700000, 5000000, 1200000).text = "Customers need faster onboarding."

    fake_agent = FakeAgent(
        content={
            "page_type": "chapter_content",
            "confidence": 0.88,
            "reason": "The slide contains explanatory content.",
        }
    )
    classifier = PageTypeClassifier(agent_factory=lambda model: fake_agent)

    summaries = classifier.summarize_slide(slide)
    result = await classifier.classify_with_agent(
        model=object(),
        page_index=2,
        slide_count=5,
        summaries=summaries,
    )

    assert result.page_type == PageType.CHAPTER_CONTENT
    assert result.confidence == 0.88
    assert result.method == "llm"
    assert fake_agent.prompts


@pytest.mark.anyio
async def test_agent_low_confidence_falls_back_to_unknown():
    _, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Ambiguous"

    fake_agent = FakeAgent(content={"page_type": "catalog", "confidence": 0.2, "reason": "Unclear"})
    classifier = PageTypeClassifier(agent_factory=lambda model: fake_agent)

    result = await classifier.classify_with_agent(
        model=object(),
        page_index=1,
        slide_count=5,
        summaries=classifier.summarize_slide(slide),
    )

    assert result.page_type == PageType.UNKNOWN
    assert result.method == "fallback"
    assert "below threshold" in result.reason


@pytest.mark.anyio
async def test_agent_invalid_output_falls_back_to_unknown():
    _, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Ambiguous"

    fake_agent = FakeAgent(content={"page_type": "bad_type", "confidence": 0.8, "reason": "Invalid enum"})
    classifier = PageTypeClassifier(agent_factory=lambda model: fake_agent)

    result = await classifier.classify_with_agent(
        model=object(),
        page_index=1,
        slide_count=5,
        summaries=classifier.summarize_slide(slide),
    )

    assert result.page_type == PageType.UNKNOWN
    assert result.method == "fallback"


@pytest.mark.anyio
async def test_agent_exception_falls_back_to_unknown():
    _, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Ambiguous"

    fake_agent = FakeAgent(exc=RuntimeError("provider failed"))
    classifier = PageTypeClassifier(agent_factory=lambda model: fake_agent)

    result = await classifier.classify_with_agent(
        model=object(),
        page_index=1,
        slide_count=5,
        summaries=classifier.summarize_slide(slide),
    )

    assert result.page_type == PageType.UNKNOWN
    assert result.method == "fallback"
    assert "provider failed" in result.reason


# --- Task 4: Batch classify_pages ---


@pytest.mark.anyio
async def test_classify_pages_returns_one_result_per_slide_without_real_llm(tmp_path):
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Launch Plan"
    cover.placeholders[1].text = "May 2026"

    content = prs.slides.add_slide(prs.slide_layouts[6])
    content.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Customer Signals"
    content.shapes.add_textbox(100000, 700000, 5000000, 1200000).text = "Customers ask for faster onboarding."

    end = prs.slides.add_slide(prs.slide_layouts[6])
    end.shapes.add_textbox(100000, 100000, 4000000, 600000).text = "Thank you"

    pptx_path = tmp_path / "deck.pptx"
    prs.save(pptx_path)

    fake_agent = FakeAgent(
        content={
            "page_type": "chapter_content",
            "confidence": 0.91,
            "reason": "Actual content slide.",
        }
    )
    classifier = PageTypeClassifier(agent_factory=lambda model: fake_agent)

    results = await classifier.classify_pages(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),
    )

    assert [result.page_type for result in results] == [
        PageType.COVER,
        PageType.CHAPTER_CONTENT,
        PageType.END,
    ]
    assert [result.method for result in results] == ["rule", "llm", "rule"]
    assert len(fake_agent.prompts) == 1


@pytest.mark.anyio
async def test_classify_pages_raises_for_missing_pptx(tmp_path):
    classifier = PageTypeClassifier(agent_factory=lambda model: FakeAgent())

    missing_path = tmp_path / "missing.pptx"

    try:
        await classifier.classify_pages(
            pptx_path=str(missing_path),
            user_id=uuid.uuid4(),
            model=object(),
        )
    except FileNotFoundError as exc:
        assert str(missing_path) in str(exc)
    else:
        raise AssertionError("Expected FileNotFoundError")
