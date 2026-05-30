"""Tests for ContentStyleImporter (Phase 2)."""

import base64
import hashlib
import json
import re
import uuid

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

from slidegen.services.presentation.component_importer import (
    _SHORT_TEXT_THRESHOLD,
    AgentShapeAssignment,
    AgentShapeRoleOutput,
    ContentStyleImporter,
    ImportSlideStatus,
    LocalShapeRoleClassifier,
    ShapeAssignment,
    ShapeRoleAgent,
    StyleBuilder,
    atomic_save_json,
    check_fingerprint,
    compute_pptx_fingerprint,
    identify_layout_type,
    store_fingerprint,
    validate_compatibility,
)
from slidegen.services.presentation.components import (
    ChapterLayout,
    ComponentContentType,
    ComponentsManager,
    CShape,
    Location,
    Style,
)
from slidegen.services.presentation.page_classifier import (
    PageClassification,
    PageType,
    PageTypeClassifier,
)

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _blank_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _two_point_slide_with_font():
    """Create a synthetic two-point slide where titles have font_size for local rules."""
    prs, slide = _blank_slide()

    # Group 0
    number0 = slide.shapes.add_textbox(Emu(200000), Emu(400000), Emu(600000), Emu(300000))
    number0.text = "01"

    title0 = slide.shapes.add_textbox(Emu(900000), Emu(400000), Emu(3000000), Emu(400000))
    tf0 = title0.text_frame
    tf0.paragraphs[0].text = "Customer Signals"
    tf0.paragraphs[0].font.size = Pt(20)
    tf0.paragraphs[0].font.bold = True

    content0 = slide.shapes.add_textbox(Emu(900000), Emu(900000), Emu(5000000), Emu(1200000))
    content0.text = "Customers ask for faster onboarding and better tools."

    # Group 1 (positioned lower with large Y gap)
    number1 = slide.shapes.add_textbox(Emu(200000), Emu(2800000), Emu(600000), Emu(300000))
    number1.text = "02"

    title1 = slide.shapes.add_textbox(Emu(900000), Emu(2800000), Emu(3000000), Emu(400000))
    tf1 = title1.text_frame
    tf1.paragraphs[0].text = "Market Drivers"
    tf1.paragraphs[0].font.size = Pt(20)
    tf1.paragraphs[0].font.bold = True

    content1 = slide.shapes.add_textbox(Emu(900000), Emu(3300000), Emu(5000000), Emu(1200000))
    content1.text = "The market is shifting toward AI-first products and solutions."

    return prs, slide


def _two_point_slide():
    """Create a synthetic two-point chapter content slide."""
    prs, slide = _blank_slide()

    # Group 0: number + title + content
    slide.shapes.add_textbox(Emu(200000), Emu(400000), Emu(600000), Emu(300000)).text = "01"
    slide.shapes.add_textbox(Emu(900000), Emu(400000), Emu(3000000), Emu(400000)).text = "Customer Signals"
    slide.shapes.add_textbox(
        Emu(900000), Emu(900000), Emu(5000000), Emu(1200000)
    ).text = "Customers ask for faster onboarding."

    # Group 1: number + title + content (positioned lower)
    slide.shapes.add_textbox(Emu(200000), Emu(2800000), Emu(600000), Emu(300000)).text = "02"
    slide.shapes.add_textbox(Emu(900000), Emu(2800000), Emu(3000000), Emu(400000)).text = "Market Drivers"
    slide.shapes.add_textbox(
        Emu(900000), Emu(3300000), Emu(5000000), Emu(1200000)
    ).text = "The market is shifting toward AI-first products."

    return prs, slide


# --- Step 1: Core types ---


def test_import_slide_status_enum_values():
    assert ImportSlideStatus.IMPORTED.value == "imported"
    assert ImportSlideStatus.SKIPPED.value == "skipped"
    assert ImportSlideStatus.DRY_RUN.value == "dry_run"
    assert ImportSlideStatus.FAILED.value == "failed"


def test_agent_shape_assignment_pydantic_validation():
    a = AgentShapeAssignment(
        shape_id=1,
        content_type="title",
        group_index=0,
        confidence=0.9,
        reason="Title text",
    )
    assert a.shape_id == 1
    assert a.content_type == "title"
    assert a.group_index == 0


def test_agent_shape_assignment_rejects_invalid_group_index():
    with pytest.raises(ValueError):
        AgentShapeAssignment(
            shape_id=1,
            content_type="title",
            group_index=5,
            confidence=0.9,
            reason="Bad group",
        )


def test_agent_shape_role_output_validation():
    output = AgentShapeRoleOutput(
        point_count=2,
        assignments=[
            AgentShapeAssignment(
                shape_id=1,
                content_type="title",
                group_index=0,
                confidence=0.85,
                reason="Section heading",
            ),
        ],
        confidence=0.85,
        reason="Two-point layout detected.",
    )
    assert output.point_count == 2


# --- Step 2: Local classifier ---


def test_local_classifier_skips_placeholders():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Main Title"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)

    placeholder_assignments = [a for a in assignments if not a.include and a.reason.startswith("Placeholder")]
    assert len(placeholder_assignments) >= 1


def test_local_classifier_classifies_pictures():
    prs, slide = _blank_slide()

    image_path = "/tmp/test_importer_pixel.png"
    with open(image_path, "wb") as f:
        f.write(PNG_BYTES)

    # Small picture (icon) — 500000 x 500000 EMU < threshold
    slide.shapes.add_picture(image_path, 100000, 100000, 500000, 500000)

    # Large picture — 5000000 x 5000000 EMU > threshold
    slide.shapes.add_picture(image_path, 2000000, 2000000, 5000000, 5000000)

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)
    picture_types = [
        a.content_type
        for a in assignments
        if a.include and a.content_type in (ComponentContentType.PICTURE, ComponentContentType.ICON)
    ]
    assert ComponentContentType.ICON in picture_types
    assert ComponentContentType.PICTURE in picture_types


def test_local_classifier_classifies_medium_wide_picture_as_picture():
    prs, slide = _blank_slide()

    image_path = "/tmp/test_importer_pixel.png"
    with open(image_path, "wb") as f:
        f.write(PNG_BYTES)

    slide.shapes.add_picture(image_path, 100000, 100000, 2880449, 1520002)

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)

    assert assignments[0].content_type == ComponentContentType.PICTURE


def test_local_classifier_classifies_text_shapes():
    prs, slide = _blank_slide()
    slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(3000000), Emu(300000)).text = "01"
    # Title with bold
    title_box = slide.shapes.add_textbox(Emu(100000), Emu(500000), Emu(3000000), Emu(400000))
    title_box.text_frame.paragraphs[0].text = "Section title"
    title_box.text_frame.paragraphs[0].font.bold = True
    # Content with longer text
    slide.shapes.add_textbox(
        Emu(100000), Emu(1000000), Emu(5000000), Emu(1200000)
    ).text = "This is a longer body paragraph that describes the content in detail."

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)
    content_types = [a.content_type for a in assignments if a.include]
    content_type_values = [ct.value if isinstance(ct, ComponentContentType) else ct for ct in content_types]
    assert "number" in content_type_values
    assert "title" in content_type_values
    assert "content" in content_type_values


def test_local_classifier_returns_skip_for_ambiguous_text():
    prs, slide = _blank_slide()
    # Short text that's not clearly number or title — ambiguous
    slide.shapes.add_textbox(100000, 100000, 1000000, 300000).text = "Note"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)
    skip_assignments = [a for a in assignments if a.content_type == "skip" and not a.include]
    assert len(skip_assignments) >= 1


def test_local_classifier_assigns_decoration_for_no_text_shapes():
    prs, slide = _blank_slide()
    # Empty text frame → decoration
    slide.shapes.add_textbox(100000, 100000, 500000, 500000).text = ""

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.classify(summaries)
    decoration_assignments = [a for a in assignments if a.content_type == ComponentContentType.DECORATION]
    assert len(decoration_assignments) >= 1


def test_local_classifier_groups_horizontal_columns_by_visual_group():
    prs, slide = _blank_slide()

    left_title = slide.shapes.add_textbox(Emu(400000), Emu(500000), Emu(2500000), Emu(400000))
    left_title.text_frame.paragraphs[0].text = "Customer Signals"
    left_title.text_frame.paragraphs[0].font.size = Pt(20)
    left_title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_textbox(
        Emu(400000), Emu(1050000), Emu(3000000), Emu(800000)
    ).text = "Customers ask for faster onboarding and better tools."

    right_title = slide.shapes.add_textbox(Emu(4500000), Emu(500000), Emu(2500000), Emu(400000))
    right_title.text_frame.paragraphs[0].text = "Market Drivers"
    right_title.text_frame.paragraphs[0].font.size = Pt(20)
    right_title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_textbox(
        Emu(4500000), Emu(1050000), Emu(3000000), Emu(800000)
    ).text = "The market is shifting toward AI-first products and solutions."

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)
    local = LocalShapeRoleClassifier()

    assignments = local.assign_group_indices(local.classify(summaries), summaries)

    assert identify_layout_type(assignments) == ChapterLayout.TWO_POINTS
    valid, reason = validate_compatibility(assignments, ChapterLayout.TWO_POINTS)
    assert valid, reason


# --- Step 4: Layout identification ---


def test_identify_layout_type_one_point():
    assignments = [
        ShapeAssignment(
            shape_id=1, content_type=ComponentContentType.TITLE, group_index=0, include=True, reason="", confidence=0.9
        ),
        ShapeAssignment(
            shape_id=2,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="",
            confidence=0.9,
        ),
    ]
    layout = identify_layout_type(assignments)
    assert layout == ChapterLayout.ONE_POINT


def test_identify_layout_type_two_points():
    assignments = [
        ShapeAssignment(
            shape_id=1, content_type=ComponentContentType.TITLE, group_index=0, include=True, reason="", confidence=0.9
        ),
        ShapeAssignment(
            shape_id=2,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="",
            confidence=0.9,
        ),
        ShapeAssignment(
            shape_id=3, content_type=ComponentContentType.TITLE, group_index=1, include=True, reason="", confidence=0.9
        ),
        ShapeAssignment(
            shape_id=4,
            content_type=ComponentContentType.CONTENT,
            group_index=1,
            include=True,
            reason="",
            confidence=0.9,
        ),
    ]
    layout = identify_layout_type(assignments)
    assert layout == ChapterLayout.TWO_POINTS


def test_identify_layout_type_returns_none_for_no_content():
    assignments = [
        ShapeAssignment(
            shape_id=1,
            content_type=ComponentContentType.DECORATION,
            group_index=None,
            include=True,
            reason="",
            confidence=0.9,
        ),
    ]
    layout = identify_layout_type(assignments)
    assert layout is None


# --- Step 4: Compatibility validation ---


def test_validate_compatibility_passes_for_valid_assignments():
    assignments = [
        ShapeAssignment(
            shape_id=1, content_type=ComponentContentType.TITLE, group_index=0, include=True, reason="", confidence=0.9
        ),
        ShapeAssignment(
            shape_id=2,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="",
            confidence=0.9,
        ),
    ]
    layout = ChapterLayout.ONE_POINT
    valid, reason = validate_compatibility(assignments, layout)
    assert valid
    assert reason == ""


def test_validate_compatibility_fails_for_mismatched_title_count():
    assignments = [
        ShapeAssignment(
            shape_id=1, content_type=ComponentContentType.TITLE, group_index=0, include=True, reason="", confidence=0.9
        ),
        ShapeAssignment(
            shape_id=2,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="",
            confidence=0.9,
        ),
        ShapeAssignment(
            shape_id=3,
            content_type=ComponentContentType.CONTENT,
            group_index=1,
            include=True,
            reason="",
            confidence=0.9,
        ),
    ]
    layout = ChapterLayout.TWO_POINTS
    valid, reason = validate_compatibility(assignments, layout)
    assert not valid
    assert "TITLE" in reason


def test_validate_compatibility_fails_for_missing_content():
    assignments = [
        ShapeAssignment(
            shape_id=1, content_type=ComponentContentType.TITLE, group_index=0, include=True, reason="", confidence=0.9
        ),
    ]
    layout = ChapterLayout.ONE_POINT
    valid, reason = validate_compatibility(assignments, layout)
    assert not valid
    assert "CONTENT" in reason


# --- Step 6: Atomic write ---


def test_atomic_save_json_creates_file(tmp_path):
    target = tmp_path / "test.json"
    data = {"one_point": {"style0": {"xml": "test", "zorder": 0, "content_type": "content", "location": []}}}

    atomic_save_json(data, target, backup=False)

    assert target.exists()
    with open(target) as f:
        loaded = json.load(f)
    assert loaded == data


def test_atomic_save_json_creates_backup(tmp_path):
    target = tmp_path / "test.json"
    # Pre-existing file
    target.write_text(json.dumps({"old": True}))

    data = {"new": True}
    atomic_save_json(data, target, backup=True)

    bak_path = target.with_suffix(".json.bak")
    assert bak_path.exists()
    with open(bak_path) as f:
        old_data = json.load(f)
    assert old_data == {"old": True}

    with open(target) as f:
        new_data = json.load(f)
    assert new_data == data


# --- Step 6: Fingerprint ---


def test_compute_pptx_fingerprint_deterministic(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    fp1 = compute_pptx_fingerprint(pptx_path)
    fp2 = compute_pptx_fingerprint(pptx_path)
    assert fp1 == fp2
    assert len(fp1) == 64  # SHA-256 hex length


def test_check_fingerprint_detects_duplicate(tmp_path):
    json_path = tmp_path / "shapes.json"
    json_path.write_text("{}")

    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(pptx_path))

    # First import
    store_fingerprint(pptx_path, json_path)

    # Check should detect duplicate
    assert check_fingerprint(pptx_path, json_path)


def test_check_fingerprint_allows_new_file(tmp_path):
    json_path = tmp_path / "shapes.json"
    json_path.write_text("{}")

    pptx1 = tmp_path / "test1.pptx"
    prs1 = Presentation()
    prs1.slides.add_slide(prs1.slide_layouts[0])
    prs1.save(str(pptx1))

    pptx2 = tmp_path / "test2.pptx"
    prs2 = Presentation()
    prs2.slides.add_slide(prs2.slide_layouts[6])
    prs2.save(str(pptx2))

    store_fingerprint(pptx1, json_path)
    assert not check_fingerprint(pptx2, json_path)


# --- Step 5: StyleBuilder ---


def test_style_builder_produces_valid_style_from_assignments():
    prs, slide = _two_point_slide()
    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)

    # Manually construct assignments that represent a two-point layout
    # because local rules can't determine TITLE without font/bold info
    [s.shape_id for s in summaries if s.text in ("01", "02")]
    [s.shape_id for s in summaries if len(s.text) > _SHORT_TEXT_THRESHOLD]
    [s.shape_id for s in summaries if s.text in ("Customer Signals", "Market Drivers")]

    # Create manual assignments
    manual_assignments: list[ShapeAssignment] = []

    # Numbers: group 0 and 1
    for s in summaries:
        if s.text == "01":
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.NUMBER,
                    group_index=0,
                    include=True,
                    reason="Number",
                    confidence=0.9,
                )
            )
        elif s.text == "02":
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.NUMBER,
                    group_index=1,
                    include=True,
                    reason="Number",
                    confidence=0.9,
                )
            )

    # Titles
    for s in summaries:
        if s.text == "Customer Signals":
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.TITLE,
                    group_index=0,
                    include=True,
                    reason="Title",
                    confidence=0.8,
                )
            )
        elif s.text == "Market Drivers":
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.TITLE,
                    group_index=1,
                    include=True,
                    reason="Title",
                    confidence=0.8,
                )
            )

    # Contents
    for s in summaries:
        if s.text.startswith("Customers"):
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.CONTENT,
                    group_index=0,
                    include=True,
                    reason="Content",
                    confidence=0.8,
                )
            )
        elif s.text.startswith("The market"):
            manual_assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type=ComponentContentType.CONTENT,
                    group_index=1,
                    include=True,
                    reason="Content",
                    confidence=0.8,
                )
            )

    builder = StyleBuilder()
    style = builder.build_style_from_assignments(
        slide=slide,
        _layout_type=ChapterLayout.TWO_POINTS,
        style_name="test_style",
        assignments=manual_assignments,
    )

    assert style.name == "test_style"
    # Should have TITLE and CONTENT shapes
    title_shapes = [s for s in style.shape_list if s.content_type == ComponentContentType.TITLE]
    content_shapes = [s for s in style.shape_list if s.content_type == ComponentContentType.CONTENT]
    assert len(title_shapes) >= 1
    assert len(content_shapes) >= 1

    # TITLE location count should equal point count (2)
    for ts in title_shapes:
        assert len(ts.location) == 2

    # CONTENT location count should equal point count (2)
    for cs in content_shapes:
        assert len(cs.location) == 2


def test_style_builder_excludes_skip_shapes():
    prs, slide = _blank_slide()
    # Add a placeholder-like shape that will be skipped
    slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "Section"

    # Manually create assignments with a skip
    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)

    # Get the text shape's id
    text_shape_id = None
    for s in summaries:
        if s.text == "Section":
            text_shape_id = s.shape_id
            break

    assignments = [
        ShapeAssignment(
            shape_id=text_shape_id,
            content_type="skip",
            group_index=None,
            include=False,
            reason="Skip this",
            confidence=0.0,
        ),
    ]

    builder = StyleBuilder()
    style = builder.build_style_from_assignments(
        slide=slide,
        _layout_type=ChapterLayout.ONE_POINT,
        style_name="skip_test",
        assignments=assignments,
    )

    assert len(style) == 0  # no shapes included


def test_style_builder_replaces_imported_text_in_shape_xml():
    prs, slide = _blank_slide()
    shape = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(4000000), Emu(600000))
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.add_run().text = "Original sensitive sentence"
    paragraph.add_run().text = "Second imported sentence"

    assignments = [
        ShapeAssignment(
            shape_id=shape.shape_id,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="Content",
            confidence=0.9,
        )
    ]

    style = StyleBuilder().build_style_from_assignments(
        slide=slide,
        _layout_type=ChapterLayout.ONE_POINT,
        style_name="text_placeholder_test",
        assignments=assignments,
    )

    content_shape = style.get_shape("content_0")
    assert content_shape is not None
    assert content_shape.xml is not None
    assert "Original sensitive sentence" not in content_shape.xml
    assert "Second imported sentence" not in content_shape.xml

    root = etree.fromstring(content_shape.xml)
    text_values = [elem.text for elem in root.findall(".//a:t", namespaces=root.nsmap)]
    assert text_values == ["Text", "Text"]


def test_style_builder_replaces_chinese_shape_names_in_shape_xml():
    prs, slide = _blank_slide()
    shape = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(4000000), Emu(600000))
    shape.name = "矩形 3"
    shape.text = "Content"

    assignments = [
        ShapeAssignment(
            shape_id=shape.shape_id,
            content_type=ComponentContentType.CONTENT,
            group_index=0,
            include=True,
            reason="Content",
            confidence=0.9,
        )
    ]

    style = StyleBuilder().build_style_from_assignments(
        slide=slide,
        _layout_type=ChapterLayout.ONE_POINT,
        style_name="shape_name_test",
        assignments=assignments,
    )

    content_shape = style.get_shape("content_0")
    assert content_shape is not None
    assert content_shape.xml is not None
    assert 'name="矩形 3"' not in content_shape.xml
    assert 'name="rectangle_3"' in content_shape.xml


# --- Step 7: Full import pipeline ---


class FakeAgentForImport:
    """Fake agno Agent that returns a shape role output."""

    def __init__(self, content=None, exc=None):
        self.content = content
        self.exc = exc
        self.prompts = []

    async def arun(self, prompt):
        self.prompts.append(prompt)
        if self.exc is not None:
            raise self.exc
        return type("RunOutput", (), {"content": self.content})()


class FakePageClassifierAgent:
    """Fake agent for PageTypeClassifier that returns chapter_content."""

    def __init__(self, content=None, exc=None):
        self.content = content or {
            "page_type": "chapter_content",
            "confidence": 0.91,
            "reason": "Content slide detected.",
        }
        self.exc = exc
        self.prompts = []

    async def arun(self, prompt):
        self.prompts.append(prompt)
        if self.exc is not None:
            raise self.exc
        return type("RunOutput", (), {"content": self.content})()


@pytest.mark.anyio
async def test_import_from_pptx_skips_non_chapter_content_slides(tmp_path):
    prs = Presentation()
    # Slide 0: cover
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Launch Plan"

    # Slide 1: chapter_content (two points)
    content = prs.slides.add_slide(prs.slide_layouts[6])
    content.shapes.add_textbox(Emu(200000), Emu(400000), Emu(600000), Emu(300000)).text = "01"
    content.shapes.add_textbox(Emu(900000), Emu(400000), Emu(3000000), Emu(400000)).text = "Signals"
    content.shapes.add_textbox(
        Emu(900000), Emu(900000), Emu(5000000), Emu(1200000)
    ).text = "Longer content text about onboarding."
    content.shapes.add_textbox(Emu(200000), Emu(2800000), Emu(600000), Emu(300000)).text = "02"
    content.shapes.add_textbox(Emu(900000), Emu(2800000), Emu(3000000), Emu(400000)).text = "Drivers"
    content.shapes.add_textbox(
        Emu(900000), Emu(3300000), Emu(5000000), Emu(1200000)
    ).text = "The market is shifting toward AI products."

    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    # Use a ComponentsManager with existing data
    json_path = tmp_path / "shapes.json"
    json_path.write_text(
        json.dumps(
            {
                "one_point": {},
                "two_points": {},
                "three_points": {},
            }
        )
    )
    cm = ComponentsManager(str(json_path))

    fake_page_agent = FakePageClassifierAgent()
    page_classifier = PageTypeClassifier(agent_factory=lambda model: fake_page_agent)

    fake_shape_agent = FakeAgentForImport(
        content={
            "point_count": 2,
            "assignments": [],
            "confidence": 0.9,
            "reason": "Two-point layout.",
        }
    )
    shape_role_agent = ShapeRoleAgent(agent_factory=lambda model: fake_shape_agent)

    importer = ContentStyleImporter(
        cm,
        page_classifier=page_classifier,
        shape_role_agent=shape_role_agent,
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),  # inject fake model, skip DB
        target_json_path=str(json_path),
        dry_run=True,
    )

    # Cover page should be skipped
    cover_slide = [s for s in report.slides if s.page_type == PageType.COVER]
    assert len(cover_slide) == 1
    assert cover_slide[0].status == ImportSlideStatus.SKIPPED


@pytest.mark.anyio
async def test_import_from_pptx_dry_run_does_not_modify_json(tmp_path):
    prs, slide = _two_point_slide_with_font()
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    json_path = tmp_path / "shapes.json"
    original_data = {"one_point": {}, "two_points": {}, "three_points": {}}
    json_path.write_text(json.dumps(original_data))

    cm = ComponentsManager(str(json_path))

    # Use FakeAgent for page classifier (returns chapter_content for all pages)
    fake_page_agent = FakePageClassifierAgent()
    page_classifier = PageTypeClassifier(agent_factory=lambda model: fake_page_agent)

    fake_shape_agent = FakeAgentForImport(
        content={
            "point_count": 2,
            "assignments": [],
            "confidence": 0.95,
            "reason": "Clear two-point layout.",
        }
    )
    shape_role_agent = ShapeRoleAgent(agent_factory=lambda model: fake_shape_agent)

    importer = ContentStyleImporter(
        cm,
        page_classifier=page_classifier,
        shape_role_agent=shape_role_agent,
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),  # inject fake model, skip DB
        target_json_path=str(json_path),
        dry_run=True,
    )

    # JSON file should be unchanged
    with open(json_path) as f:
        data = json.load(f)
    assert data == original_data
    assert cm.get_layout_type(ChapterLayout.TWO_POINTS).get_style("upload_p1") is None

    # Report should have DRY_RUN status for content slides
    content_slides = [s for s in report.slides if s.status == ImportSlideStatus.DRY_RUN]
    assert len(content_slides) >= 1


@pytest.mark.anyio
async def test_import_from_pptx_uses_unique_name_when_generated_collision_exists(tmp_path):
    prs, slide = _two_point_slide_with_font()
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    existing_suffix = hashlib.sha256(b"upload_p1").hexdigest()[:6]
    existing_names = {"upload_p1", f"upload_p1_{existing_suffix}"}
    json_path = tmp_path / "shapes.json"
    json_path.write_text(
        json.dumps(
            {
                "one_point": {},
                "two_points": {name: {} for name in existing_names},
                "three_points": {},
            }
        )
    )
    cm = ComponentsManager(str(json_path))

    importer = ContentStyleImporter(
        cm,
        page_classifier=PageTypeClassifier(agent_factory=lambda model: FakePageClassifierAgent()),
        shape_role_agent=ShapeRoleAgent(
            agent_factory=lambda model: FakeAgentForImport(
                content={
                    "point_count": 2,
                    "assignments": [],
                    "confidence": 0.95,
                    "reason": "Clear two-point layout.",
                }
            )
        ),
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),
        target_json_path=str(json_path),
        dry_run=False,
        overwrite_existing=False,
    )

    imported = [s for s in report.slides if s.status == ImportSlideStatus.IMPORTED]
    assert len(imported) == 1
    assert imported[0].style_name not in existing_names


def test_generate_style_name_uses_english_stem_for_chinese_pptx(tmp_path):
    prs, _slide = _blank_slide()
    pptx_path = tmp_path / "深度学习原理架构与应用.pptx"
    prs.save(str(pptx_path))

    importer = ContentStyleImporter(ComponentsManager(), style_name_stem="deck")

    style_name = importer._generate_style_name(pptx_path, 0)

    assert re.fullmatch(r"upload_deck_[0-9a-f]{8}_p1", style_name)


@pytest.mark.anyio
async def test_import_from_pptx_skips_duplicate_style_in_same_layout(tmp_path):
    prs, slide = _two_point_slide_with_font()
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    json_path = tmp_path / "shapes.json"
    json_path.write_text(json.dumps({"one_point": {}, "two_points": {}, "three_points": {}}))
    cm = ComponentsManager(str(json_path))

    local_classifier = LocalShapeRoleClassifier()
    summaries = PageTypeClassifier().summarize_slide(slide)
    assignments = local_classifier.assign_group_indices(
        local_classifier.classify(summaries),
        summaries,
    )
    existing_style = StyleBuilder().build_style_from_assignments(
        slide=slide,
        _layout_type=ChapterLayout.TWO_POINTS,
        style_name="existing_duplicate",
        assignments=assignments,
    )
    cm.get_layout_type(ChapterLayout.TWO_POINTS).add_style(existing_style)

    importer = ContentStyleImporter(
        cm,
        page_classifier=PageTypeClassifier(agent_factory=lambda model: FakePageClassifierAgent()),
        shape_role_agent=ShapeRoleAgent(
            agent_factory=lambda model: FakeAgentForImport(
                content={
                    "point_count": 2,
                    "assignments": [],
                    "confidence": 0.95,
                    "reason": "Clear two-point layout.",
                }
            )
        ),
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),
        target_json_path=str(json_path),
        dry_run=False,
        overwrite_existing=False,
    )

    assert report.imported_count == 0
    assert report.skipped_count == 1
    assert report.slides[0].status == ImportSlideStatus.SKIPPED
    assert report.slides[0].style_name == "existing_duplicate"
    assert "Duplicate style" in report.slides[0].reason
    assert cm.get_layout_type(ChapterLayout.TWO_POINTS).style_names == ["existing_duplicate"]


@pytest.mark.anyio
async def test_import_from_pptx_honors_min_role_confidence(tmp_path):
    prs, slide = _blank_slide()
    slide.shapes.add_textbox(Emu(900000), Emu(900000), Emu(3000000), Emu(400000)).text = "Signal"
    slide.shapes.add_textbox(
        Emu(900000), Emu(1400000), Emu(5000000), Emu(1000000)
    ).text = "Customers need a faster onboarding path with clearer operational support."
    slide.shapes.add_textbox(Emu(5000000), Emu(900000), Emu(1000000), Emu(300000)).text = "Note"
    slide.shapes.add_textbox(Emu(5000000), Emu(1300000), Emu(1000000), Emu(300000)).text = "Tag"
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    json_path = tmp_path / "shapes.json"
    json_path.write_text(json.dumps({"one_point": {}, "two_points": {}, "three_points": {}}))
    cm = ComponentsManager(str(json_path))

    saved_slide = Presentation(str(pptx_path)).slides[0]
    shape_id = next(s.shape_id for s in PageTypeClassifier().summarize_slide(saved_slide) if s.text == "Signal")
    fake_shape_agent = FakeAgentForImport(
        content={
            "point_count": 1,
            "assignments": [
                {
                    "shape_id": shape_id,
                    "content_type": "title",
                    "group_index": 0,
                    "confidence": 0.6,
                    "reason": "Short text is the content group title.",
                }
            ],
            "confidence": 0.6,
            "reason": "Resolved with lower confidence accepted by caller threshold.",
        }
    )
    importer = ContentStyleImporter(
        cm,
        page_classifier=PageTypeClassifier(agent_factory=lambda model: FakePageClassifierAgent()),
        shape_role_agent=ShapeRoleAgent(agent_factory=lambda model: fake_shape_agent),
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),
        target_json_path=str(json_path),
        min_role_confidence=0.5,
        dry_run=True,
    )

    assert report.slides[0].status == ImportSlideStatus.DRY_RUN
    assert report.slides[0].layout == ChapterLayout.ONE_POINT


@pytest.mark.anyio
async def test_import_from_pptx_reports_excluded_non_placeholder_shapes(tmp_path):
    prs, slide = _two_point_slide_with_font()
    title = slide.shapes.add_textbox(Emu(100000), Emu(10000), Emu(8000000), Emu(400000))
    tf = title.text_frame
    tf.paragraphs[0].text = "Slide Level Title"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    json_path = tmp_path / "shapes.json"
    json_path.write_text(json.dumps({"one_point": {}, "two_points": {}, "three_points": {}}))
    cm = ComponentsManager(str(json_path))
    importer = ContentStyleImporter(
        cm,
        page_classifier=PageTypeClassifier(agent_factory=lambda model: FakePageClassifierAgent()),
        shape_role_agent=ShapeRoleAgent(
            agent_factory=lambda model: FakeAgentForImport(
                content={
                    "point_count": 2,
                    "assignments": [],
                    "confidence": 0.95,
                    "reason": "Keep local assignments.",
                }
            )
        ),
    )

    report = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),
        target_json_path=str(json_path),
        dry_run=True,
    )

    assert any("Slide-level title excluded" in warning for warning in report.slides[0].warnings)


@pytest.mark.anyio
async def test_import_missing_pptx_raises_file_not_found(tmp_path):
    json_path = tmp_path / "shapes.json"
    json_path.write_text("{}")
    cm = ComponentsManager(str(json_path))

    importer = ContentStyleImporter(cm)

    with pytest.raises(FileNotFoundError):
        await importer.import_from_pptx(
            pptx_path=str(tmp_path / "missing.pptx"),
            user_id=uuid.uuid4(),
            model=object(),  # inject fake model
            target_json_path=str(json_path),
        )


@pytest.mark.anyio
async def test_import_fingerprint_dedup_prevents_reimport(tmp_path):
    prs, slide = _two_point_slide_with_font()
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    json_path = tmp_path / "shapes.json"
    json_path.write_text(json.dumps({"one_point": {}, "two_points": {}, "three_points": {}}))
    cm = ComponentsManager(str(json_path))

    fake_page_agent = FakePageClassifierAgent()
    page_classifier = PageTypeClassifier(agent_factory=lambda model: fake_page_agent)
    fake_shape_agent = FakeAgentForImport(
        content={
            "point_count": 2,
            "assignments": [],
            "confidence": 0.9,
            "reason": "Two-point layout.",
        }
    )
    shape_role_agent = ShapeRoleAgent(agent_factory=lambda model: fake_shape_agent)

    importer = ContentStyleImporter(
        cm,
        page_classifier=page_classifier,
        shape_role_agent=shape_role_agent,
    )

    # First import (not dry_run)
    await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),  # inject fake model, skip DB
        target_json_path=str(json_path),
        dry_run=False,
    )

    # Second import of same PPT — should be skipped via fingerprint
    report2 = await importer.import_from_pptx(
        pptx_path=str(pptx_path),
        user_id=uuid.uuid4(),
        model=object(),  # inject fake model, skip DB
        target_json_path=str(json_path),
        dry_run=False,
    )

    assert report2.imported_count == 0
    assert any(s.status == ImportSlideStatus.SKIPPED for s in report2.slides)


# --- Step 3: Shape role agent ---


@pytest.mark.anyio
async def test_shape_role_agent_assigns_roles_with_fake_agent():
    prs, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Note"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)

    local = LocalShapeRoleClassifier()
    local_assignments = local.classify(summaries)

    # Find the ambiguous shape
    ambiguous_ids = [a.shape_id for a in local_assignments if a.content_type == "skip" and not a.include]
    assert len(ambiguous_ids) >= 1, "Need at least one ambiguous shape for this test"

    fake_agent = FakeAgentForImport(
        content={
            "point_count": 1,
            "assignments": [
                {
                    "shape_id": ambiguous_ids[0],
                    "content_type": "title",
                    "group_index": 0,
                    "confidence": 0.8,
                    "reason": "Short text near top looks like a title.",
                },
            ],
            "confidence": 0.8,
            "reason": "Single-point layout with title.",
        }
    )

    agent = ShapeRoleAgent(agent_factory=lambda model: fake_agent)
    page_classification = PageClassification(
        page_index=0,
        page_type=PageType.CHAPTER_CONTENT,
        confidence=0.9,
        reason="Content slide",
        method="rule",
    )

    merged = await agent.assign_roles(
        model=object(),
        summaries=summaries,
        local_assignments=local_assignments,
        page_classification=page_classification,
    )

    # The ambiguous shape should now have a resolved role
    resolved = [a for a in merged if a.shape_id == ambiguous_ids[0] and a.include]
    assert len(resolved) == 1
    assert resolved[0].content_type == ComponentContentType.TITLE


@pytest.mark.anyio
async def test_shape_role_agent_fallback_on_llm_failure():
    prs, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Note"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)

    local = LocalShapeRoleClassifier()
    local_assignments = local.classify(summaries)

    fake_agent = FakeAgentForImport(exc=RuntimeError("provider failed"))
    agent = ShapeRoleAgent(agent_factory=lambda model: fake_agent)

    merged = await agent.assign_roles(
        model=object(),
        summaries=summaries,
        local_assignments=local_assignments,
        page_classification=PageClassification(
            page_index=0,
            page_type=PageType.CHAPTER_CONTENT,
            confidence=0.9,
            reason="test",
            method="rule",
        ),
    )

    # On failure, local assignments should be returned unchanged
    assert merged == local_assignments


@pytest.mark.anyio
async def test_shape_role_agent_ignores_invalid_shape_id():
    prs, slide = _blank_slide()
    slide.shapes.add_textbox(100000, 100000, 3000000, 400000).text = "Note"

    classifier = PageTypeClassifier()
    summaries = classifier.summarize_slide(slide)

    local = LocalShapeRoleClassifier()
    local_assignments = local.classify(summaries)

    # Find ambiguous ids
    [a.shape_id for a in local_assignments if a.content_type == "skip" and not a.include]

    fake_agent = FakeAgentForImport(
        content={
            "point_count": 1,
            "assignments": [
                {
                    "shape_id": 99999,  # nonexistent shape id
                    "content_type": "title",
                    "group_index": 0,
                    "confidence": 0.8,
                    "reason": "Invalid shape id.",
                },
            ],
            "confidence": 0.8,
            "reason": "Bad data.",
        }
    )

    agent = ShapeRoleAgent(agent_factory=lambda model: fake_agent)

    merged = await agent.assign_roles(
        model=object(),
        summaries=summaries,
        local_assignments=local_assignments,
        page_classification=PageClassification(
            page_index=0,
            page_type=PageType.CHAPTER_CONTENT,
            confidence=0.9,
            reason="test",
            method="rule",
        ),
    )

    # No changes should be made for invalid shape_id
    assert merged == local_assignments


# --- Validation render round-trip tests ---


def _make_style_with_one_text_shape(name: str = "test_style") -> Style:
    """Build a minimal valid one_point style with one TITLE + one CONTENT shape."""
    style = Style(name)
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr><p:cNvPr id="1" name="title"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="100000" y="100000"/>'
        '<a:ext cx="5000000" cy="600000"/></a:xfrm></p:spPr>'
        "<p:txBody><a:bodyPr/><a:p><a:r><a:t>placeholder</a:t></a:r></a:p></p:txBody>"
        "</p:sp>"
    )

    title_shape = CShape(
        xml=xml,
        zorder=0,
        content_type=ComponentContentType.TITLE,
        location=[Location(x=100000, y=100000, width=5000000, height=600000)],
    )
    content_shape = CShape(
        xml=xml,
        zorder=1,
        content_type=ComponentContentType.CONTENT,
        location=[Location(x=100000, y=800000, width=5000000, height=2000000)],
    )
    style.add_shape("title_0", title_shape)
    style.add_shape("content_0", content_shape)
    return style


@pytest.mark.anyio
async def test_validation_passes_for_valid_one_point_style():
    """A style with one TITLE + one CONTENT location should render successfully."""
    style = _make_style_with_one_text_shape()
    cm = ComponentsManager()
    importer = ContentStyleImporter(cm)

    result = await importer._validate_render_roundtrip(
        style=style,
        layout_type=ChapterLayout.ONE_POINT,
    )

    assert result.ok, f"Expected OK but got: {result.reason}"
    assert "OK" in result.reason


@pytest.mark.anyio
async def test_validation_fails_when_location_count_mismatches_point_count():
    """A one_point style with 2 TITLE locations should fail — counts don't match."""
    style = Style("bad_style")
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr><p:cNvPr id="1" name="title"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="100000" y="100000"/>'
        '<a:ext cx="5000000" cy="600000"/></a:xfrm></p:spPr>'
        "<p:txBody><a:bodyPr/><a:p><a:r><a:t>placeholder</a:t></a:r></a:p></p:txBody>"
        "</p:sp>"
    )
    title_shape = CShape(
        xml=xml,
        zorder=0,
        content_type=ComponentContentType.TITLE,
        location=[
            Location(x=100000, y=100000, width=5000000, height=600000),
            Location(x=100000, y=2000000, width=5000000, height=600000),
        ],
    )
    content_shape = CShape(
        xml=xml,
        zorder=1,
        content_type=ComponentContentType.CONTENT,
        location=[Location(x=100000, y=800000, width=5000000, height=2000000)],
    )
    style.add_shape("title_0", title_shape)
    style.add_shape("content_0", content_shape)

    cm = ComponentsManager()
    importer = ContentStyleImporter(cm)

    result = await importer._validate_render_roundtrip(
        style=style,
        layout_type=ChapterLayout.ONE_POINT,
    )

    assert not result.ok


@pytest.mark.anyio
async def test_validation_fails_with_broken_xml():
    """A style whose shape XML is not valid p:sp should fail."""
    style = Style("broken")
    broken = CShape(
        xml="<not>valid</not>",
        zorder=0,
        content_type=ComponentContentType.TITLE,
        location=[Location(x=100000, y=100000, width=5000000, height=600000)],
    )
    content = CShape(
        xml=(
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvPr id="2" name="content"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            '<p:spPr><a:xfrm><a:off x="100000" y="800000"/>'
            '<a:ext cx="5000000" cy="2000000"/></a:xfrm></p:spPr>'
            "<p:txBody><a:bodyPr/><a:p><a:r><a:t>placeholder</a:t></a:r></a:p></p:txBody>"
            "</p:sp>"
        ),
        zorder=1,
        content_type=ComponentContentType.CONTENT,
        location=[Location(x=100000, y=800000, width=5000000, height=2000000)],
    )
    style.add_shape("broken_title", broken)
    style.add_shape("content_0", content)

    cm = ComponentsManager()
    importer = ContentStyleImporter(cm)

    result = await importer._validate_render_roundtrip(
        style=style,
        layout_type=ChapterLayout.ONE_POINT,
    )

    assert not result.ok
