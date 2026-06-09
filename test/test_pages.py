import asyncio
import os
import sys
import zipfile
from collections import Counter
from io import BytesIO
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.append(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slidegen"))

import slidegen.services.presentation.pages as pages_module
from slidegen.schemas.theme import PresentationTheme, ThemeColorMapping
from slidegen.services.document import MarkdownDocument
from slidegen.services.document.markdown.elements import Heading
from slidegen.services.presentation.components import ComponentContentType, CShape, Location, Style
from slidegen.services.presentation.orchestrator import PresentationOrchestrator
from slidegen.services.presentation.pages import (
    CatalogPage,
    ChapterContentPage,
    ChapterHomePage,
    CoverPage,
    Page,
    ShapeFactory,
)


class TestPages:
    """test pages"""

    @pytest.fixture
    def markdown_path(self):
        """test markdown file path"""
        return os.path.join(os.path.dirname(__file__), "data", "report.md")

    @pytest.fixture
    def template_path(self):
        """test template file path"""
        return os.path.join(os.path.dirname(__file__), "data", "template_0.pptx")

    @pytest.fixture
    def presentation(self, template_path):
        """test presentation"""
        return Presentation(template_path)

    def test_markdown_document_parse(self, markdown_path):
        """test MarkdownDocument"""
        doc = MarkdownDocument(markdown_path)

        assert doc is not None
        headings = [elem for elem in doc.main.children]
        assert len(headings) > 0

    @pytest.fixture
    def markdown_document(self, markdown_path):
        doc = MarkdownDocument(markdown_path)
        return doc

    @pytest.fixture
    def heading_list(self, markdown_document):
        return [elem for elem in markdown_document.main.children]

    async def test_cover_page_generation(self, presentation):
        """test CoverPage"""
        title = Heading(level=1, text="Hello World!")

        await CoverPage.generate_slide(presentation, title, cover_page_index=0)

        temp_output = os.path.join(os.path.dirname(__file__), "test_cover.pptx")
        presentation.save(temp_output)

        assert os.path.exists(temp_output)

    def test_cover_page_generation_keeps_title_word_wrap_disabled(self, presentation):
        """Cover title placeholders should keep the previous no-wrap behavior."""
        title = Heading(level=1, text="A Long Cover Title")

        asyncio.run(CoverPage.generate_slide(presentation, title, cover_page_index=0))

        title_placeholder = next(
            shape
            for shape in presentation.slides[0].shapes.placeholders
            if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        )
        assert title_placeholder.text_frame.word_wrap is False

    def test_cover_page_generation_expands_narrow_title_box_first(self, presentation):
        """Long no-wrap titles should first expand their text box."""
        title_placeholder = next(
            shape
            for shape in presentation.slides[0].shapes.placeholders
            if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        )
        title_placeholder.width = title_placeholder.width // 5
        narrow_width = title_placeholder.width

        title = Heading(level=1, text="A Very Long Strategy Report Title That Needs More Horizontal Space")

        asyncio.run(CoverPage.generate_slide(presentation, title, cover_page_index=0))

        assert title_placeholder.text_frame.word_wrap is False
        assert title_placeholder.width > narrow_width
        assert title_placeholder.left + title_placeholder.width <= presentation.slide_width

    def test_injected_page_placeholder_scales_to_current_slide_size(self, presentation, monkeypatch):
        """Fallback page placeholders from shape.json should scale to the active template size."""
        source_width = int(presentation.slide_width)
        source_height = int(presentation.slide_height)
        presentation.slide_width = Emu(source_width * 2)
        presentation.slide_height = Emu(source_height * 3)

        blank_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        source_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_text_frame)
        fallback_location = Location(x=1000, y=2000, width=3000000, height=400000)
        fallback_placeholder = SimpleNamespace(
            xml=source_shape.element.xml,
            location=fallback_location,
        )

        class FakeComponentsManager:
            metadata = {"slide_width": source_width, "slide_height": source_height}

            def get_page_placeholder(self, _page_type, _role):
                return fallback_placeholder

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())

        injected = Page._find_or_inject_placeholder(
            slide=blank_slide,
            page_type="cover",
            role="title",
            text="T",
            placeholder_types=(),
        )

        assert injected is not None
        assert injected.left == 2000
        assert injected.top == 6000
        assert injected.width == 6000000
        assert injected.height == 1200000

    def test_shape_factory_creates_text_and_image_shapes(self, tmp_path):
        """ShapeFactory should centralize XML text and image shape creation."""
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        source_shape = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000))
        source_shape.text = "placeholder"
        source_xml = source_shape.element.xml

        text_shape = ShapeFactory.create_text_shape(
            slide,
            source_xml,
            "Factory text",
            Location(x=100000, y=200000, width=1200000, height=300000),
            shape_id=7,
            shape_name="factory_text",
        )

        assert text_shape.text == "Factory text"
        assert text_shape.left == 100000
        assert text_shape.top == 200000
        assert text_shape.width == 1200000
        assert text_shape.height == 300000

        image_path = tmp_path / "factory.png"
        Image.new("RGBA", (8, 8), (10, 20, 30, 255)).save(image_path)
        image_shape = ShapeFactory.create_image_shape(
            slide,
            str(image_path),
            Location(x=300000, y=400000, width=500000, height=600000),
        )

        assert image_shape.shape_type.name == "PICTURE"
        assert image_shape.left == 300000
        assert image_shape.top == 400000
        assert image_shape.width == 500000
        assert image_shape.height == 600000

    def test_page_shape_alignment_justifies_text_at_top(self):
        """Text alignment should be a shared Page-level shape operation."""
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000))
        shape.text = "Aligned"

        Page._shape_alignment(shape)

        assert shape.text_frame.vertical_anchor == MSO_ANCHOR.TOP
        assert all(paragraph.alignment == PP_ALIGN.JUSTIFY for paragraph in shape.text_frame.paragraphs)

    def _add_catalog_item(
        self,
        slide,
        *,
        number: str,
        text: str,
        left: int,
        top: int,
        dx_text: int = 700000,
        dy_text: int = 0,
    ):
        number_shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(420000), Emu(300000))
        number_shape.text = number
        text_shape = slide.shapes.add_textbox(Emu(left + dx_text), Emu(top + dy_text), Emu(2500000), Emu(300000))
        text_shape.text = text
        return number_shape, text_shape

    def _catalog_template_from_shapes(self, number_shape, text_shape, background_shape=None):
        return SimpleNamespace(
            number=SimpleNamespace(
                xml=number_shape._element.xml,
                zorder=1,
                text=number_shape.text,
                location=Location(
                    x=int(number_shape.left),
                    y=int(number_shape.top),
                    width=int(number_shape.width),
                    height=int(number_shape.height),
                ),
            ),
            text=SimpleNamespace(
                xml=text_shape._element.xml,
                zorder=2,
                text=text_shape.text,
                location=Location(
                    x=int(text_shape.left),
                    y=int(text_shape.top),
                    width=int(text_shape.width),
                    height=int(text_shape.height),
                ),
            ),
            background=None
            if background_shape is None
            else SimpleNamespace(
                xml=background_shape._element.xml,
                zorder=0,
                text=background_shape.text if background_shape.has_text_frame else None,
                location=Location(
                    x=int(background_shape.left),
                    y=int(background_shape.top),
                    width=int(background_shape.width),
                    height=int(background_shape.height),
                ),
            ),
        )

    def _install_catalog_library(self, monkeypatch, presentation, catalog_items):
        class FakeComponentsManager:
            metadata = {"slide_width": int(presentation.slide_width), "slide_height": int(presentation.slide_height)}

            def get_catalog_items(self, template_name):
                assert template_name == "general"
                return catalog_items

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())

    async def test_catalog_page_generation(self, presentation, heading_list):
        """test CatalogPage"""

        await CatalogPage.generate_slide(presentation, heading_list, catalog_page_index=1)
        # save generated PPT file for test
        temp_output = os.path.join(os.path.dirname(__file__), "test_catalog.pptx")
        presentation.save(temp_output)

        # verify if the file is successfully saved
        assert os.path.exists(temp_output)

    async def test_catalog_page_generation_keeps_all_chapters_across_recursive_pages(self):
        """CatalogPage should keep every chapter when recursive pages are needed."""
        presentation = Presentation(os.path.join(os.path.dirname(__file__), "data", "template_0.pptx"))
        heading_list = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 15)]

        last_catalog_index = await CatalogPage.generate_slide(presentation, heading_list, catalog_page_index=1)

        catalog_texts = []
        for slide_index in range(1, last_catalog_index + 1):
            catalog_texts.extend(
                shape.text.strip()
                for shape in presentation.slides[slide_index].shapes
                if shape.has_text_frame and shape.text.strip()
            )
        assert {f"Chapter {i}" for i in range(1, 15)} <= set(catalog_texts)

    async def test_catalog_page_cloned_shapes_keep_unique_ids_and_background_order(self):
        """Cloned catalog shapes should be valid and keep backgrounds under text."""
        presentation = Presentation(os.path.join(os.path.dirname(__file__), "data", "template_0.pptx"))
        catalog_slide = presentation.slides[1]
        catalog_items = CatalogPage._get_catalog_items(catalog_slide)

        for item in catalog_items[2:]:
            Page.remove_shapes(
                catalog_slide.shapes._spTree,
                [
                    item.number_shape["shape"],
                    item.text_shape["shape"],
                    *([item.background_shape["shape"]] if item.background_shape else []),
                ],
            )

        heading_list = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 5)]
        await CatalogPage.generate_slide(presentation, heading_list, catalog_page_index=1)

        shape_ids = [shape.shape_id for shape in catalog_slide.shapes]
        duplicate_ids = [shape_id for shape_id, count in Counter(shape_ids).items() if count > 1]
        assert duplicate_ids == []

        refreshed_items = CatalogPage._get_catalog_items(catalog_slide)
        shape_order = {shape._element: index for index, shape in enumerate(catalog_slide.shapes)}
        for item in refreshed_items:
            if item.background_shape is None:
                continue
            background_index = shape_order[item.background_shape["shape"]._element]
            assert background_index < shape_order[item.number_shape["shape"]._element]
            assert background_index < shape_order[item.text_shape["shape"]._element]

    @pytest.mark.anyio
    async def test_catalog_page_single_template_item_expands_on_same_slide(self):
        """A single catalog item should seed more items instead of one item per page."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        number_shape = catalog_slide.shapes.add_textbox(Emu(900000), Emu(1400000), Emu(420000), Emu(320000))
        number_shape.text = "01"
        text_shape = catalog_slide.shapes.add_textbox(Emu(1500000), Emu(1400000), Emu(4500000), Emu(320000))
        text_shape.text = "Existing"
        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 4)]

        last_catalog_index = await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0)

        assert last_catalog_index == 0
        catalog_texts = {
            shape.text.strip()
            for shape in catalog_slide.shapes
            if shape.has_text_frame and shape.text.strip()
        }
        assert {"01", "02", "03", "Chapter 1", "Chapter 2", "Chapter 3"} <= catalog_texts

    @pytest.mark.anyio
    async def test_catalog_page_blank_slide_uses_default_fallback_items(self):
        """A blank catalog slide should generate usable default catalog items."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 5)]

        last_catalog_index = await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0)

        assert last_catalog_index == 0
        catalog_texts = {
            shape.text.strip()
            for shape in catalog_slide.shapes
            if shape.has_text_frame and shape.text.strip()
        }
        assert {"01", "02", "03", "04", "Chapter 1", "Chapter 2", "Chapter 3", "Chapter 4"} <= catalog_texts

    @pytest.mark.anyio
    async def test_catalog_page_blank_slide_uses_shape_json_catalog_items(self, monkeypatch):
        """A blank catalog slide should use template-specific catalog items from shapes.json."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        number_shape = source_slide.shapes.add_textbox(Emu(1234567), Emu(2000000), Emu(600000), Emu(350000))
        number_shape.text = "01"
        text_shape = source_slide.shapes.add_textbox(Emu(2034567), Emu(2000000), Emu(3000000), Emu(350000))
        text_shape.text = ""

        catalog_item = SimpleNamespace(
            number=SimpleNamespace(
                xml=number_shape._element.xml,
                zorder=0,
                text="01",
                location=Location(
                    x=number_shape.left,
                    y=number_shape.top,
                    width=number_shape.width,
                    height=number_shape.height,
                ),
            ),
            text=SimpleNamespace(
                xml=text_shape._element.xml,
                zorder=1,
                text="",
                location=Location(
                    x=text_shape.left,
                    y=text_shape.top,
                    width=text_shape.width,
                    height=text_shape.height,
                ),
            ),
            background=None,
        )

        class FakeComponentsManager:
            metadata = {"slide_width": int(presentation.slide_width), "slide_height": int(presentation.slide_height)}

            def get_catalog_items(self, template_name):
                assert template_name == "general"
                return [catalog_item]

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 3)]

        await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0, template_name="general")

        catalog_items = CatalogPage._get_catalog_items(catalog_slide)
        assert catalog_items[0].number_shape["left"] == number_shape.left
        assert catalog_items[0].text_shape["left"] == text_shape.left
        catalog_texts = {
            shape.text.strip()
            for shape in catalog_slide.shapes
            if shape.has_text_frame and shape.text.strip()
        }
        assert {"01", "02", "Chapter 1", "Chapter 2"} <= catalog_texts

    @pytest.mark.anyio
    async def test_catalog_page_single_template_item_prefers_shape_json_catalog_items(self, monkeypatch):
        """A weak single-item slide should be replaced by template library catalog items."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_catalog_item(catalog_slide, number="01", text="Old single", left=3000000, top=2600000)

        source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        library_templates = []
        for index, top in enumerate([900000, 1400000, 1900000], start=1):
            number_shape, text_shape = self._add_catalog_item(
                source_slide,
                number=str(index).zfill(2),
                text="",
                left=900000,
                top=top,
            )
            library_templates.append(self._catalog_template_from_shapes(number_shape, text_shape))
        self._install_catalog_library(monkeypatch, presentation, library_templates)

        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 4)]
        last_catalog_index = await CatalogPage.generate_slide(
            presentation,
            headings,
            catalog_page_index=0,
            template_name="general",
        )

        assert last_catalog_index == 0
        catalog_items = CatalogPage._get_catalog_items(catalog_slide)
        assert [item.number_shape["left"] for item in catalog_items] == [900000, 900000, 900000]
        assert [item.number_shape["top"] for item in catalog_items] == [900000, 1400000, 1900000]
        visible_texts = {
            shape.text.strip()
            for shape in catalog_slide.shapes
            if shape.has_text_frame and shape.text.strip()
        }
        assert "Old single" not in visible_texts
        assert {"01", "02", "03", "Chapter 1", "Chapter 2", "Chapter 3"} <= visible_texts

    @pytest.mark.anyio
    async def test_catalog_page_insufficient_template_items_prefers_shape_json_catalog_items(self, monkeypatch):
        """A slide with too few catalog items should be replaced instead of locally cloned."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_catalog_item(catalog_slide, number="01", text="Old one", left=400000, top=400000)
        self._add_catalog_item(catalog_slide, number="02", text="Old two", left=900000, top=900000)

        source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        library_templates = []
        for index, top in enumerate([1200000, 1700000, 2200000, 2700000], start=1):
            number_shape, text_shape = self._add_catalog_item(
                source_slide,
                number=str(index).zfill(2),
                text="",
                left=1600000,
                top=top,
            )
            library_templates.append(self._catalog_template_from_shapes(number_shape, text_shape))
        self._install_catalog_library(monkeypatch, presentation, library_templates)

        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 5)]
        await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0, template_name="general")

        catalog_items = CatalogPage._get_catalog_items(catalog_slide)
        assert [item.number_shape["left"] for item in catalog_items] == [1600000, 1600000, 1600000, 1600000]
        visible_texts = {
            shape.text.strip()
            for shape in catalog_slide.shapes
            if shape.has_text_frame and shape.text.strip()
        }
        assert "Old one" not in visible_texts
        assert "Old two" not in visible_texts
        assert {"Chapter 1", "Chapter 2", "Chapter 3", "Chapter 4"} <= visible_texts

    @pytest.mark.anyio
    async def test_catalog_page_default_fallback_items_are_centered_and_larger(self):
        """Default fallback catalog items should read as a centered group."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 5)]

        await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0)

        catalog_items = CatalogPage._get_catalog_items(catalog_slide)
        item_shapes = [item.number_shape for item in catalog_items] + [item.text_shape for item in catalog_items]
        group_left = min(shape["left"] for shape in item_shapes)
        group_top = min(shape["top"] for shape in item_shapes)
        group_right = max(shape["left"] + shape["width"] for shape in item_shapes)
        group_bottom = max(shape["top"] + shape["height"] for shape in item_shapes)
        group_center_x = (group_left + group_right) / 2
        group_center_y = (group_top + group_bottom) / 2

        assert abs(group_center_x - int(presentation.slide_width) / 2) < 100000
        assert abs(group_center_y - int(presentation.slide_height) / 2) < 100000
        assert group_right - group_left < int(presentation.slide_width) * 0.7
        for item in catalog_items:
            number_shape = item.number_shape["shape"]
            text_shape = item.text_shape["shape"]
            assert number_shape.text_frame.paragraphs[0].font.size.pt >= 24
            assert text_shape.text_frame.paragraphs[0].font.size.pt >= 24

    @pytest.mark.anyio
    async def test_catalog_page_removes_empty_placeholders(self):
        """Empty layout placeholders should not remain visible in PowerPoint edit view."""
        presentation = Presentation()
        catalog_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        nonempty_placeholder = catalog_slide.shapes.placeholders[1]
        nonempty_placeholder.text = "目录"
        headings = [Heading(level=2, text=f"Chapter {i}") for i in range(1, 3)]

        await CatalogPage.generate_slide(presentation, headings, catalog_page_index=0)

        placeholders = list(catalog_slide.shapes.placeholders)
        assert all(shape.text.strip() for shape in placeholders if shape.has_text_frame)
        assert nonempty_placeholder in placeholders

    async def test_chapter_home_page_generation(self, presentation, heading_list):
        """test ChapterHomePage"""
        title = Heading(level=2, text="Hello World! This is a test title.")
        await ChapterHomePage.generate_slide(
            presentation, title, chapter_home_page_index=2, chapter_number=1, slide_index=2
        )
        temp_output = os.path.join(os.path.dirname(__file__), "test_chapter_home.pptx")
        presentation.save(temp_output)
        assert os.path.exists(temp_output)

    async def test_chapter_content_page_generation(self, presentation, heading_list):
        """test ChapterContentPage"""
        content = heading_list[0]
        await ChapterContentPage.generate_slide(presentation, content, chapter_page_index=4, slide_index=4)
        temp_output = os.path.join(os.path.dirname(__file__), "test_chapter_content.pptx")
        presentation.save(temp_output)
        assert os.path.exists(temp_output)

    @pytest.mark.anyio
    async def test_chapter_content_page_title_is_above_imported_style_shapes(self, monkeypatch):
        """Page-level titles should be brought above shapes imported from the style."""
        presentation = Presentation()
        for _ in range(5):
            presentation.slides.add_slide(presentation.slide_layouts[6])

        xml_source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        xml_source_shape = xml_source_slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000))
        xml_source_shape.text = "placeholder"

        fallback_placeholder = SimpleNamespace(
            xml=xml_source_shape.element.xml,
            location=Location(x=100000, y=100000, width=2000000, height=400000),
        )
        style = Style("covering_decoration")
        style.add_shape(
            "decoration_0",
            CShape(
                xml=xml_source_shape.element.xml,
                zorder=999,
                content_type=ComponentContentType.DECORATION,
                location=[Location(x=0, y=0, width=3000000, height=1000000)],
            ),
        )

        class FakeComponentsManager:
            def get_page_placeholder(self, _page_type, _role):
                return fallback_placeholder

            def get_random_style(self, _chapter_layout):
                return style

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())

        content = Heading(level=2, text="Validation Test")
        content.append(Heading(level=3, text="Point 1"))

        await ChapterContentPage.generate_slide(
            presentation,
            content,
            chapter_page_index=4,
            slide_index=4,
            style_override=style,
        )

        generated_slide = presentation.slides[4]
        assert generated_slide.shapes[-1].name == "content_title_title"

    @pytest.mark.anyio
    async def test_chapter_content_picture_uses_placeholder_when_generation_fails(self, presentation, monkeypatch):
        """Picture styles should not call add_picture with None when generation fails."""
        style = Style("picture_only")
        style.add_shape(
            "picture",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.PICTURE,
                location=[Location(x=Emu(900000), y=Emu(1200000), width=Emu(2500000), height=Emu(1600000))],
            ),
        )

        class FakeComponentsManager:
            def get_random_style(self, _chapter_layout):
                return style

            def get_page_placeholder(self, _page_type, _role):
                return None

        class FailingImageGenerator:
            async def generate_image(self, _prompt):
                return SimpleNamespace(path=None)

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "image_generator", FailingImageGenerator())

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))

        await ChapterContentPage.generate_slide(presentation, content, chapter_page_index=4, slide_index=4)

        generated_slide = presentation.slides[4]
        assert any(shape.shape_type.name == "PICTURE" for shape in generated_slide.shapes)

    @pytest.mark.anyio
    async def test_prepare_slide_data_uses_style_override_and_generated_picture(self, tmp_path, monkeypatch):
        """Preparation should resolve generated picture assets without mutating a slide."""
        generated_image = tmp_path / "generated.png"
        Image.new("RGBA", (8, 8), (255, 0, 0, 255)).save(generated_image)

        style = Style("picture_only")
        style.add_shape(
            "picture",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.PICTURE,
                location=[Location(x=Emu(900000), y=Emu(1200000), width=Emu(2500000), height=Emu(1600000))],
            ),
        )

        class ExplodingComponentsManager:
            def get_random_style(self, _chapter_layout):
                raise AssertionError("style_override should avoid random style selection")

        class FakeImageGenerator:
            async def generate_image(self, prompt):
                assert prompt.prompt == "Customer Signals"
                return SimpleNamespace(path=str(generated_image))

        monkeypatch.setattr(pages_module, "components_manager", ExplodingComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "image_generator", FakeImageGenerator())

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))

        slide_data = await ChapterContentPage._prepare_slide_data(content, style_override=style, theme=None)

        assert slide_data.content_title == "Market Context"
        assert slide_data.section_titles == ["Customer Signals"]
        assert slide_data.section_texts == [""]
        assert slide_data.style is style
        asset = slide_data.assets[("picture", 0, ComponentContentType.PICTURE)]
        assert asset.path == str(generated_image)

    @pytest.mark.anyio
    async def test_prepare_slide_data_resolves_icons_before_rendering(self, tmp_path, monkeypatch):
        """Preparation should search icons and apply theme recoloring before rendering."""
        source_icon = tmp_path / "source.png"
        recolored_icon = tmp_path / "recolored.png"
        Image.new("RGBA", (8, 8), (0, 0, 0, 255)).save(source_icon)
        Image.new("RGBA", (8, 8), (231, 111, 81, 255)).save(recolored_icon)

        style = Style("icon_pair")
        style.add_shape(
            "icon",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.ICON,
                location=[
                    Location(x=Emu(900000), y=Emu(1200000), width=Emu(250000), height=Emu(250000)),
                    Location(x=Emu(1300000), y=Emu(1200000), width=Emu(250000), height=Emu(250000)),
                ],
            ),
        )

        class FakeIconSearcher:
            async def search_icons(self, query, k=1):
                assert k == 1
                assert query in {"Customer Signals", "Competitive Position"}
                return [str(source_icon)]

        prepared_calls = []

        def fake_prepare_icon(icon_path, theme, icon_index):
            prepared_calls.append((icon_path, theme, icon_index))
            return str(recolored_icon)

        monkeypatch.setattr(ChapterContentPage, "icon_searcher", FakeIconSearcher())
        monkeypatch.setattr(ChapterContentPage, "_prepare_icon_for_theme", staticmethod(fake_prepare_icon))

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))
        content.append(Heading(level=3, text="Competitive Position"))
        theme = PresentationTheme(
            name="Icon Theme",
            colors=ThemeColorMapping(accent1="E76F51", accent2="0066FF"),
        )

        slide_data = await ChapterContentPage._prepare_slide_data(content, style_override=style, theme=theme)

        assert prepared_calls == [(str(source_icon), theme, 0), (str(source_icon), theme, 1)]
        assert slide_data.assets[("icon", 0, ComponentContentType.ICON)].path == str(recolored_icon)
        assert slide_data.assets[("icon", 1, ComponentContentType.ICON)].path == str(recolored_icon)

    def test_render_slide_is_sync_and_uses_prepared_assets(self, tmp_path, monkeypatch):
        """Rendering should mutate the slide without calling async asset providers."""
        prepared_image = tmp_path / "prepared.png"
        Image.new("RGBA", (8, 8), (0, 102, 255, 255)).save(prepared_image)

        presentation = Presentation()
        for _ in range(5):
            presentation.slides.add_slide(presentation.slide_layouts[6])
        new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        xml_source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_shape = xml_source_slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000))
        text_shape.text = "placeholder"

        style = Style("render_style")
        style.add_shape(
            "title",
            CShape(
                xml=text_shape.element.xml,
                zorder=0,
                content_type=ComponentContentType.TITLE,
                location=[Location(x=Emu(100000), y=Emu(100000), width=Emu(1000000), height=Emu(300000))],
            ),
        )
        style.add_shape(
            "content",
            CShape(
                xml=text_shape.element.xml,
                zorder=1,
                content_type=ComponentContentType.CONTENT,
                location=[Location(x=Emu(100000), y=Emu(500000), width=Emu(1000000), height=Emu(600000))],
            ),
        )
        style.add_shape(
            "picture",
            CShape(
                xml=None,
                zorder=2,
                content_type=ComponentContentType.PICTURE,
                location=[Location(x=Emu(1200000), y=Emu(500000), width=Emu(600000), height=Emu(600000))],
            ),
        )

        class FakeComponentsManager:
            def get_page_placeholder(self, _page_type, _role):
                return None

        class ExplodingImageGenerator:
            async def generate_image(self, _prompt):
                raise AssertionError("rendering should not generate images")

        class ExplodingIconSearcher:
            async def search_icons(self, _query, k=1):
                raise AssertionError("rendering should not search icons")

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "image_generator", ExplodingImageGenerator())
        monkeypatch.setattr(ChapterContentPage, "icon_searcher", ExplodingIconSearcher())

        slide_data = pages_module.ChapterSlideData(
            content_title="Market Context",
            section_titles=["Customer Signals"],
            section_texts=["Signals body"],
            style=style,
            assets={
                ("picture", 0, ComponentContentType.PICTURE): pages_module.ChapterSlideAsset(
                    shape_name="picture",
                    location_index=0,
                    content_type=ComponentContentType.PICTURE,
                    path=str(prepared_image),
                )
            },
        )

        result = ChapterContentPage._render_slide(presentation, new_slide, slide_data)

        assert result is None
        assert any(shape.shape_type.name == "PICTURE" for shape in new_slide.shapes)
        assert any(getattr(shape, "text", "") == "Customer Signals" for shape in new_slide.shapes)
        assert any(getattr(shape, "text", "") == "Signals body" for shape in new_slide.shapes)

    def test_render_slide_preserves_title_count_validation(self):
        """Rendering should keep the existing PPTGenError for mismatched title locations."""
        presentation = Presentation()
        for _ in range(5):
            presentation.slides.add_slide(presentation.slide_layouts[6])
        new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        xml_source_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_shape = xml_source_slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000))
        text_shape.text = "placeholder"

        style = Style("bad_title_count")
        style.add_shape(
            "title",
            CShape(
                xml=text_shape.element.xml,
                zorder=0,
                content_type=ComponentContentType.TITLE,
                location=[
                    Location(x=Emu(100000), y=Emu(100000), width=Emu(1000000), height=Emu(300000)),
                    Location(x=Emu(100000), y=Emu(500000), width=Emu(1000000), height=Emu(300000)),
                ],
            ),
        )

        slide_data = pages_module.ChapterSlideData(
            content_title="Market Context",
            section_titles=["Customer Signals"],
            section_texts=["Signals body"],
            style=style,
            assets={},
        )

        with pytest.raises(Exception, match="Title must be equal to the number of locations"):
            ChapterContentPage._render_slide(presentation, new_slide, slide_data)

    def test_recolor_png_icon_uses_theme_color_and_preserves_alpha(self, tmp_path):
        """Icon recoloring should replace visible RGB values while keeping the source alpha mask."""
        source_icon = tmp_path / "source.png"
        Image.new("RGBA", (2, 2), (0, 0, 0, 0)).save(source_icon)
        image = Image.open(source_icon).convert("RGBA")
        image.putpixel((0, 0), (0, 0, 0, 255))
        image.putpixel((1, 0), (0, 0, 0, 128))
        image.save(source_icon)

        recolored = ChapterContentPage._recolor_png_icon(str(source_icon), "E76F51", output_dir=tmp_path)
        recolored_image = Image.open(recolored).convert("RGBA")

        assert [pixel[3] for pixel in recolored_image.getdata()] == [255, 128, 0, 0]
        assert {pixel[:3] for pixel in recolored_image.getdata() if pixel[3] > 0} == {(231, 111, 81)}

    @pytest.mark.anyio
    async def test_chapter_content_icons_use_accent1_and_accent2_theme_colors(self, tmp_path, monkeypatch):
        """Generated icon pictures should be recolored with accent1/accent2 before insertion."""
        source_icon = tmp_path / "source.png"
        Image.new("RGBA", (8, 8), (0, 0, 0, 255)).save(source_icon)

        style = Style("icon_pair")
        style.add_shape(
            "icon",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.ICON,
                location=[
                    Location(x=Emu(900000), y=Emu(1200000), width=Emu(250000), height=Emu(250000)),
                    Location(x=Emu(1300000), y=Emu(1200000), width=Emu(250000), height=Emu(250000)),
                ],
            ),
        )

        class FakeComponentsManager:
            def get_random_style(self, _chapter_layout):
                return style

            def get_page_placeholder(self, _page_type, _role):
                return None

        class FakeIconSearcher:
            async def search_icons(self, _query, k=1):
                return [str(source_icon)]

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "icon_searcher", FakeIconSearcher())

        presentation = Presentation()
        for _ in range(5):
            presentation.slides.add_slide(presentation.slide_layouts[6])

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))
        content.append(Heading(level=3, text="Competitive Position"))
        theme = PresentationTheme(
            name="Icon Theme",
            colors=ThemeColorMapping(accent1="E76F51", accent2="0066FF"),
        )

        await ChapterContentPage.generate_slide(
            presentation,
            content,
            chapter_page_index=4,
            slide_index=4,
            theme=theme,
        )

        output_path = tmp_path / "icons.pptx"
        presentation.save(output_path)
        embedded_rgbs = []
        with zipfile.ZipFile(output_path) as pptx:
            for media_name in sorted(name for name in pptx.namelist() if name.startswith("ppt/media/")):
                image = Image.open(BytesIO(pptx.read(media_name))).convert("RGBA")
                visible_rgbs = {pixel[:3] for pixel in image.getdata() if pixel[3] > 0}
                embedded_rgbs.append(next(iter(visible_rgbs)))

        assert embedded_rgbs == [(231, 111, 81), (0, 102, 255)]

    @pytest.mark.anyio
    async def test_chapter_content_shapes_scale_to_current_slide_size(self, presentation, monkeypatch):
        """Shape library locations should scale from the source canvas to the active template size."""
        source_width = int(presentation.slide_width)
        source_height = int(presentation.slide_height)
        presentation.slide_width = Emu(source_width * 2)
        presentation.slide_height = Emu(source_height * 3)

        style = Style("scaled_picture")
        style.add_shape(
            "picture",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.PICTURE,
                location=[Location(x=1000, y=2000, width=3000, height=4000)],
            ),
        )

        class FakeComponentsManager:
            metadata = {"slide_width": source_width, "slide_height": source_height}

            def get_random_style(self, _chapter_layout):
                return style

            def get_page_placeholder(self, _page_type, _role):
                return None

        class FailingImageGenerator:
            async def generate_image(self, _prompt):
                return SimpleNamespace(path=None)

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "image_generator", FailingImageGenerator())

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))

        await ChapterContentPage.generate_slide(presentation, content, chapter_page_index=4, slide_index=4)

        generated_picture = next(shape for shape in presentation.slides[4].shapes if shape.shape_type.name == "PICTURE")
        assert generated_picture.left == 2000
        assert generated_picture.top == 6000
        assert generated_picture.width == 6000
        assert generated_picture.height == 12000

    @pytest.mark.anyio
    async def test_chapter_content_square_shape_scaling_keeps_square(self, presentation, monkeypatch):
        """Square shape locations should remain square when source and target slide ratios differ."""
        source_width = int(presentation.slide_width)
        source_height = int(presentation.slide_height)
        presentation.slide_width = Emu(source_width * 2)
        presentation.slide_height = Emu(source_height * 3)

        style = Style("scaled_square_picture")
        style.add_shape(
            "picture",
            CShape(
                xml=None,
                zorder=0,
                content_type=ComponentContentType.PICTURE,
                location=[Location(x=1000, y=2000, width=3000, height=3000)],
            ),
        )

        class FakeComponentsManager:
            metadata = {"slide_width": source_width, "slide_height": source_height}

            def get_random_style(self, _chapter_layout):
                return style

            def get_page_placeholder(self, _page_type, _role):
                return None

        class FailingImageGenerator:
            async def generate_image(self, _prompt):
                return SimpleNamespace(path=None)

        monkeypatch.setattr(pages_module, "components_manager", FakeComponentsManager())
        monkeypatch.setattr(ChapterContentPage, "image_generator", FailingImageGenerator())

        content = Heading(level=2, text="Market Context")
        content.append(Heading(level=3, text="Customer Signals"))

        await ChapterContentPage.generate_slide(presentation, content, chapter_page_index=4, slide_index=4)

        generated_picture = next(shape for shape in presentation.slides[4].shapes if shape.shape_type.name == "PICTURE")
        assert generated_picture.left == 2000
        assert generated_picture.top == 7500
        assert generated_picture.width == 6000
        assert generated_picture.height == 6000

    async def test_ppt_generation(self, presentation, markdown_document):
        """test PresentationOrchestrator"""
        ppt_gen = PresentationOrchestrator()
        template_prs = Presentation(os.path.join(os.path.dirname(__file__), "data", "template_0.pptx"))
        await ppt_gen.generate(template_prs, markdown_document)
        temp_output = os.path.join(os.path.dirname(__file__), "test_ppt.pptx")
        template_prs.save(temp_output)
        assert os.path.exists(temp_output)

    async def test_nested_ppt_generation(self):
        """test PresentationOrchestrator with chapter/section/topic markdown"""
        markdown_document = MarkdownDocument(
            "# Product Strategy\n"
            "## Market Context\n"
            "### Customer Signals\n"
            "#### Demand\n"
            "Customers want faster onboarding.\n"
            "### Competitive Position\n"
            "#### Differentiation\n"
            "Workflow depth matters most.\n"
            "### Buying Motion\n"
            "#### Procurement\n"
            "Enterprise cycles are longer.\n"
            "### Expansion Path\n"
            "#### Accounts\n"
            "Existing accounts need reporting.\n"
            "### Risk Review\n"
            "#### Delivery\n"
            "Scope needs tight sequencing.\n"
        )
        ppt_gen = PresentationOrchestrator()
        template_prs = Presentation(os.path.join(os.path.dirname(__file__), "data", "template_0.pptx"))

        await ppt_gen.generate(template_prs, markdown_document)

        temp_output = os.path.join(os.path.dirname(__file__), "test_nested_ppt.pptx")
        template_prs.save(temp_output)
        assert os.path.exists(temp_output)

    @pytest.mark.anyio
    async def test_presentation_orchestrator_passes_template_name_to_catalog_page(self, monkeypatch):
        """Template name should flow to CatalogPage so shape.json catalog items can be selected."""
        markdown_document = MarkdownDocument("# Product Strategy\n## Market Context\n### Customer Signals\n")
        presentation = Presentation()
        for _ in range(5):
            presentation.slides.add_slide(presentation.slide_layouts[6])
        captured = {}

        async def fake_cover(*_args, **_kwargs):
            return None

        async def fake_catalog(*_args, **kwargs):
            captured["template_name"] = kwargs.get("template_name")
            return kwargs.get("catalog_page_index", 1)

        async def fake_page(*_args, **_kwargs):
            return None

        monkeypatch.setattr(CoverPage, "generate_slide", fake_cover)
        monkeypatch.setattr(CatalogPage, "generate_slide", fake_catalog)
        monkeypatch.setattr(ChapterHomePage, "generate_slide", fake_page)
        monkeypatch.setattr(ChapterContentPage, "generate_slide", fake_page)
        monkeypatch.setattr(pages_module.EndPage, "generate_slide", fake_page)

        await PresentationOrchestrator().generate(presentation, markdown_document, template_name="purple")

        assert captured["template_name"] == "purple"
