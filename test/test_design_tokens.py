from pptx import Presentation

from slidegen.services.presentation.design_tokens import (
    DesignTokens,
    DEFAULT_TOKENS,
    PRESET_TOKENS,
    extract_design_tokens_from_presentation,
    _dominant_color,
)


class TestDesignTokens:
    def test_default_tokens_is_frozen(self):
        import dataclasses
        assert dataclasses.is_dataclass(DEFAULT_TOKENS)

    def test_preset_tokens_contain_expected_themes(self):
        assert "general" in PRESET_TOKENS
        assert "minimal" in PRESET_TOKENS
        assert "academic" in PRESET_TOKENS
        assert PRESET_TOKENS["minimal"].primary == "#2C2C2C"


class TestDominantColor:
    def test_empty_returns_none(self):
        assert _dominant_color([]) is None

    def test_ignores_black_and_white(self):
        samples = [("000000", 100.0), ("FFFFFF", 200.0)]
        assert _dominant_color(samples) is None

    def test_returns_weight_dominant(self):
        samples = [("FF0000", 100.0), ("00FF00", 300.0)]
        assert _dominant_color(samples) == "00FF00"

    def test_skips_black_white_and_returns_colored(self):
        samples = [("000000", 100.0), ("FFFFFF", 200.0), ("123456", 50.0)]
        assert _dominant_color(samples) == "123456"


class TestExtractDesignTokens:
    def test_extract_from_empty_presentation_returns_preset_tokens(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        tokens = extract_design_tokens_from_presentation(prs, "general")
        assert tokens.primary == DEFAULT_TOKENS.primary
        assert tokens.accent == DEFAULT_TOKENS.accent

    def test_extract_respects_theme_name(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        tokens = extract_design_tokens_from_presentation(prs, "academic")
        assert tokens.title_font == "Times New Roman"

    def test_extract_slide_dimensions(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        tokens = extract_design_tokens_from_presentation(prs, "general")
        assert tokens.slide_width == prs.slide_width / 914400.0
        assert tokens.slide_height == prs.slide_height / 914400.0
