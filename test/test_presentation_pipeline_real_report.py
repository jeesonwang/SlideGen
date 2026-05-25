from pathlib import Path

import pytest
from pptx import Presentation

from slidegen.schemas.gen_request import ExportFormat
from slidegen.services.presentation.generator import PresentationGenerator
from slidegen.services.presentation.post_render_validator import PostRenderValidator


@pytest.mark.anyio
async def test_report_markdown_runs_full_pipeline_to_real_pptx(tmp_path: Path) -> None:
    markdown = Path("test/data/report.md").read_text(encoding="utf-8")
    output_path = tmp_path / "report_pipeline.pptx"

    result = await PresentationGenerator().generate_from_markdown(
        markdown_content=markdown,
        template_name="general",
        output_path=str(output_path),
        export_as=ExportFormat.PPTX,
    )

    assert result == str(output_path)
    assert output_path.exists()
    assert output_path.stat().st_size > 0

    generated = Presentation(str(output_path))
    assert len(generated.slides) == 21

    slide_texts = [
        "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame).strip() for slide in generated.slides
    ]
    assert all(slide_texts)

    all_text = "\n".join(slide_texts)
    assert "抖音&快手 竞品分析报告" in all_text
    assert "一、产品概述" in all_text
    assert "五、未来发展" in all_text

    issues = PostRenderValidator(mode="fail").validate(generated)
    assert not [issue for issue in issues if issue.level == "error"]
