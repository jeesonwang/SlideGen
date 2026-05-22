import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from slidegen.exceptions import PPTTemplateError
from slidegen.services.document import MarkdownDocument
from slidegen.services.document.markdown.elements import Heading
from slidegen.services.presentation.converter import MarkdownToPresentation
from slidegen.services.presentation.render_plan import build_presentation_render_plan
from slidegen.services.presentation.semantic import BlockKind, SlideKind, build_content_slide_spec
from slidegen.services.presentation.template_profile import (
    TemplateRole,
    _has_keyword_match,
    profile_presentation_template,
)
from slidegen.services.slidegen.outline_structure import iter_chapter_slide_groups
from test.helpers import add_catalog_slide


def _template_path() -> str:
    return os.path.join(os.path.dirname(__file__), "data", "template_0.pptx")


def _markdown_document(markdown: str) -> MarkdownDocument:
    return MarkdownDocument(markdown)


def _add_text_slide(presentation: Presentation, lines: list[str], *, layout_index: int = 6) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    if layout_index == 0:
        slide.shapes.title.text = lines[0]
        if len(lines) > 1 and len(slide.placeholders) > 1:
            slide.placeholders[1].text = lines[1]
        return

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(4.5))
    frame = textbox.text_frame
    frame.text = lines[0]
    for line in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = line


def test_profile_curated_template_detects_roles():
    presentation = Presentation(_template_path())

    profile = profile_presentation_template(presentation)

    assert profile.slide_count >= 5
    assert profile.role_index(TemplateRole.COVER) == 0
    assert profile.role_index(TemplateRole.CATALOG) == 1
    assert profile.role_index(TemplateRole.CHAPTER) == 2
    assert profile.role_index(TemplateRole.CONTENT) == 3
    assert profile.role_index(TemplateRole.END) == 4
    assert profile.has_role(TemplateRole.CONTENT)
    assert profile.status == "ready"
    assert profile.missing_roles == []
    assert profile.warnings == []


def test_profile_scores_textbox_only_template_by_heuristics():
    presentation = Presentation()
    for lines in (
        ["Quarterly Business Review", "2026 Strategy Update"],
        ["Agenda", "1. Market", "2. Product", "3. Finance"],
        ["Chapter 1", "Market Landscape"],
        [
            "Market analysis",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand remains strongest in regulated industries.",
            "Retention stays above target.",
        ],
        ["Thank You", "Questions"],
    ):
        _add_text_slide(presentation, lines)

    profile = profile_presentation_template(presentation)

    assert profile.slide_count == 5
    # With heuristic scoring only (no legacy renderer structural check),
    # text content and position alone can assign roles on textbox-only templates.
    assert profile.role_index(TemplateRole.COVER) is not None


def test_profile_one_slide_template_warns_missing_roles_without_rejecting():
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])

    profile = profile_presentation_template(presentation)

    assert profile.slide_count == 1
    assert profile.role_index(TemplateRole.COVER) == 0
    assert profile.role_index(TemplateRole.CATALOG) is None
    assert profile.role_index(TemplateRole.CHAPTER) is None
    assert profile.role_index(TemplateRole.CONTENT) is None
    assert profile.role_index(TemplateRole.END) is None
    assert profile.status == "review_required"
    assert TemplateRole.CATALOG.value in profile.missing_roles
    assert TemplateRole.CONTENT.value in profile.missing_roles


def test_profile_empty_presentation_rejects_as_unusable():
    presentation = Presentation()

    with pytest.raises(PPTTemplateError, match="must contain at least one slide"):
        profile_presentation_template(presentation)


def test_keyword_matching_avoids_english_substring_false_positives():
    text = "department partnership discovery recovery coverage"

    assert not _has_keyword_match(text, ("part", "cover"))
    assert _has_keyword_match("Part 1: Market landscape", ("part",))
    assert _has_keyword_match("封面设计", ("封面",))


def test_profile_uses_global_assignment_when_cover_and_catalog_compete():
    presentation = Presentation()
    add_catalog_slide(presentation)
    _add_text_slide(presentation, ["Title", "Annual Business Review"], layout_index=0)
    _add_text_slide(presentation, ["Chapter 1", "Market Landscape"], layout_index=1)
    _add_text_slide(
        presentation,
        [
            "Market analysis",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand remains strongest in regulated industries.",
            "Retention stays above target.",
        ],
        layout_index=1,
    )
    _add_text_slide(presentation, ["Thank You", "Questions"], layout_index=1)

    profile = profile_presentation_template(presentation)

    assert profile.role_index(TemplateRole.CATALOG) == 0
    assert profile.role_index(TemplateRole.COVER) == 1
    assert profile.status == "ready"


def test_render_plan_builds_from_chapter_groups():
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    document = _markdown_document("# Deck\n## Chapter A\n### Point\nBody")
    assert document.main is not None
    groups = list(iter_chapter_slide_groups(document.main))

    plan = build_presentation_render_plan(groups, catalog_last_index=1)

    assert plan.chapters[0].home_slide_index == 2
    assert plan.chapters[0].content_slides[0].slide_index == 3
    assert plan.end_slide_index == 4


def test_markdown_to_presentation_is_stateless():
    converter = MarkdownToPresentation()

    assert not hasattr(converter, "slide_index")
    assert not hasattr(converter, "chapter_index")


@pytest.mark.anyio
async def test_converter_generates_from_one_slide_template():
    converter = MarkdownToPresentation()
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    presentation.slides.add_slide(presentation.slide_layouts[0])
    markdown_document = _markdown_document("# Deck\n## Chapter A\n### Point\nBody")

    result = await converter.generate(presentation, markdown_document)

    slide_texts = [
        shape.text.strip()
        for slide in result.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    assert "Deck" in slide_texts
    assert any("Point" in t for t in slide_texts)
    assert any("谢谢" in t for t in slide_texts)


@pytest.mark.anyio
async def test_converter_generates_from_textbox_template():
    converter = MarkdownToPresentation()
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    for lines in (
        ["Quarterly Business Review", "2026 Strategy Update"],
        ["Agenda", "1. Market", "2. Product", "3. Finance"],
        ["Chapter 1", "Market Landscape"],
        [
            "Market analysis",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand remains strongest in regulated industries.",
            "Retention stays above target.",
        ],
        ["Thank You", "Questions"],
    ):
        _add_text_slide(presentation, lines)
    markdown_document = _markdown_document("# Deck\n## Chapter A\n### Point\nBody")

    result = await converter.generate(presentation, markdown_document)

    slide_texts = [
        shape.text.strip()
        for slide in result.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    assert "Deck" in slide_texts
    assert any("Point" in t for t in slide_texts)
    assert any("谢谢" in t for t in slide_texts)
    # Recipe renderer creates a new presentation, so original template text is not present.
    assert "Quarterly Business Review" not in slide_texts


@pytest.mark.anyio
async def test_converter_creates_new_presentation_not_mutating_template():
    converter = MarkdownToPresentation()
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    _add_text_slide(presentation, ["Decorative cover", "No placeholder contract"])
    _add_text_slide(presentation, ["Intro page", "No catalog number shapes"])
    _add_text_slide(presentation, ["Chapter 1", "Market Landscape"], layout_index=1)
    _add_text_slide(
        presentation,
        [
            "Market analysis",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand remains strongest in regulated industries.",
            "Retention stays above target.",
        ],
        layout_index=1,
    )
    _add_text_slide(presentation, ["Thank You", "Questions"], layout_index=1)
    markdown_document = _markdown_document("# Deck\n## Chapter A\n### Point\nBody")

    result = await converter.generate(presentation, markdown_document)

    slide_texts = [
        shape.text.strip()
        for slide in result.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    assert "Deck" in slide_texts
    assert "Chapter A" in slide_texts
    assert any("谢谢" in t for t in slide_texts)
    assert "Decorative cover" not in slide_texts


@pytest.mark.anyio
async def test_converter_generates_from_curated_template():
    converter = MarkdownToPresentation()
    presentation = Presentation(_template_path())
    markdown_document = _markdown_document("# Deck\n## Chapter A\n### Point\nBody")

    result = await converter.generate(presentation, markdown_document)

    assert len(result.slides) >= 5


def test_build_content_slide_spec_maps_heading_children_to_point_blocks():
    document = _markdown_document("# Deck\n## Chapter A\n### Point 1\nBody 1\n### Point 2\nBody 2\n")
    assert document.main is not None
    group = next(iter_chapter_slide_groups(document.main))

    spec = build_content_slide_spec(group.slides[0])

    assert spec.kind is SlideKind.CONTENT_POINTS
    assert spec.title == "Chapter A"
    assert [block.kind for block in spec.blocks] == [BlockKind.POINT, BlockKind.POINT]
    assert [block.title for block in spec.blocks] == ["Point 1", "Point 2"]
    assert [block.text for block in spec.blocks] == ["Body 1", "Body 2"]


@pytest.mark.anyio
async def test_recipe_renderer_produces_full_deck():
    template = Presentation()
    template.slide_width = 12192000
    template.slide_height = 6858000
    template.slides.add_slide(template.slide_layouts[6])
    document = MarkdownDocument(
        "# Deck\n"
        "## Chapter A\n"
        "### Point One\n"
        "#### Detail\n"
        "Body A\n"
        "### Point Two\n"
        "#### Detail\n"
        "Body B\n"
    )

    prs = await MarkdownToPresentation().generate(template, document)

    assert prs is not template
    assert prs.slide_width == template.slide_width
    assert prs.slide_height == template.slide_height
    assert len(prs.slides) == 6  # cover, agenda, chapter home, 2 content slides, closing
    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame)
    assert "Deck" in all_text
    assert "Chapter A" in all_text
    assert "Point One" in all_text
    assert "Body A" in all_text
