import uuid
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, Mock

import pytest
from fastapi import FastAPI, UploadFile
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from slidegen.api.deps import get_current_user
from slidegen.api.routers.slidegen import router as slidegen_router
from slidegen.core.database import get_db_session
from slidegen.models.presentation_template import PresentationTemplateModel
from slidegen.schemas.template import (
    Template,
    TemplateProfileResponse,
    TemplateRoleAssignmentResponse,
    TemplateSource,
    UserTemplateUploadResponse,
)
from slidegen.services.presentation.generator import PresentationGenerator
from slidegen.services.presentation.user_templates import (
    USER_TEMPLATE_KEY_PREFIX,
    UploadedTemplateService,
    UserTemplateStorage,
    parse_user_template_key,
    template_key_for_id,
)


def test_user_template_model_defaults_to_review_required() -> None:
    user_id = uuid.uuid4()
    template = PresentationTemplateModel(
        user_id=user_id,
        name="Board Template",
        original_filename="board.pptx",
        file_path="/tmp/board.pptx",
        file_size=1024,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_hash="a" * 64,
        template_key="user_1234567890abcdef1234567890abcdef",
        slide_count=5,
        role_profile={"assignments": []},
        warnings=[],
    )

    assert template.user_id == user_id
    assert template.status == "review_required"
    assert template.is_deleted is False


def test_template_schema_exposes_profile_for_uploaded_templates() -> None:
    assignment = TemplateRoleAssignmentResponse(
        role="cover",
        slide_index=0,
        confidence=0.95,
        reason="first slide",
    )
    profile = TemplateProfileResponse(
        slide_count=5,
        status="ready",
        assignments=[assignment],
        warnings=[],
        missing_roles=[],
    )

    template = Template(
        id="user_1234567890abcdef1234567890abcdef",
        name="Board Template",
        thumbnail=None,
        source=TemplateSource.USER,
        profile_status="ready",
        role_profile=profile,
    )

    assert template.source is TemplateSource.USER
    assert template.role_profile is profile


def test_user_template_upload_response_contains_template_key() -> None:
    template_id = uuid.uuid4()
    response = UserTemplateUploadResponse(
        id=template_id,
        template_key=f"user_{template_id.hex}",
        name="Board Template",
        original_filename="board.pptx",
        file_size=1024,
        profile=TemplateProfileResponse(
            slide_count=5,
            status="ready",
            assignments=[],
            warnings=[],
            missing_roles=[],
        ),
        message="Template uploaded successfully",
    )

    assert response.template_key == f"user_{template_id.hex}"


def _pptx_upload_bytes() -> bytes:
    prs = Presentation()
    for lines in (
        ["Quarterly Business Review", "2026 Strategy Update"],
        ["Agenda", "1. Market", "2. Product", "3. Finance"],
        ["Chapter 1", "Market"],
        ["Market", "Revenue grew 18 percent.", "Enterprise demand increased.", "Retention is stable."],
        ["Thank You", "Questions"],
    ):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(4.5))
        frame = textbox.text_frame
        frame.text = lines[0]
        for line in lines[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = line

    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def _upload_file(filename: str, content: bytes) -> UploadFile:
    return UploadFile(
        filename=filename,
        file=BytesIO(content),
        headers={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    )


class _FakeDbSession:
    def __init__(self) -> None:
        self.added = None
        self.add = Mock(side_effect=self._add)
        self.commit = AsyncMock()
        self.refresh = AsyncMock()

    def _add(self, obj) -> None:
        self.added = obj


def test_template_key_round_trips_uuid() -> None:
    template_id = uuid.uuid4()

    key = template_key_for_id(template_id)

    assert key == f"{USER_TEMPLATE_KEY_PREFIX}{template_id.hex}"
    assert parse_user_template_key(key) == template_id
    assert parse_user_template_key("general") is None


def test_storage_rejects_non_pptx_extension(tmp_path: Path) -> None:
    storage = UserTemplateStorage(tmp_path, max_file_size=10_000_000)

    with pytest.raises(ValueError, match="Only .pptx template files are supported"):
        storage.validate_upload("template.pdf", b"fake")


@pytest.mark.anyio
async def test_upload_service_saves_pptx_profiles_and_persists_model(tmp_path: Path) -> None:
    user_id = uuid.uuid4()
    db_session = _FakeDbSession()
    service = UploadedTemplateService(
        storage=UserTemplateStorage(tmp_path, max_file_size=10_000_000),
    )
    upload = _upload_file("board template.pptx", _pptx_upload_bytes())

    response = await service.upload_template(
        db_session=db_session,
        user_id=user_id,
        upload_file=upload,
        display_name="Board Template",
    )

    assert response.template_key.startswith(USER_TEMPLATE_KEY_PREFIX)
    assert response.name == "Board Template"
    assert response.profile.status == "ready"
    assert db_session.added is not None
    assert db_session.added.user_id == user_id
    assert db_session.added.template_key == response.template_key
    assert Path(db_session.added.file_path).exists()
    db_session.add.assert_called_once()
    db_session.commit.assert_awaited_once()
    db_session.refresh.assert_awaited_once()


def test_upload_template_endpoint_returns_profile(monkeypatch: pytest.MonkeyPatch) -> None:
    user_id = uuid.uuid4()
    uploaded_id = uuid.uuid4()

    async def override_get_current_user() -> SimpleNamespace:
        return SimpleNamespace(id=user_id)

    async def override_get_db_session() -> SimpleNamespace:
        return SimpleNamespace()

    async def fake_upload_template(db_session=None, user_id=None, upload_file=None, display_name=None):  # noqa: ARG001
        return UserTemplateUploadResponse(
            id=uploaded_id,
            template_key=f"user_{uploaded_id.hex}",
            name=display_name,
            original_filename=upload_file.filename,
            file_size=1234,
            profile=TemplateProfileResponse(
                slide_count=5,
                status="ready",
                assignments=[
                    TemplateRoleAssignmentResponse(
                        role="cover",
                        slide_index=0,
                        confidence=0.9,
                        reason="first slide",
                    )
                ],
                warnings=[],
                missing_roles=[],
            ),
        )

    monkeypatch.setattr(
        "slidegen.api.routers.slidegen.uploaded_template_service.upload_template",
        fake_upload_template,
    )

    app = FastAPI()
    app.include_router(slidegen_router, prefix="/slidegen")
    app.dependency_overrides[get_current_user] = override_get_current_user
    app.dependency_overrides[get_db_session] = override_get_db_session
    client = TestClient(app)

    response = client.post(
        "/slidegen/templates/upload",
        data={"display_name": "Board Template"},
        files={
            "file": (
                "board.pptx",
                _pptx_upload_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 200
    body = response.json()
    assert body["template_key"] == f"user_{uploaded_id.hex}"
    assert body["profile"]["status"] == "ready"


def test_templates_endpoint_merges_builtin_and_uploaded(monkeypatch: pytest.MonkeyPatch) -> None:
    user_id = uuid.uuid4()
    uploaded_id = uuid.uuid4()

    async def override_get_current_user() -> SimpleNamespace:
        return SimpleNamespace(id=user_id)

    async def override_get_db_session() -> SimpleNamespace:
        return SimpleNamespace()

    async def fake_list_templates(db_session=None, user_id=None):  # noqa: ARG001
        return [
            Template(
                id=f"user_{uploaded_id.hex}",
                name="Board Template",
                thumbnail=None,
                source=TemplateSource.USER,
                profile_status="ready",
            )
        ]

    monkeypatch.setattr(
        "slidegen.api.routers.slidegen.presentation_generator.list_templates",
        lambda: ["general"],
    )
    monkeypatch.setattr(
        "slidegen.api.routers.slidegen.uploaded_template_service.list_templates",
        fake_list_templates,
    )

    app = FastAPI()
    app.include_router(slidegen_router, prefix="/slidegen")
    app.dependency_overrides[get_current_user] = override_get_current_user
    app.dependency_overrides[get_db_session] = override_get_db_session
    client = TestClient(app)

    response = client.get("/slidegen/templates")

    assert response.status_code == 200
    body = response.json()
    assert body[0]["id"] == "general"
    assert body[0]["source"] == "builtin"
    assert body[1]["id"] == f"user_{uploaded_id.hex}"
    assert body[1]["source"] == "user"


def test_generator_resolves_uploaded_template_path(tmp_path: Path) -> None:
    user_id = uuid.uuid4()
    template_id = uuid.uuid4()
    storage = UserTemplateStorage(tmp_path, max_file_size=10_000_000)
    template_dir = storage.template_dir(user_id, template_id)
    template_dir.mkdir(parents=True)
    template_path = storage.template_path(user_id, template_id)
    template_path.write_bytes(_pptx_upload_bytes())

    generator = PresentationGenerator(templates_dir=str(tmp_path / "builtins"))
    generator.user_template_storage = storage

    resolved = generator.get_template_path(template_key_for_id(template_id), user_id=user_id)

    assert resolved == str(template_path)


def test_generator_rejects_uploaded_template_without_user_id(tmp_path: Path) -> None:
    template_id = uuid.uuid4()
    generator = PresentationGenerator(templates_dir=str(tmp_path / "builtins"))
    generator.user_template_storage = UserTemplateStorage(tmp_path, max_file_size=10_000_000)

    with pytest.raises(FileNotFoundError, match="requires user_id"):
        generator.get_template_path(template_key_for_id(template_id), user_id=None)
