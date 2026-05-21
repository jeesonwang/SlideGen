from io import BytesIO

import pytest
from pptx import Presentation
from pptx.util import Inches

from slidegen.exceptions import PPTTemplateError
from slidegen.services.presentation.template_profile import (
    TemplateRole,
    _has_keyword_match,
    profile_presentation_template,
)


def _add_text_slide(prs: Presentation, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(4.5))
    frame = textbox.text_frame
    frame.text = lines[0]
    for line in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = line


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _ordinary_template() -> Presentation:
    prs = Presentation()
    _add_text_slide(prs, ["Quarterly Business Review", "2026 Strategy Update"])
    _add_text_slide(prs, ["Agenda", "1. Market landscape", "2. Product strategy", "3. Financial outlook"])
    _add_text_slide(prs, ["Chapter 1", "Market Landscape"])
    _add_text_slide(
        prs,
        [
            "Market landscape",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand is strongest in regulated industries.",
            "Customer retention remains above the target range.",
        ],
    )
    _add_text_slide(prs, ["Thank You", "Questions and discussion"])
    return prs


def test_profile_detects_five_roles_from_ordinary_pptx() -> None:
    profile = profile_presentation_template(_ordinary_template())

    assert profile.slide_count == 5
    assert profile.role_index(TemplateRole.COVER) == 0
    assert profile.role_index(TemplateRole.CATALOG) == 1
    assert profile.role_index(TemplateRole.CHAPTER) == 2
    assert profile.role_index(TemplateRole.CONTENT) == 3
    assert profile.role_index(TemplateRole.END) == 4
    assert profile.status == "ready"
    assert profile.missing_roles == []


def test_profile_marks_missing_roles_as_review_required() -> None:
    prs = Presentation()
    _add_text_slide(prs, ["Only One Slide"])

    profile = profile_presentation_template(prs)

    assert profile.status == "review_required"
    assert profile.role_index(TemplateRole.COVER) == 0
    assert TemplateRole.CATALOG.value in profile.missing_roles
    assert TemplateRole.CONTENT.value in profile.missing_roles
    assert "catalog role not detected" in " ".join(profile.warnings)


def test_profile_rejects_empty_presentation() -> None:
    prs = Presentation()

    with pytest.raises(PPTTemplateError, match="must contain at least one slide"):
        profile_presentation_template(prs)


def test_keyword_matching_avoids_english_substring_false_positives() -> None:
    text = "department partnership discovery recovery coverage"

    assert not _has_keyword_match(text, ("part", "cover"))
    assert _has_keyword_match("Part 1: Market landscape", ("part",))
    assert _has_keyword_match("封面设计", ("封面",))


def test_profile_uses_global_assignment_when_cover_and_catalog_compete() -> None:
    prs = Presentation()
    _add_text_slide(prs, ["Agenda", "1. Market", "2. Product", "3. Finance"])
    _add_title_slide(prs, "Title", "Annual Business Review")
    _add_text_slide(prs, ["Chapter 1", "Market Landscape"])
    _add_text_slide(
        prs,
        [
            "Market analysis",
            "Revenue grew 18 percent year over year.",
            "Enterprise demand remains strongest in regulated industries.",
            "Retention stays above target.",
        ],
    )
    _add_text_slide(prs, ["Thank You", "Questions"])

    profile = profile_presentation_template(prs)

    assert profile.role_index(TemplateRole.CATALOG) == 0
    assert profile.role_index(TemplateRole.COVER) == 1
    assert profile.status == "ready"


def test_profile_accepts_reloaded_uploaded_bytes() -> None:
    prs = _ordinary_template()
    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)

    reloaded = Presentation(stream)
    profile = profile_presentation_template(reloaded)

    assert profile.status == "ready"
    assert {assignment.role for assignment in profile.assignments} == {
        TemplateRole.COVER,
        TemplateRole.CATALOG,
        TemplateRole.CHAPTER,
        TemplateRole.CONTENT,
        TemplateRole.END,
    }
