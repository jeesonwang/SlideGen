import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from slidegen.services.presentation.post_render_validator import PostRenderValidator


def test_empty_slide_passes():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    validator = PostRenderValidator()
    issues = validator.validate(prs)
    assert len(issues) == 0


def test_in_bounds_shape_passes():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
    validator = PostRenderValidator()
    issues = validator.validate(prs)
    out_of_bounds = [i for i in issues if "boundaries" in i.message.lower()]
    assert len(out_of_bounds) == 0


def test_tiny_font_detected():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.paragraphs[0].font.size = Pt(6)
    validator = PostRenderValidator()
    issues = validator.validate(prs)
    readability = [i for i in issues if "smaller than 8pt" in i.message.lower()]
    assert len(readability) >= 1


def test_overlapping_text_shapes_detected():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    first.text = "First"
    second = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(4), Inches(2))
    second.text = "Second"

    validator = PostRenderValidator(mode="fail")
    issues = validator.validate(prs)

    overlaps = [i for i in issues if "overlap" in i.message.lower()]
    assert len(overlaps) == 1
    assert overlaps[0].level == "error"


def test_mode_off_returns_empty():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.paragraphs[0].font.size = Pt(6)
    validator = PostRenderValidator(mode="off")
    issues = validator.validate(prs)
    assert len(issues) == 0
