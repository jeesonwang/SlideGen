from pathlib import Path

import pytest
from pptx import Presentation

from slidegen.schemas.gen_request import ExportFormat
from slidegen.services.presentation.design_tokens import DEFAULT_TOKENS
from slidegen.services.presentation.generator import PresentationGenerator
from slidegen.services.presentation.semantic import BlockKind, BlockSpec, SlideKind, SlideSpec
from slidegen.services.presentation.visual_plan import (
    VisualObjectKind,
    VisualPlanRenderer,
    build_composite_points_visual_plan,
)


@pytest.mark.anyio
async def test_composite_points_visual_plan_renders_layered_objects(tmp_path: Path) -> None:
    spec = SlideSpec(
        kind=SlideKind.CONTENT_POINTS,
        title="生成式排版 MVP",
        source_level=3,
        blocks=(
            BlockSpec(kind=BlockKind.POINT, title="对象树", text="一个观点由背景、编号、标题、正文和装饰线叠加。"),
            BlockSpec(kind=BlockKind.POINT, title="可校验", text="每个对象都有坐标和层级，可以在渲染前检查越界。"),
        ),
    )
    plan = build_composite_points_visual_plan(spec, DEFAULT_TOKENS)

    block_zero_objects = [obj for obj in plan.objects if obj.group_id == "point_0"]
    assert len(block_zero_objects) >= 5
    assert {obj.kind for obj in block_zero_objects} >= {VisualObjectKind.SHAPE, VisualObjectKind.TEXT}

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    await VisualPlanRenderer(DEFAULT_TOKENS).render(slide, plan, spec)

    output_path = tmp_path / "visual_plan_mvp.pptx"
    prs.save(output_path)

    generated = Presentation(str(output_path))
    slide_text = "\n".join(shape.text for shape in generated.slides[0].shapes if shape.has_text_frame)
    assert "生成式排版 MVP" in slide_text
    assert "对象树" in slide_text
    assert "可校验" in slide_text
    assert len(generated.slides[0].shapes) >= len(plan.objects)


@pytest.mark.anyio
async def test_report_pipeline_can_opt_into_visual_plan_mvp(tmp_path: Path) -> None:
    markdown = Path("test/data/report.md").read_text(encoding="utf-8")
    output_path = tmp_path / "report_visual_plan_mvp.pptx"

    result = await PresentationGenerator(enable_visual_plan_mvp=True).generate_from_markdown(
        markdown_content=markdown,
        template_name="general",
        output_path=str(output_path),
        export_as=ExportFormat.PPTX,
    )

    assert result == str(output_path)
    assert output_path.exists()
    assert output_path.stat().st_size > 0

    generated = Presentation(str(output_path))
    assert len(generated.slides) == 21
    shape_counts = [len(slide.shapes) for slide in generated.slides]
    assert max(shape_counts) >= 15

    all_text = "\n".join(shape.text for slide in generated.slides for shape in slide.shapes if shape.has_text_frame)
    assert "抖音&快手 竞品分析报告" in all_text
    assert "一、产品概述" in all_text
