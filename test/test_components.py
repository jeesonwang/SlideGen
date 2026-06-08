import json
import os
import sys

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.append(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slidegen"))

from pptx import Presentation

from slidegen.services.presentation.components import ComponentsManager, components_manager


def test_components_manager():
    print(f"Available layout types: {components_manager.layout_types_names}")
    layout = components_manager.get_layout_type("two_points")
    if layout:
        print(f"Available styles for two_points: {layout.style_names}")
        print(f"Available styles for two_points: {layout.style_list}")
    assert "two_points" in components_manager.layout_types_names
    random_style = components_manager.get_random_style("two_points")
    if random_style:
        print(f"Randomly selected style: {random_style.name}")

        for shape_name, shape in random_style.shapes.items():
            print(f"  Shape: {shape_name}, Content type: {shape.content_type}")


def test_components_manager_loads_metadata_without_layout_pollution(tmp_path):
    json_path = tmp_path / "shapes.json"
    json_path.write_text(
        """
        {
          "metadata": {
            "slide_width": 12192000,
            "slide_height": 6858000
          },
          "one_point": {
            "style0": {}
          }
        }
        """,
        encoding="utf-8",
    )

    manager = ComponentsManager(json_path)

    assert manager.metadata == {"slide_width": 12192000, "slide_height": 6858000}
    assert "metadata" not in manager.layout_types_names
    assert "one_point" in manager.layout_types_names


def test_components_manager_loads_catalog_items_without_layout_pollution(tmp_path):
    json_path = tmp_path / "shapes.json"
    json_path.write_text(
        json.dumps(
            {
                "metadata": {"slide_width": 12192000, "slide_height": 6858000},
                "catalog_items": {
                    "general": {
                        "items": [
                            {
                                "number": {
                                    "xml": "number-xml",
                                    "zorder": 0,
                                    "text": "01",
                                    "location": {"x": 1, "y": 2, "width": 3, "height": 4},
                                },
                                "text": {
                                    "xml": "text-xml",
                                    "zorder": 1,
                                    "text": "Title",
                                    "location": {"x": 5, "y": 6, "width": 7, "height": 8},
                                },
                            }
                        ]
                    }
                },
                "one_point": {"style0": {}},
            }
        ),
        encoding="utf-8",
    )

    manager = ComponentsManager(json_path)

    catalog_items = manager.get_catalog_items("general")
    assert catalog_items is not None
    assert len(catalog_items) == 1
    assert catalog_items[0].number.text == "01"
    assert catalog_items[0].text.location.x == 5
    assert "catalog_items" not in manager.layout_types_names
    assert "one_point" in manager.layout_types_names


def test_add_style():
    path = "data/深度学习原理架构与应用.pptx"
    presentation = Presentation(path)
    slide = presentation.slides[16]
    components_manager.add_style_from_slide(slide, "one_point", "style1")
    components_manager.save_to_json("components/shapes/shapes.json")
    components_manager.reload("components/shapes/shapes.json")

    layout = components_manager.get_layout_type("one_points")
    assert "style1" in layout.style_names
    print(f"Available styles for one_points: {layout.style_names}")
    print(f"Available styles for one_points: {layout.style_list}")
