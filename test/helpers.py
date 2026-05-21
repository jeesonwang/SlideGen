from pptx import Presentation
from pptx.util import Inches


def add_catalog_slide(prs: Presentation, *, titles: tuple[str, ...] = ("Market", "Product", "Finance")) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.5)).text = "Agenda"
    for index, title in enumerate(titles, start=1):
        top = 1.0 + index * 0.5
        slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(0.5), Inches(0.3)).text = str(index)
        slide.shapes.add_textbox(Inches(1.8), Inches(top), Inches(6.0), Inches(0.3)).text = title
