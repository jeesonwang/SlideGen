import pytest
from pptx import Presentation

from slidegen.services.presentation.slide_renderer import SlideRenderer, AssetProvider
from slidegen.services.presentation.design_tokens import DEFAULT_TOKENS
from slidegen.services.presentation.default_recipes import (
    agenda_recipe,
    cover_recipe,
    grid_cards_recipe,
    title_body_recipe,
    two_column_recipe,
)
from slidegen.services.presentation.recipes import LayoutRecipe
from slidegen.services.presentation.region import Region, RegionRole, RepeatRule
from slidegen.services.presentation.semantic import SlideSpec, SlideKind, BlockSpec, BlockKind


def _make_spec(texts: list[str]) -> SlideSpec:
    blocks = tuple(BlockSpec(kind=BlockKind.POINT, title=t[:20], text=t) for t in texts)
    return SlideSpec(kind=SlideKind.CONTENT_POINTS, title="Test Slide", source_level=2, blocks=blocks)


@pytest.mark.anyio
async def test_renderer_creates_shapes_for_title_body():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = _make_spec(["Body content here"])
    recipe = title_body_recipe(DEFAULT_TOKENS)
    renderer = SlideRenderer(DEFAULT_TOKENS)

    await renderer.render(slide, recipe, spec)

    assert len(slide.shapes) >= 2
    texts = [s.text for s in slide.shapes if s.has_text_frame]
    assert "Test Slide" in texts
    assert "Body content here" in texts


@pytest.mark.anyio
async def test_renderer_grid_cards_4_blocks():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = _make_spec(["Block A", "Block B", "Block C", "Block D"])
    recipe = grid_cards_recipe(DEFAULT_TOKENS, n_blocks=4)
    renderer = SlideRenderer(DEFAULT_TOKENS)

    await renderer.render(slide, recipe, spec)

    assert len(slide.shapes) >= 5  # title + 4 card bodies


@pytest.mark.anyio
async def test_renderer_maps_text_for_cover_two_column_and_agenda():
    prs = Presentation()
    renderer = SlideRenderer(DEFAULT_TOKENS)

    cover = SlideSpec(
        kind=SlideKind.COVER,
        title="Quarterly Review",
        source_level=1,
        blocks=(BlockSpec(kind=BlockKind.SUBTITLE, title="", text="Q4 Executive Summary"),),
    )
    cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
    await renderer.render(cover_slide, cover_recipe(DEFAULT_TOKENS), cover)
    cover_texts = [s.text for s in cover_slide.shapes if s.has_text_frame]
    assert "Quarterly Review" in cover_texts
    assert "Q4 Executive Summary" in cover_texts

    comparison = SlideSpec(
        kind=SlideKind.COMPARISON,
        title="Build vs Buy",
        source_level=3,
        blocks=(
            BlockSpec(kind=BlockKind.POINT, title="Build", text="Control roadmap"),
            BlockSpec(kind=BlockKind.POINT, title="Buy", text="Faster launch"),
        ),
    )
    comparison_slide = prs.slides.add_slide(prs.slide_layouts[6])
    await renderer.render(comparison_slide, two_column_recipe(DEFAULT_TOKENS), comparison)
    comparison_text = "\n".join(s.text for s in comparison_slide.shapes if s.has_text_frame)
    assert "Build\x0bControl roadmap" in comparison_text
    assert "Buy\x0bFaster launch" in comparison_text

    agenda = SlideSpec(
        kind=SlideKind.AGENDA,
        title="Agenda",
        source_level=1,
        blocks=(
            BlockSpec(kind=BlockKind.POINT, title="Market", text="Trends"),
            BlockSpec(kind=BlockKind.POINT, title="Plan", text="Next steps"),
        ),
    )
    agenda_slide = prs.slides.add_slide(prs.slide_layouts[6])
    await renderer.render(agenda_slide, agenda_recipe(DEFAULT_TOKENS, n_blocks=2), agenda)
    agenda_text = "\n".join(s.text for s in agenda_slide.shapes if s.has_text_frame)
    assert "01" in agenda_text
    assert "02" in agenda_text
    assert "Agenda" in agenda_text
    assert "Market\x0bTrends" in agenda_text
    assert "Plan\x0bNext steps" in agenda_text


@pytest.mark.anyio
async def test_renderer_applies_repeat_rule_bindings_to_expanded_regions():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = _make_spec(["First body", "Second body"])
    recipe = LayoutRecipe(
        name="RepeatedBodyRecipe",
        regions=(
            Region(region_id="title", x_frac=0.08, y_frac=0.05, w_frac=0.84, h_frac=0.10),
        ),
        repeats=(
            RepeatRule(
                seed=Region(region_id="item", x_frac=0.08, y_frac=0.22, w_frac=0.84, h_frac=0.12),
                step_x=0.0,
                step_y=0.15,
                role=RegionRole.BODY,
            ),
        ),
        region_roles={"title": RegionRole.TITLE},
        region_block_indexes={"item": 0},
        region_text_sources={"title": "slide_title", "item": "block_text"},
    )

    await SlideRenderer(DEFAULT_TOKENS).render(slide, recipe, spec)

    texts = [s.text for s in slide.shapes if s.has_text_frame]
    assert "Test Slide" in texts
    assert "First body" in texts
    assert "Second body" in texts


@pytest.mark.anyio
async def test_slide_renderer_full_loop_generates_valid_pptx(tmp_path):
    """集成测试：SlideSpec → PresetRecipe → SlideRenderer → PPTX"""
    from slidegen.services.presentation.preset_recipes import PresetRecipeFallback

    prs = Presentation()
    fallback = PresetRecipeFallback()

    # Slide 1: Cover
    cover_spec = SlideSpec(
        kind=SlideKind.COVER, title="My Deck", source_level=1,
        blocks=(BlockSpec(kind=BlockKind.TITLE, title="My Deck", text="Subtitle here"),)
    )
    cover_recipe = fallback.select(cover_spec, DEFAULT_TOKENS)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    await SlideRenderer(DEFAULT_TOKENS).render(slide1, cover_recipe, cover_spec)

    # Slide 2: Content
    content_spec = _make_spec(["Point A text", "Point B text", "Point C text", "Point D text"])
    content_recipe = fallback.select(content_spec, DEFAULT_TOKENS)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    await SlideRenderer(DEFAULT_TOKENS).render(slide2, content_recipe, content_spec)

    # 验证 PPTX 可保存和读取
    output = tmp_path / "test_output.pptx"
    prs.save(str(output))
    assert output.exists()

    reloaded = Presentation(str(output))
    assert len(reloaded.slides) == 2
    slide1_texts = [s.text for s in reloaded.slides[0].shapes if s.has_text_frame]
    assert "My Deck" in slide1_texts
