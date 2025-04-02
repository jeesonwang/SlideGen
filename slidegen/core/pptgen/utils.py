import re
from enum import Enum
from typing import Any, Optional
from types import SimpleNamespace
from copy import deepcopy

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.text.text import _Paragraph, _Run
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from exception.custom_exception import PPTGenError

FONT_SIZE_CASE = {
    PP_PLACEHOLDER.TITLE: 54,
    PP_PLACEHOLDER.SUBTITLE: 40,
    PP_PLACEHOLDER.BODY: 14,
    PP_PLACEHOLDER.FOOTER: 24,
    PP_PLACEHOLDER.HEADER: 32,
    PP_PLACEHOLDER.CHART: 12,
    PP_PLACEHOLDER.MIXED: 18
}

# scale factor for emu to pt
SCALE_FACTOR = 12700
SCALE_FACTOR_CH = 16000

class CatalogLayout(Enum):
    """
    Catalog layout enum
    """
    VERTICAL = "vertical"
    HORIZONTAL = "horizontal"
    UNDEFINED = "undefined"


class CatalogItem:
    """
    Catalog item including number shape, text shape and background shape.
    """
    def __init__(self, number_shape: dict, text_shape: dict, background_shape: dict | None = None):
        self.number_shape = number_shape
        self.text_shape = text_shape
        self.background_shape = background_shape

    def asdict(self):
        return {
            "number_shape": self.number_shape,
            "text_shape": self.text_shape,
            "background_shape": self.background_shape
        }

class CatalogList(list):
    """
    Catalog list including a list of `CatalogItem`.
    """
    def asdict(self):
        return [item.asdict() for item in self]

def is_english(texts):
    eng = 0
    if not texts:
        return False
    for t in texts:
        if re.match(r"[ `a-zA-Z.,':;/\"?<>!\(\)-]", t.strip()):
            eng += 1
    if eng / len(texts) > 0.8:
        return True
    return False

def is_chinese(text):
    if not text:
        return False
    chinese = 0
    for ch in text:
        if '\u4e00' <= ch <= '\u9fff':
            chinese += 1
    if chinese / len(text) > 0.2:
        return True
    return False

def get_font_style(font: dict[str, Any]) -> str:
    """
    Convert a font dictionary to a CSS style string.

    Args:
        font (Dict[str, Any]): The font dictionary.

    Returns:
        str: The CSS style string.
    """
    font = SimpleNamespace(**font)
    styles = []

    if hasattr(font, "size") and font.size:
        styles.append(f"font-size: {font.size}pt")

    if hasattr(font, "color") and font.color:
        if all(c in "0123456789abcdefABCDEF" for c in font.color):
            styles.append(f"color: #{font.color}")
        else:
            styles.append(f"color: {font.color}")

    if hasattr(font, "bold") and font.bold:
        styles.append("font-weight: bold")

    if hasattr(font, "italic") and font.italic:
        styles.append("font-style: italic")

    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph) -> Optional[_Run]:
    """
    Merge all runs in a paragraph into a single run.

    Args:
        paragraph (_Paragraph): The paragraph to merge runs in.

    Returns:
        Optional[_Run]: The merged run, or None if there are no runs.
    """
    runs = paragraph.runs

    # Handle field codes
    if len(runs) == 0:
        runs = tuple(
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        )
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None

    # Find the run with the most text
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    # Remove other runs
    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def del_para(paragraph_id: int, shape: BaseShape):
    """
    Delete a paragraph from a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    para._element.getparent().remove(para._element)

def add_para_by_xml(shape: Shape, xml: str):
    """
    Add a paragraph in a text frame by XML string.
    """
    if not shape.has_text_frame:
        raise PPTGenError("Shape does not have a text frame.")
    shape.text_frame.clear()
    shape.text_frame.add_paragraph()
    shape.text_frame.paragraphs[0]._element.addnext(parse_xml(xml))
    del_para(0, shape)
    return shape

def get_theme_colors(presentation):
    theme_part = presentation.slide_master.part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)
    color_elements = theme.xpath('a:themeElements/a:clrScheme/*')
    result = {}
    for element in color_elements:
        namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        root = etree.fromstring(etree.tostring(element))
        theme_name = element.tag.replace('{http://schemas.openxmlformats.org/drawingml/2006/main}', '')
        result[theme_name] = root.xpath('//a:srgbClr/@val', namespaces=namespaces)[0]
    return result

def modify_shape_xml(xml_str: str, shape_id: int | str, shape_name: str, text_content: str) -> str:
    """
    Modify the XML of a PPTX shape: update the shape ID, name, and text content.
    
    Args:
        xml_str (str): The input XML string.
        shape_id (int | str): The new shape ID.
        shape_name (str): The new shape name.
        text_content (str): The new text content.
    
    Returns:
        str: The modified XML string.
    """
    root = etree.fromstring(xml_str)
    nsmap = root.nsmap
    
    cNvPr = root.find('.//p:cNvPr', namespaces=nsmap)
    if cNvPr is not None:
        cNvPr.set('id', str(shape_id))
        cNvPr.set('name', shape_name)
    
    t_element = root.find('.//a:t', namespaces=nsmap)
    if t_element is not None:
        # Keep the original rPr format
        r_element = t_element.getparent()
        if r_element is not None:
            r_pr = r_element.find('.//a:rPr', namespaces=nsmap)
            if r_pr is not None:
                new_r = etree.Element(f"{{{nsmap['a']}}}r")
                new_r.append(deepcopy(r_pr))
                new_t = etree.Element(f"{{{nsmap['a']}}}t")
                new_t.text = text_content
                new_r.append(new_t)
                
                p_element = r_element.getparent()
                p_element.replace(r_element, new_r)
    return etree.tostring(root, encoding='unicode', pretty_print=True)

def add_shape_by_xml(slide: Slide, shape_id: int, shape_name: str, text_content: str, shape_xml: str):
    """
    Add a shape by XML string.

    Args:
        slide (Slide): The slide to add the shape to.
        shape_id (int): The ID of the shape.
        shape_name (str): The name of the shape.
        text_content (str): The text content of the shape.
        shape_xml (str): The XML string of the shape.

    Returns:
        Shape: The added shape.
    """ 
    shape_xml = modify_shape_xml(shape_xml, shape_id, shape_name, text_content)

    new_shape = slide.shapes._shape_factory(
    slide.shapes._spTree.insert_element_before(parse_xml(shape_xml), 'p:extLst'))
    return new_shape

def clone_para(paragraph_id: int, shape: BaseShape):
    """
    Clone a paragraph in a shape.
    """
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))

def convert_paragraph_xml(paragraph_xml: str, text_content: str) -> str:
    """
    Convert paragraph xml to add text content and keep the original paragraph style.
    Args:
        paragraph_xml (str): The paragraph xml. It usually comes from an empty text shape.
        text_content (str): The text content to add.

    Returns:
        str: The converted paragraph xml.
    """
    root = etree.fromstring(paragraph_xml)
    drawingml_ns = root.nsmap.get('a')
    if root.tag == f"{{{drawingml_ns}}}p":
        p_element = root
    else:
        p_element = root.find(".//a:p", namespaces=root.nsmap)
    if p_element is None:
        return etree.tostring(root, encoding="unicode", pretty_print=True)
    
    end_para_rpr = p_element.find(".//a:endParaRPr", namespaces=root.nsmap)
    if end_para_rpr is not None:
        r_pr = etree.Element(f"{{{drawingml_ns}}}rPr")
        for attr, value in end_para_rpr.attrib.items():
            r_pr.set(attr, value)
        for child in end_para_rpr:
            r_pr.append(deepcopy(child))
        r_element = etree.Element(f"{{{drawingml_ns}}}r")
        r_element.append(r_pr)
        t_element = etree.Element(f"{{{drawingml_ns}}}t")
        t_element.text = text_content
        r_element.append(t_element)
        p_element.remove(end_para_rpr)

        p_pr = p_element.find(".//a:pPr", namespaces=root.nsmap)
        if p_pr is not None:
            p_pr.addnext(r_element)
        else:
            p_element.insert(0, r_element)

    return etree.tostring(root, encoding="unicode", pretty_print=True)


def remove_custDataLst(xml_str: str) -> str:
    """
    Remove the <p:custDataLst> part in the XML and return the processed XML string.
    
    Args:
        xml_str (str): The input XML string, containing the <p:custDataLst> part.
        
    Returns:
        str: The processed XML string, without the <p:custDataLst> part.
    """
    root = etree.fromstring(xml_str)
    ns = root.nsmap
    
    cust_data_list = root.find(".//p:custDataLst", namespaces=ns)
    
    if cust_data_list is not None:
        parent = cust_data_list.getparent()
        if parent is not None:
            parent.remove(cust_data_list)
    
    return etree.tostring(
        root,
        encoding="unicode",
        pretty_print=True,
        xml_declaration=False,
    )
