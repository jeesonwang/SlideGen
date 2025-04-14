import copy
import io
import os
import random
from enum import Enum
from typing import Any

from loguru import logger
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from slidegen.core.docparse.markdown_parser import Heading
from slidegen.core.pptgen.components import ChapterLayout, ContentType, components_manager
from slidegen.core.pptgen.utils import (
    add_para_by_xml,
    add_shape_by_xml,
    convert_paragraph_xml,
    is_image_path,
    runs_merge,
)
from slidegen.exception import PPTGenError, PPTTemplateError


class Page:
    """PPT pages base class"""

    @staticmethod
    def _set_text(shape: Shape, text: str) -> None:
        """Set the text of the shape, keep the original paragraph style."""
        assert shape.has_text_frame, "Shape must have a text frame"
        if shape.is_placeholder:
            shape.text = text
            return
        tf = shape.text_frame
        # if the shape has text, merge the runs and set the run text
        if shape.text:
            para = tf.paragraphs[0]
            run = runs_merge(para)
            if run:
                run.text = text
        # if the shape has no text, add a new paragraph and keep the original paragraph style
        else:
            para = tf.paragraphs[0]
            para_xml = convert_paragraph_xml(para._element.xml, text)
            shape = add_para_by_xml(shape, para_xml)

    @staticmethod
    def remove_slide(prs: Presentation, index: int) -> None:
        """Delete the slide at the given index"""
        rId = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rId)
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(list(xml_slides)[index])

    @staticmethod
    def remove_shapes(sp_tree: CT_GroupShape, shapes: list[Shape]) -> None:
        for shp in shapes:
            if shp:
                sp_tree.remove(shp.element)

    @staticmethod
    def _set_text_style(shape: Shape, style: dict[str, Any]) -> None:
        """Set the input `Shape` text style"""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        paragraph = tf.paragraphs[0]
        for key, value in style.items():
            if value is None:
                continue
            setattr(paragraph.font, key, value)

    @staticmethod
    def move_slide(pres: Presentation, slide: Slide, index: int) -> None:
        """
        Move the slide to the new index

        Args:
            pres: Presentation object
            slide: Slide object to be moved
            index: New index of the slide
        """
        old_index = pres.slides.index(slide)
        xml_slides = pres.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(index, slides[old_index])

    @staticmethod
    def duplicate_slide(pres: Presentation, index: int) -> Slide:
        """Duplicate the slide at the given index

        Args:
            pres: Presentation object
            index: Index of the slide to be duplicated

        Returns:
            Copied slide
        """
        template = pres.slides[index]
        copied_slide = pres.slides.add_slide(template.slide_layout)
        # Delete the existing shapes that are part of the layout
        for shp in copied_slide.shapes:
            copied_slide.shapes.element.remove(shp.element)

        # Perform a deep copy of the shapes from the template
        for shp in template.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = io.BytesIO(shp.image.blob)
                copied_slide.shapes.add_picture(
                    image_file=img,
                    left=shp.left,
                    top=shp.top,
                    width=shp.width,
                    height=shp.height,
                )
            else:
                el = shp.element
                newel = copy.deepcopy(el)
                custDataLst = newel.xpath(".//p:custDataLst")
                for cd in custDataLst:
                    cd.getparent().remove(cd)
                copied_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

        return copied_slide


class CoverPage(Page):
    """Presentation cover page"""

    @staticmethod
    def generate_slide(prs: Presentation, content: Heading, *, cover_page_index: int = 0) -> None:
        """
        Generate the cover page

        Args:
            prs: Presentation object
            content: Heading object, the main heading of the markdown document(level 1)
            cover_page_index: index of the cover page
        """
        cover_page = prs.slides[cover_page_index]
        assert content.level == 1, "Cover page must have a level 1 heading"
        main_title = content.element_text
        if not main_title.strip():
            main_title = "Presentation Title"
        title_found = False
        # TODO: add subtitle
        for placeholder in cover_page.shapes.placeholders:
            placeholder_type = placeholder.placeholder_format.type
            if placeholder_type == PP_PLACEHOLDER.TITLE or placeholder_type == PP_PLACEHOLDER.CENTER_TITLE:
                CoverPage._set_text(placeholder, main_title)
                placeholder.text_frame.word_wrap = False
                title_found = True
                break
        if not title_found:
            raise PPTTemplateError(f"{CoverPage.__name__}: No title placeholder found in cover slide")


class CatalogLayout(Enum):
    """
    Catalog layout enum
    """

    VERTICAL = "vertical"
    HORIZONTAL = "horizontal"
    UNDEFINED = "undefined"


class CatalogItem:
    """
    Catalog item including number shape, text shape and background shape.
    """

    def __init__(
        self, number_shape: dict[str, Any], text_shape: dict[str, Any], background_shape: dict[str, Any] | None = None
    ):
        self.number_shape = number_shape
        self.text_shape = text_shape
        self.background_shape = background_shape

    def asdict(self) -> dict[str, Any]:
        return {
            "number_shape": self.number_shape,
            "text_shape": self.text_shape,
            "background_shape": self.background_shape,
        }


class CatalogList(list[CatalogItem]):
    """
    Catalog list including a list of `CatalogItem`.
    """

    def asdict(self) -> list[dict[str, Any]]:
        return [item.asdict() for item in self]


class CatalogPage(Page):
    """Presentation catalog page"""

    # vertical tolerance coefficient
    vertical_tolerance = 1.5

    @staticmethod
    def _calculate_distance(shape1: dict[str, Any], shape2: dict[str, Any]) -> float:
        """Calculate the distance between two shapes"""
        return ((shape1["left"] - shape2["left"]) ** 2 + (shape1["top"] - shape2["top"]) ** 2) ** 0.5

    @staticmethod
    def _layout_direction(number_shapes: list[dict[str, Any]]) -> CatalogLayout:
        """Judge the layout direction of the catalog page"""
        if len(number_shapes) < 2:
            raise PPTTemplateError("To judge the layout direction, catalog page must have at least two chapter numbers")

        sorted_numbers = sorted(number_shapes, key=lambda x: (x["left"], x["top"]))
        horizontal_diffs = []
        vertical_diffs = []

        for i in range(len(sorted_numbers) - 1):
            horizontal_diffs.append(abs(sorted_numbers[i + 1]["left"] - sorted_numbers[i]["left"]))
            vertical_diffs.append(abs(sorted_numbers[i + 1]["top"] - sorted_numbers[i]["top"]))
        avg_horizontal_diff = sum(horizontal_diffs) / len(horizontal_diffs) if horizontal_diffs else 0
        avg_vertical_diff = sum(vertical_diffs) / len(vertical_diffs) if vertical_diffs else 0

        if avg_horizontal_diff > avg_vertical_diff:
            return CatalogLayout.HORIZONTAL
        else:
            return CatalogLayout.VERTICAL

    @staticmethod
    def _get_catalog_items(slide: Slide) -> CatalogList:
        number_shapes: list[dict[str, Any]] = []
        text_shapes: list[dict[str, Any]] = []
        all_shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                # placeholder shapes are not included in the all_shapes list
                continue
            shape_info: dict[str, Any] = {
                "text": shape.text.strip() if shape.has_text_frame else None,  # type: ignore
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "shape_type": shape.shape_type,
                "shape_id": shape.shape_id,
                "shape": shape,
            }
            if shape.has_text_frame:
                text_shapes.append(shape_info)
            all_shapes.append(shape_info)
        for shape_info in text_shapes:
            # check if the text is a chapter number
            text = shape_info["text"].strip()  # type: ignore
            if len(text) > 3:
                continue
            if text.isdigit() or (text.endswith(".") and text[:-1].isdigit()):
                # TODO: Optimize judgment conditions
                if int(text.replace(".", "")) > 49:
                    continue
                number_shapes.append(shape_info)
        try:
            number_shapes.sort(key=lambda shape: int(shape["text"]))  # type: ignore
        except ValueError:
            raise PPTTemplateError("Chapter number must be a number")

        match len(number_shapes):
            case 0:
                raise PPTTemplateError("Catalog page must have at least one chapter numbers")
            case 1:
                layout_direction = CatalogLayout.UNDEFINED
            case _:
                layout_direction = CatalogPage._layout_direction(number_shapes)

        except_number_shapes = [shape for shape in text_shapes if shape not in number_shapes]
        catalog_list = CatalogList()
        # Find the closest text shape for each number shape
        for number_shape in number_shapes:
            min_distance = float("inf")
            closest_text_shape = None

            for text_shape in except_number_shapes:
                if layout_direction == CatalogLayout.HORIZONTAL:
                    # For horizontal layout, find the text shape below the number shape
                    if text_shape["top"] > number_shape["top"]:
                        distance = CatalogPage._calculate_distance(number_shape, text_shape)
                        horizontal_overlap = min(
                            number_shape["left"] + number_shape["width"],
                            text_shape["left"] + text_shape["width"],
                        ) - max(number_shape["left"], text_shape["left"])
                        if horizontal_overlap > 0 and distance < min_distance:
                            min_distance = distance
                            closest_text_shape = text_shape
                elif layout_direction == CatalogLayout.VERTICAL:
                    # For vertical layout, find the text shape to the right of the number shape
                    if text_shape["left"] > number_shape["left"]:
                        distance = CatalogPage._calculate_distance(number_shape, text_shape)
                        vertical_overlap = min(
                            number_shape["top"] + number_shape["height"],
                            text_shape["top"] + text_shape["height"],
                        ) - max(number_shape["top"], text_shape["top"])
                        if vertical_overlap > 0 and distance < min_distance:
                            min_distance = distance
                            closest_text_shape = text_shape
                else:
                    distance = CatalogPage._calculate_distance(number_shape, text_shape)
                    if distance < min_distance:
                        min_distance = distance
                        closest_text_shape = text_shape

            if closest_text_shape:
                catalog_list.append(CatalogItem(number_shape, closest_text_shape))
                try:
                    all_shapes.remove(number_shape)
                    if closest_text_shape is not None:
                        all_shapes.remove(closest_text_shape)
                except ValueError:
                    raise PPTTemplateError(
                        f"all shape: {all_shapes}\n current closest_text_shape: {closest_text_shape} \n current number_shape: {number_shape}"
                    )
        assert len(number_shapes) == len(catalog_list), (
            "The number of chapter numbers and chapter titles must be the same"
        )

        if len(all_shapes) >= len(number_shapes):
            # continue to calculate the distance to find the background shape
            for i, number_shape in enumerate(number_shapes):
                min_distance = float("inf")
                closest_background_shape = None
                for shape_info in all_shapes:
                    distance = CatalogPage._calculate_distance(number_shape, shape_info)
                    if distance < min_distance:
                        min_distance = distance
                        if min_distance < shape_info["height"] * CatalogPage.vertical_tolerance:
                            closest_background_shape = shape_info
                if closest_background_shape:
                    catalog_list[i].background_shape = closest_background_shape

        return catalog_list

    @staticmethod
    def generate_slide(
        prs: Presentation,
        content: list[Heading],
        *,
        catalog_page_index: int = 1,
        begin_number: int = 1,
    ) -> int:
        """
        Generate the catalog page

        Args:
            prs: Presentation object
            content: list of Heading objects
            catalog_page_index: index of the catalog page
            begin_number: starting number of the catalog page
        Returns:
            index of the catalog page
        """
        if not content:
            raise PPTGenError("Catalog page must have content.")
        catalog_num = len(content)
        catalog_slide = prs.slides[catalog_page_index]
        catalog_items = CatalogPage._get_catalog_items(catalog_slide)

        if len(catalog_items) > catalog_num:
            sp_tree = catalog_slide.shapes._spTree
            # delete the excess shape pairs from the slide
            excess_items = catalog_items[catalog_num:]
            for item in excess_items:
                # delete the excess chapter number and chapter title shapes
                Page.remove_shapes(
                    sp_tree,
                    [
                        item.number_shape["shape"],
                        item.text_shape["shape"],
                        *([item.background_shape["shape"]] if item.background_shape else []),
                    ],
                )

            catalog_items = catalog_items[:catalog_num]  # type: ignore
        # TODO: Add catalog items to the slide
        for i in range(len(catalog_items)):
            cur_content = content[i].element_text
            cur_number = begin_number
            text_shape = catalog_items[i].text_shape["shape"]
            number_shape = catalog_items[i].number_shape["shape"]

            catalog_items[i].text_shape["text"] = cur_content
            CatalogPage._set_text(text_shape, cur_content)
            # The chapter number is formatted as "01"
            chapter_number = str(cur_number).zfill(2)
            catalog_items[i].number_shape["text"] = chapter_number
            CatalogPage._set_text(number_shape, chapter_number)
            begin_number += 1
        # If the number of catalog_num is less than the number of catalog_items, generate a new catalog page
        if begin_number - 1 < catalog_num:
            new_catalog_slide = CatalogPage.duplicate_slide(prs, catalog_page_index)
            catalog_page_index += 1
            CatalogPage.move_slide(prs, new_catalog_slide, catalog_page_index)
            # Recursively generate the new catalog page
            CatalogPage.generate_slide(
                prs,
                content[len(catalog_items) :],
                catalog_page_index=catalog_page_index,
                begin_number=begin_number,
            )

        return catalog_page_index


class ChapterHomePage(Page):
    """Chapter home page"""

    selected_style: int | None = None

    @staticmethod
    def generate_slide(
        prs: Presentation,
        content: Heading,
        *,
        chapter_home_page_index: int = 2,
        chapter_number: int = 1,
        slide_index: int = 2,
    ) -> None:
        """
        Generate the chapter home page

        Args:
            prs: Presentation object
            content: Heading object
            chapter_home_page_index: index of the chapter home page
            chapter_number: current chapter number, begin from 1
            slide_index: index of the slide to be generated
        """
        assert content.level == 2, f"{ChapterHomePage.__name__}: Chapter home page must input a level 2 heading"
        template_slide = prs.slides[chapter_home_page_index]
        chapter_home_slide = prs.slides.add_slide(template_slide.slide_layout)

        title = content.element_text
        title_placeholder = None

        for placeholder in chapter_home_slide.shapes.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                ChapterHomePage._set_text(placeholder, title)
                tf = placeholder.text_frame
                tf.word_wrap = False
                title_placeholder = placeholder
                break

        if not title_placeholder:
            raise PPTTemplateError(f"{ChapterHomePage.__name__}: No title placeholder found in chapter home slide")
        chapter_number_shape = None
        min_distance = float("inf")
        for shape in chapter_home_slide.shapes:
            if shape == title_placeholder:
                continue
            if shape.has_text_frame:
                if shape.top < title_placeholder.top:
                    distance = title_placeholder.top - shape.top
                    if distance < min_distance:
                        shape_text = shape.text.strip()
                        if (
                            (shape_text.startswith("0") and shape_text[1:].isdigit())
                            or shape_text.lower().startswith("part")
                            or (shape_text.endswith(".") and shape_text[:-1].isdigit())
                        ):
                            chapter_number_shape = shape
                            break
                        min_distance = distance
                        chapter_number_shape = shape
        if chapter_number_shape:
            chapter_index = ChapterHomePage.convert_chapter_number(chapter_number)
            ChapterHomePage._set_text(chapter_number_shape, chapter_index)
        ChapterHomePage.move_slide(prs, chapter_home_slide, slide_index)

    @staticmethod
    def convert_chapter_number(chapter_number: int) -> str:
        """Randomly select a style for the chapter number"""
        import inflect

        # If the style has been selected, use it
        if ChapterHomePage.selected_style is not None:
            style_type = ChapterHomePage.selected_style
        else:
            style_type = random.randint(1, 3)
            ChapterHomePage.selected_style = style_type

        p = inflect.engine()
        if style_type == 1:
            return str(chapter_number).zfill(2)  # 01, 02, 03, ...
        elif style_type == 2:
            return f"PART {str(chapter_number).zfill(2)}"  # PART 01, PART 02, PART 03, ...
        else:
            # PART ONE, PART TWO, PART THREE, ...
            return f"PART {p.number_to_words(chapter_number).upper()}"  # type: ignore


class ChapterContentPage(Page):
    """
    Chapter content page

    Divide the chapter content slides into one-point, two-point, three-point, and four-point slides.
    """

    # add static variable to track used pictures
    used_pictures: dict[str, set[str]] = {"opaque": set(), "transparent": set()}

    @staticmethod
    def _get_picture(picture_path: str) -> str:
        """
        Get a random picture path

        Args:
            picture_type: picture type, 'opaque' or 'transparent'

        Returns:
            str: the random chosen picture path
        """
        if picture_path.endswith("opaque"):
            picture_type = "opaque"
        elif picture_path.endswith("transparent"):
            picture_type = "transparent"
        else:
            raise PPTGenError(f"{ChapterContentPage.__name__}: Invalid picture path: {picture_path}")
        picture_dir = picture_path
        available_pictures = [
            p for p in os.listdir(picture_dir) if p not in ChapterContentPage.used_pictures[picture_type]
        ]
        if not available_pictures:
            ChapterContentPage.used_pictures[picture_type] = set()
            available_pictures = os.listdir(picture_dir)

        chosen_picture = random.choice(available_pictures)
        ChapterContentPage.used_pictures[picture_type].add(chosen_picture)
        return os.path.join(picture_dir, chosen_picture)

    @staticmethod
    def _get_slide_type(content: Heading) -> int:
        """
        Get the slide type of the chapter content page
        """
        return len(content)

    @staticmethod
    def _shape_alignment(shape: BaseShape | Shape) -> None:
        """Set the alignment of the shape. Uniformly justify the text in the shape"""
        if shape.has_text_frame:
            tf = shape.text_frame  # type: ignore
            tf.vertical_anchor = MSO_ANCHOR.TOP
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY

    @staticmethod
    def generate_slide(
        prs: Presentation,
        content: Heading,
        *,
        chapter_page_index: int = 3,
        slide_index: int = 3,
    ) -> None:
        """
        Generate the chapter content page

        Args:
            prs: Presentation object
            content: Heading object
            chapter_page_index: index of the template chapter content slide
            slide_index: index of the slide to be generated
        """
        assert content.level == 2, f"{ChapterContentPage.__name__}: Chapter content page must have a level 2 heading"

        slide_type = ChapterContentPage._get_slide_type(content)
        if slide_type > 4:
            raise PPTGenError(f"{ChapterContentPage.__name__}: Invalid slide type: {slide_type}")
        titles = [child.element_text for child in content.children]  # type: ignore
        section_texts = [child.text for child in content.children]

        chapter_page = prs.slides[chapter_page_index]
        new_slide = prs.slides.add_slide(chapter_page.slide_layout)
        # set the title of the new slide
        for placeholder in new_slide.shapes.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                placeholder.text = content.element_text
                placeholder.text_frame.word_wrap = False
                break

        index = 0
        chapter_layout = ChapterLayout(slide_type)  # type: ignore
        style = components_manager.get_random_style(chapter_layout)
        logger.debug(f"{ChapterContentPage.__name__}: {chapter_layout} {style.name if style else 'None'}")

        # Sort by zorder
        sorted_shapes = sorted(style.shapes.items(), key=lambda x: x[1].zorder)

        for shape_name, shape in sorted_shapes:
            # locs must be in order
            locs = shape.location
            for idx, loc in enumerate(locs):
                if shape.content_type == ContentType.CONTENT:
                    if len(section_texts) != len(locs):
                        raise PPTGenError(
                            f"{ChapterContentPage.__name__}: \
                                          Text content must be equal to the number of locations: {len(section_texts)} != {len(locs)}"
                        )
                    added_shape = add_shape_by_xml(
                        slide=new_slide,
                        shape_xml=shape.xml,  # type: ignore
                        shape_id=index,
                        shape_name=shape_name,
                        text_content=section_texts[idx],
                        location=loc,
                    )
                    ChapterContentPage._shape_alignment(added_shape)
                elif shape.content_type == ContentType.TITLE:
                    if len(titles) != len(locs):
                        raise PPTGenError(
                            f"{ChapterContentPage.__name__}: \
                                          Title must be equal to the number of locations: {len(titles)} != {len(locs)}"
                        )
                    added_shape = add_shape_by_xml(
                        slide=new_slide,
                        shape_xml=shape.xml,  # type: ignore
                        shape_id=index,
                        shape_name=shape_name,
                        text_content=titles[idx],
                        location=loc,
                    )
                    ChapterContentPage._shape_alignment(added_shape)
                elif shape.content_type == ContentType.PICTURE:
                    if shape.path is None:
                        raise PPTGenError(f"{ChapterContentPage.__name__}: Picture path is None: {shape.path}")
                    if is_image_path(shape.path):
                        image_path = shape.path
                    else:
                        image_path = ChapterContentPage._get_picture(shape.path)
                    added_shape = new_slide.shapes.add_picture(image_path, loc.x, loc.y, loc.width, loc.height)
                elif shape.content_type == ContentType.NUMBER:
                    added_shape = add_shape_by_xml(
                        slide=new_slide,
                        shape_xml=shape.xml,  # type: ignore
                        shape_id=index,
                        shape_name=shape_name,
                        text_content=str(idx + 1).zfill(2),
                        location=loc,
                    )
                else:
                    added_shape = add_shape_by_xml(
                        slide=new_slide,
                        shape_xml=shape.xml,  # type: ignore
                        shape_id=index,
                        shape_name=shape_name,
                        location=loc,
                    )
            index += 1
        ChapterContentPage.move_slide(prs, new_slide, slide_index)


class EndPage(Page):
    """End page"""

    @staticmethod
    def generate_slide(
        prs: Presentation,
        content: Heading | None = None,
        *,
        end_page_index: int = 4,
        slide_index: int = 4,
    ) -> None:
        if content is None:
            content = Heading(text="Thank you!", level=2)
        template_slide = prs.slides[end_page_index]
        end_slide = prs.slides.add_slide(template_slide.slide_layout)
        title_found = False
        for placeholder in end_slide.shapes.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                placeholder.text = content.element_text
                placeholder.text_frame.word_wrap = False
                title_found = True
                break
        if not title_found:
            raise PPTTemplateError(
                f"{EndPage.__name__}: No title placeholder found in end slide, end slide index: {end_page_index}"
            )
        EndPage.move_slide(prs, end_slide, slide_index)
