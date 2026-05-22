from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from slidegen.services.presentation.design_tokens import DesignTokens
from slidegen.services.presentation.recipes import LayoutRecipe
from slidegen.services.presentation.region import Region, RegionRole, RepeatRule
from slidegen.services.presentation.semantic import SlideSpec, BlockKind
from slidegen.services.presentation.image_generator import ImageGenerator
from slidegen.services.presentation.icon_searcher import IconSearcher
from slidegen.schemas.image_prompt import ImagePrompt


class AssetProvider:
    """Interface for image generation and icon search."""

    async def get_image(self, prompt: str, width: int, height: int) -> str | None:
        return None

    async def get_icon(self, query: str) -> str | None:
        return None


class DefaultAssetProvider(AssetProvider):
    def __init__(self, image_generator: ImageGenerator, icon_searcher: IconSearcher):
        self._image_generator = image_generator
        self._icon_searcher = icon_searcher

    async def get_image(self, prompt: str, width: int = 1024, height: int = 1024) -> str | None:
        try:
            asset = await self._image_generator.generate_image(ImagePrompt(prompt=prompt))
            return asset.path
        except Exception:
            return None

    async def get_icon(self, query: str) -> str | None:
        try:
            result = await self._icon_searcher.search_icons(query, k=1)
            return result[0] if result else None
        except Exception:
            return None


_SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
}


class SlideRenderer:
    def __init__(self, tokens: DesignTokens, asset_provider: AssetProvider | None = None):
        self.tokens = tokens
        self.asset_provider = asset_provider

    async def render(self, slide: Slide, recipe: LayoutRecipe, spec: SlideSpec) -> None:
        all_regions = recipe.all_regions(len(spec.blocks))
        sorted_regions = sorted(all_regions, key=lambda r: (r.z_layer, all_regions.index(r)))

        for region in sorted_regions:
            role = recipe.region_roles.get(region.region_id)
            if role == RegionRole.DECORATION:
                self._render_decoration(slide, region)
            elif role in (RegionRole.TITLE, RegionRole.SUBTITLE, RegionRole.BODY, RegionRole.CARD_BODY, RegionRole.INDEX, RegionRole.FOOTER):
                self._render_text(slide, region, self._text_for_region(region.region_id, role, recipe, spec), role)
            elif role == RegionRole.CARD:
                self._render_card_background(slide, region)
            elif role == RegionRole.ICON:
                await self._render_icon(slide, region, recipe, spec)
            elif role == RegionRole.IMAGE:
                await self._render_image(slide, region, recipe, spec)

    def _text_for_region(self, region_id: str, role: RegionRole, recipe: LayoutRecipe, spec: SlideSpec) -> str:
        source = recipe.region_text_sources.get(region_id)
        if source == "slide_title" or (source is None and role == RegionRole.TITLE):
            return spec.title
        if source == "index":
            block_index = recipe.region_block_indexes.get(region_id, 0)
            return f"{block_index + 1:02d}"

        block = self._block_for_region(region_id, recipe, spec)
        if block is None:
            return ""
        if source == "block_title":
            return block.title
        if source == "block_title_text":
            return f"{block.title}\n{block.text}".strip()
        return block.text

    def _block_for_region(self, region_id: str, recipe: LayoutRecipe, spec: SlideSpec):
        block_index = recipe.region_block_indexes.get(region_id)
        if block_index is not None and 0 <= block_index < len(spec.blocks):
            return spec.blocks[block_index]
        return spec.blocks[0] if spec.blocks else None

    async def _render_icon(self, slide: Slide, region: Region, recipe: LayoutRecipe, spec: SlideSpec) -> None:
        block = self._block_for_region(region.region_id, recipe, spec)
        query = block.icon_query if block else None
        if not query or self.asset_provider is None:
            return
        icon_path = await self.asset_provider.get_icon(query)
        if icon_path:
            self._add_picture(slide, region, icon_path)

    async def _render_image(self, slide: Slide, region: Region, recipe: LayoutRecipe, spec: SlideSpec) -> None:
        block = self._block_for_region(region.region_id, recipe, spec)
        prompt = block.image_prompt if block else None
        if not prompt or self.asset_provider is None:
            return
        image_path = await self.asset_provider.get_image(prompt, width=1024, height=1024)
        if image_path:
            self._add_picture(slide, region, image_path)

    def _add_picture(self, slide: Slide, region: Region, path: str) -> None:
        left, top, width, height = region.to_absolute(self.tokens.slide_width, self.tokens.slide_height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))

    def _render_card_background(self, slide: Slide, region: Region) -> None:
        self._render_decoration(slide, region)

    def _render_text(self, slide: Slide, region: Region, text: str, role: RegionRole) -> None:
        left, top, width, height = region.to_absolute(self.tokens.slide_width, self.tokens.slide_height)
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        para = frame.paragraphs[0]
        para.text = text

        if role == RegionRole.TITLE:
            para.font.size = Pt(self.tokens.title_size)
            para.font.name = self.tokens.title_font
            para.font.color.rgb = RGBColor.from_string(self.tokens.text_primary.lstrip("#"))
        elif role == RegionRole.SUBTITLE:
            para.font.size = Pt(self.tokens.subtitle_size)
            para.font.name = self.tokens.subtitle_font
            para.font.color.rgb = RGBColor.from_string(self.tokens.text_secondary.lstrip("#"))
        else:
            para.font.size = Pt(self.tokens.body_size)
            para.font.name = self.tokens.body_font
            para.font.color.rgb = RGBColor.from_string(self.tokens.text_primary.lstrip("#"))

    def _render_decoration(self, slide: Slide, region: Region) -> None:
        shape_type = _SHAPE_MAP.get(region.decoration_shape or "rect", MSO_SHAPE.RECTANGLE)
        left, top, width, height = region.to_absolute(self.tokens.slide_width, self.tokens.slide_height)
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))

        if region.fill_role:
            token_attr = getattr(self.tokens, region.fill_role, None)
            if token_attr:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(token_attr.lstrip("#"))

        if region.line_role:
            token_attr = getattr(self.tokens, region.line_role, None)
            if token_attr:
                shape.line.color.rgb = RGBColor.from_string(token_attr.lstrip("#"))
        else:
            shape.line.fill.background()
