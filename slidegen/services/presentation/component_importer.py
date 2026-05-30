"""Content style importer — Phase 2.

Import chapter_content slides from an uploaded PPTX, assign shape roles via
local rules + LLM fallback, build Style objects compatible with
ChapterContentPage.generate_slide(), and write them into shapes.json.
"""

import hashlib
import json
import os
import re
import tempfile
import uuid
from collections.abc import Callable
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Literal

from agno.agent import Agent
from agno.models.base import Model
from loguru import logger
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pydantic import BaseModel, Field

from slidegen.core import settings
from slidegen.schemas.gen_request import BaseGenerationRequest
from slidegen.services.presentation.components import (
    ChapterLayout,
    ComponentContentType,
    ComponentsManager,
    CShape,
    LayoutType,
    Location,
    Style,
    remove_custDataLst,
)
from slidegen.services.presentation.page_classifier import (
    PageClassification,
    PageType,
    PageTypeClassifier,
    ShapeSummary,
)
from slidegen.services.slidegen.workflow import get_llm_instance

_SHAPE_NAME_TRANSLATIONS = (
    ("圆角矩形", "rounded rectangle"),
    ("任意多边形", "freeform"),
    ("矩形", "rectangle"),
    ("椭圆", "ellipse"),
    ("图形", "shape"),
    ("形状", "shape"),
)


def _sanitize_ascii_name(value: str | None, *, fallback: str) -> str:
    """Return a lowercase ASCII identifier safe for style keys and XML names."""
    text = value or ""
    for source, target in _SHAPE_NAME_TRANSLATIONS:
        text = text.replace(source, target)
    text = re.sub(r"[^A-Za-z0-9]+", "_", text).strip("_").lower()
    return text or fallback


# ---------------------------------------------------------------------------
# Step 1: Core types
# ---------------------------------------------------------------------------


class ImportSlideStatus(str, Enum):
    IMPORTED = "imported"
    SKIPPED = "skipped"
    DRY_RUN = "dry_run"
    FAILED = "failed"


@dataclass
class ContentStyleImportOptions:
    pptx_path: Path
    user_id: uuid.UUID
    llm_config_id: uuid.UUID | None = None
    target_json_path: Path = Path(settings.COMPONENTS_PATH)
    min_page_confidence: float = 0.7
    min_role_confidence: float = 0.7
    dry_run: bool = True
    overwrite_existing: bool = False
    backup: bool = True
    preview_dir: Path | None = None


@dataclass
class ShapeAssignment:
    shape_id: int
    content_type: ComponentContentType | Literal["skip"]
    group_index: int | None  # None for slide-level; 0-3 for content groups
    include: bool
    reason: str
    confidence: float


@dataclass
class ValidationResult:
    ok: bool
    reason: str
    preview_path: Path | None = None


@dataclass
class ImportedSlideReport:
    page_index: int
    page_type: PageType
    page_confidence: float
    status: ImportSlideStatus
    layout: ChapterLayout | None
    style_name: str | None
    reason: str
    warnings: list[str] = field(default_factory=list)
    validation: ValidationResult | None = None


@dataclass
class ImportReport:
    pptx_path: Path
    target_json_path: Path
    dry_run: bool
    imported_count: int
    skipped_count: int
    failed_count: int
    slides: list[ImportedSlideReport]


class AgentShapeAssignment(BaseModel):
    shape_id: int
    content_type: Literal["title", "content", "number", "picture", "icon", "decoration", "skip"]
    group_index: int | None = Field(default=None, ge=0, le=3)
    confidence: float = Field(ge=0.0, le=1.0)
    reason: str = Field(min_length=1)


class AgentShapeRoleOutput(BaseModel):
    point_count: int = Field(ge=1, le=4)
    assignments: list[AgentShapeAssignment]
    confidence: float = Field(ge=0.0, le=1.0)
    reason: str = Field(min_length=1)


# ---------------------------------------------------------------------------
# Step 6: Atomic write & fingerprint
# ---------------------------------------------------------------------------


def atomic_save_json(
    data: dict[str, Any],
    json_path: str | Path,
    backup: bool = True,
) -> None:
    """Atomically write JSON to file using temp file + os.replace."""
    path = Path(json_path)
    dir_path = path.parent

    if backup and path.exists():
        bak_path = path.with_suffix(path.suffix + ".bak")
        bak_path.write_text(path.read_text(encoding="utf-8"), encoding="utf-8")

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".json", dir=str(dir_path))
    try:
        with os.fdopen(tmp_fd, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
        os.replace(tmp_path, str(path))
    except BaseException:
        # Clean up the temp file on any failure
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise


def compute_pptx_fingerprint(pptx_path: str | Path) -> str:
    """Compute SHA-256 hash of the PPTX file content."""
    path = Path(pptx_path)
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


_FINGERPRINTS_FILE_SUFFIX = ".fingerprints"


def _fingerprints_path(json_path: str | Path) -> Path:
    """Return the companion fingerprints file path."""
    return Path(json_path).with_suffix(Path(json_path).suffix + _FINGERPRINTS_FILE_SUFFIX)


def check_fingerprint(pptx_path: str | Path, json_path: str | Path) -> bool:
    """Check if this PPTX has already been imported by comparing fingerprint.

    Returns True if already imported (should skip), False if new.
    """
    fp_path = _fingerprints_path(json_path)
    if not fp_path.exists():
        return False

    current_hash = compute_pptx_fingerprint(pptx_path)
    with open(fp_path, encoding="utf-8") as f:
        data = json.load(f)

    for entry in data.get("fingerprints", []):
        if entry.get("hash") == current_hash:
            return True
    return False


def store_fingerprint(pptx_path: str | Path, json_path: str | Path) -> None:
    """Store the fingerprint of a newly imported PPTX."""
    fp_path = _fingerprints_path(json_path)
    current_hash = compute_pptx_fingerprint(pptx_path)
    ppt_name = Path(pptx_path).stem

    existing: dict[str, Any] = {"fingerprints": []}
    if fp_path.exists():
        with open(fp_path, encoding="utf-8") as f:
            existing = json.load(f)

    existing["fingerprints"].append(
        {
            "path": ppt_name,
            "hash": current_hash,
        }
    )

    atomic_save_json(existing, fp_path, backup=False)


# ---------------------------------------------------------------------------
# Step 2: Local rule-based shape role classifier
# ---------------------------------------------------------------------------

_ICON_AREA_THRESHOLD = 3_900_000 * 3_900_000  # reuse ComponentsManager.is_icon threshold
_SHORT_TEXT_THRESHOLD = 20  # characters
_NUMBER_MAX_LENGTH = 3  # max digits for a "number" shape


class LocalShapeRoleClassifier:
    """Assign shape roles using deterministic local rules."""

    def classify(self, summaries: list[ShapeSummary]) -> list[ShapeAssignment]:
        """Return assignments for all shapes. Ambiguous shapes get content_type="skip"."""
        assignments: list[ShapeAssignment] = []

        for s in summaries:
            if s.is_placeholder:
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type="skip",
                        group_index=None,
                        include=False,
                        reason="Placeholder shapes are handled by the template.",
                        confidence=1.0,
                    )
                )
                continue

            # Picture shapes
            if s.is_picture:
                area = s.width * s.height
                if area < _ICON_AREA_THRESHOLD:
                    ct = ComponentContentType.ICON
                    reason = "Small-area picture classified as icon."
                else:
                    ct = ComponentContentType.PICTURE
                    reason = "Large-area picture classified as picture."
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ct,
                        group_index=None,  # will be assigned later during grouping
                        include=True,
                        reason=reason,
                        confidence=0.85,
                    )
                )
                continue

            # No text frame, no picture → decoration
            if not s.has_text_frame:
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ComponentContentType.DECORATION,
                        group_index=None,
                        include=True,
                        reason="Non-text, non-picture shape classified as decoration.",
                        confidence=0.9,
                    )
                )
                continue

            # Text frame with empty text → decoration
            if not s.text:
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ComponentContentType.DECORATION,
                        group_index=None,
                        include=True,
                        reason="Empty text frame classified as decoration.",
                        confidence=0.9,
                    )
                )
                continue

            # --- Text content shapes ---
            # Slide-level title: top of slide, large font, short text, wide
            is_slide_level_title = self._is_slide_level_title(s, summaries)
            if is_slide_level_title:
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type="skip",
                        group_index=None,
                        include=False,
                        reason="Slide-level title excluded from content style.",
                        confidence=0.85,
                    )
                )
                continue

            # Number: all digits, very short
            stripped = s.text.strip()
            if self._is_number_text(stripped):
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ComponentContentType.NUMBER,
                        group_index=None,
                        include=True,
                        reason=f"Short numeric text '{stripped}' classified as number.",
                        confidence=0.9,
                    )
                )
                continue

            # Title heuristic: large font or bold, short text, near top of group
            if self._is_title_like(s, summaries):
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ComponentContentType.TITLE,
                        group_index=None,
                        include=True,
                        reason="Short bold/large-font text classified as title.",
                        confidence=0.75,
                    )
                )
                continue

            # Content: longer text
            if len(stripped) > _SHORT_TEXT_THRESHOLD:
                assignments.append(
                    ShapeAssignment(
                        shape_id=s.shape_id,
                        content_type=ComponentContentType.CONTENT,
                        group_index=None,
                        include=True,
                        reason="Longer text classified as content.",
                        confidence=0.75,
                    )
                )
                continue

            # Ambiguous short text that isn't clearly number or title → skip for LLM
            assignments.append(
                ShapeAssignment(
                    shape_id=s.shape_id,
                    content_type="skip",
                    group_index=None,
                    include=False,
                    reason="Ambiguous short text; deferred to LLM agent.",
                    confidence=0.0,
                )
            )

        return assignments

    def assign_group_indices(
        self,
        assignments: list[ShapeAssignment],
        summaries: list[ShapeSummary],
    ) -> list[ShapeAssignment]:
        """Assign group_index to shapes that need it, using position clustering.

        Shapes with content_type TITLE/CONTENT/NUMBER/ICON/PICTURE are grouped.
        Group count must be 1-4; otherwise all shapes revert to "skip".
        """
        summary_by_id = {s.shape_id: s for s in summaries}

        # Collect shapes that need a group assignment
        groupable = [
            a
            for a in assignments
            if a.include and a.content_type not in ("skip", ComponentContentType.DECORATION) and a.group_index is None
        ]

        if not groupable:
            return assignments

        original_state = {
            a.shape_id: (
                a.content_type,
                a.group_index,
                a.include,
                a.reason,
                a.confidence,
            )
            for a in groupable
        }

        for axis in ("y", "x"):
            self._restore_groupable_state(groupable, original_state)
            if self._apply_axis_grouping(groupable, summary_by_id, axis=axis):
                if self._grouping_is_compatible(assignments):
                    return assignments

        self._restore_groupable_state(groupable, original_state)
        self._apply_axis_grouping(groupable, summary_by_id, axis="y")

        # Decorations get group_index=None (slide-level)
        for a in assignments:
            if a.include and a.content_type == ComponentContentType.DECORATION:
                a.group_index = None

        return assignments

    @staticmethod
    def _apply_axis_grouping(
        groupable: list[ShapeAssignment],
        summary_by_id: dict[int, ShapeSummary],
        *,
        axis: Literal["x", "y"],
    ) -> bool:
        """Assign groups by a primary visual axis and return whether count is valid."""
        primary_attr = "y" if axis == "y" else "x"
        secondary_attr = "x" if axis == "y" else "y"

        groupable.sort(
            key=lambda a: (
                getattr(summary_by_id[a.shape_id], primary_attr),
                getattr(summary_by_id[a.shape_id], secondary_attr),
            )
        )

        positions = [getattr(summary_by_id[a.shape_id], primary_attr) for a in groupable]
        max_pos = max(positions) if positions else 0
        pos_range = max_pos - min(positions) if positions else 0
        gap_threshold = pos_range * 0.3 if pos_range > 0 else 500_000

        groups: list[list[ShapeAssignment]] = []
        current_group: list[ShapeAssignment] = []

        for i, a in enumerate(groupable):
            if i == 0:
                current_group.append(a)
                continue
            gap = positions[i] - positions[i - 1]
            if gap > gap_threshold:
                groups.append(current_group)
                current_group = [a]
            else:
                current_group.append(a)
        if current_group:
            groups.append(current_group)

        group_count = len(groups)

        # If group count is not in 1-4, mark all groupable as skip
        if group_count < 1 or group_count > 4:
            for a in groupable:
                a.content_type = "skip"
                a.include = False
                a.reason = f"Group count {group_count} is outside 1-4 range."
                a.confidence = 0.0
            return False

        for gi, group in enumerate(groups):
            for a in group:
                a.group_index = gi

        return True

    @staticmethod
    def _restore_groupable_state(
        groupable: list[ShapeAssignment],
        original_state: dict[
            int,
            tuple[
                ComponentContentType | Literal["skip"],
                int | None,
                bool,
                str,
                float,
            ],
        ],
    ) -> None:
        for a in groupable:
            state = original_state[a.shape_id]
            a.content_type, a.group_index, a.include, a.reason, a.confidence = state

    @staticmethod
    def _grouping_is_compatible(assignments: list[ShapeAssignment]) -> bool:
        layout_type = identify_layout_type(assignments)
        if layout_type is None:
            return False
        valid, _ = validate_compatibility(assignments, layout_type)
        return valid

    @staticmethod
    def _is_number_text(text: str) -> bool:
        """Check if text looks like a number indicator (01, 02, 1., roman)."""
        if not text:
            return False
        stripped = text.strip()
        # Pure digits, very short
        if stripped.isdigit() and len(stripped) <= _NUMBER_MAX_LENGTH:
            return True
        # Digits with dot suffix: "1.", "2."
        if re.match(r"^\d{1,2}\.$", stripped):
            return True
        # Roman numerals: i, ii, iii, iv
        if re.match(r"^[ivx]{1,4}$", stripped.casefold()):
            return True
        # Zero-padded: 01, 02
        if re.match(r"^0\d$", stripped):
            return True
        return False

    @staticmethod
    def _is_slide_level_title(
        s: ShapeSummary,
        all_summaries: list[ShapeSummary],
    ) -> bool:
        """Detect slide-level title: wide, near top, large font, short text."""
        if not s.text:
            return False
        # Top 25% of slide
        max_y = max(ss.y for ss in all_summaries) if all_summaries else 0
        if s.y > max_y * 0.25 and max_y > 0:
            return False
        # Wide: spans > 60% of slide width
        max_width = max(ss.x + ss.width for ss in all_summaries) if all_summaries else 0
        if max_width > 0 and s.width < max_width * 0.6:
            return False
        # Short text
        if len(s.text) > 60:
            return False
        # Large font or bold
        if s.font_size is not None and s.font_size >= 24:
            return True
        if s.is_bold:
            return True
        return False

    @staticmethod
    def _is_title_like(
        s: ShapeSummary,
        _all_summaries: list[ShapeSummary],
    ) -> bool:
        """Check if text shape looks like a section/group title."""
        if not s.text:
            return False
        # Short text
        if len(s.text) > _SHORT_TEXT_THRESHOLD:
            return False
        # Large font
        if s.font_size is not None and s.font_size >= 18:
            return True
        # Bold
        if s.is_bold:
            return True
        return False


# ---------------------------------------------------------------------------
# Step 3: LLM Shape Role Agent
# ---------------------------------------------------------------------------

SHAPE_ROLE_AGENT_INSTRUCTIONS = """\
You are analyzing shapes on a PowerPoint chapter-content slide. Your task is to
assign each shape a role and a content group index.

Roles:
- title: section/subsection heading text within a content group
- content: body paragraph text or description
- number: ordinal indicator like "01", "02", "1.", roman numerals
- picture: large image or illustration serving as main visual
- icon: small decorative image near a title or number
- decoration: non-text, non-image visual element (line, circle, background shape)
- skip: shapes that should not be included in the style template (slide-level title, logo, watermark, page number, footer)

Group index (0-3):
- Each content group is a self-contained unit with title + content, optionally number, icon, picture.
- Group 0 is the first content group, group 1 the second, etc.
- Shapes that apply to the entire slide (decoration, skip) get group_index null.

You receive a list of shape summaries with position, size, text, and font info.
Use the visual layout to determine which shapes belong to which group.
"""

ShapeRoleAgentFactory = Callable[[Model | Any], Any]


class ShapeRoleAgent:
    """LLM fallback for assigning shape roles when local rules are ambiguous."""

    def __init__(
        self,
        *,
        min_role_confidence: float = 0.7,
        agent_factory: ShapeRoleAgentFactory | None = None,
    ) -> None:
        self.min_role_confidence = min_role_confidence
        self.agent_factory = agent_factory or self._default_agent_factory

    async def assign_roles(
        self,
        *,
        model: Model | Any,
        summaries: list[ShapeSummary],
        local_assignments: list[ShapeAssignment],
        page_classification: PageClassification,
        min_role_confidence: float | None = None,
    ) -> list[ShapeAssignment]:
        """Ask the LLM to resolve ambiguous shape assignments.

        Only shapes marked as content_type="skip" by local rules are sent
        to the LLM, but the full slide context is provided.
        """
        confidence_threshold = self.min_role_confidence if min_role_confidence is None else min_role_confidence
        ambiguous_ids = [
            a.shape_id
            for a in local_assignments
            if a.content_type == "skip" and a.reason != "Placeholder shapes are handled by the template."
        ]
        if not ambiguous_ids:
            return local_assignments

        prompt = self._build_prompt(
            summaries=summaries,
            local_assignments=local_assignments,
            ambiguous_ids=ambiguous_ids,
            page_classification=page_classification,
        )

        agent = self.agent_factory(model)
        try:
            response = await agent.arun(prompt)
            output = self._coerce_agent_output(getattr(response, "content", response))

            if output.confidence < confidence_threshold:
                logger.warning(
                    "Shape role agent confidence {:.2f} below threshold {:.2f}",
                    output.confidence,
                    confidence_threshold,
                )
                return local_assignments  # keep local assignments unchanged

            # Validate and merge LLM assignments
            summary_ids = {s.shape_id for s in summaries}
            merged = list(local_assignments)

            for agent_a in output.assignments:
                if agent_a.shape_id not in summary_ids:
                    logger.warning("Agent returned unknown shape_id {}", agent_a.shape_id)
                    continue
                if agent_a.shape_id not in ambiguous_ids:
                    continue  # LLM should only resolve ambiguous shapes

                ct_str = agent_a.content_type
                if ct_str == "skip":
                    # LLM also says skip — keep as skip
                    continue

                ct = ComponentContentType(ct_str)
                # Find the local assignment and update it
                for i, m in enumerate(merged):
                    if m.shape_id == agent_a.shape_id:
                        merged[i] = ShapeAssignment(
                            shape_id=agent_a.shape_id,
                            content_type=ct,
                            group_index=agent_a.group_index,
                            include=True,
                            reason=agent_a.reason,
                            confidence=agent_a.confidence,
                        )
                        break

            # Validate point_count vs group count
            group_indices = {a.group_index for a in merged if a.include and a.group_index is not None}
            if output.point_count != len(group_indices):
                logger.warning(
                    "Agent point_count {} != actual group count {}",
                    output.point_count,
                    len(group_indices),
                )
                # Fall back — keep local assignments
                return local_assignments

            return merged

        except Exception as exc:
            logger.warning("Shape role agent failed: {}", exc)
            return local_assignments

    def _build_prompt(
        self,
        *,
        summaries: list[ShapeSummary],
        local_assignments: list[ShapeAssignment],
        ambiguous_ids: list[int],
        page_classification: PageClassification,
    ) -> str:
        payload = {
            "page_classification": {
                "page_type": page_classification.page_type.value,
                "confidence": page_classification.confidence,
                "reason": page_classification.reason,
            },
            "ambiguous_shape_ids": ambiguous_ids,
            "local_assignments": [
                {
                    "shape_id": a.shape_id,
                    "content_type": (
                        a.content_type.value if isinstance(a.content_type, ComponentContentType) else a.content_type
                    ),
                    "group_index": a.group_index,
                    "include": a.include,
                    "reason": a.reason,
                }
                for a in local_assignments
            ],
            "shapes": [
                {
                    "shape_id": s.shape_id,
                    "name": s.name,
                    "text": s.text,
                    "x": s.x,
                    "y": s.y,
                    "width": s.width,
                    "height": s.height,
                    "has_text_frame": s.has_text_frame,
                    "is_placeholder": s.is_placeholder,
                    "is_picture": s.is_picture,
                    "font_size": s.font_size,
                    "is_bold": s.is_bold,
                }
                for s in summaries
            ],
        }
        return (
            "Assign roles and group indices to the ambiguous shapes on this "
            "chapter-content slide. Use the provided full context including "
            "local rule results.\n\n"
            f"{json.dumps(payload, ensure_ascii=False)}"
        )

    def _coerce_agent_output(self, content: Any) -> AgentShapeRoleOutput:
        if isinstance(content, AgentShapeRoleOutput):
            return content
        if isinstance(content, dict):
            return AgentShapeRoleOutput.model_validate(content)
        if isinstance(content, str):
            return AgentShapeRoleOutput.model_validate_json(content)
        return AgentShapeRoleOutput.model_validate(content)

    def _default_agent_factory(self, model: Model | Any) -> Agent:
        return Agent(
            name="Shape role assigner",
            instructions=[SHAPE_ROLE_AGENT_INSTRUCTIONS],
            model=model,
            output_schema=AgentShapeRoleOutput,
            structured_outputs=False,
        )


# ---------------------------------------------------------------------------
# Step 4: Layout type identification & compatibility validation
# ---------------------------------------------------------------------------


def identify_layout_type(assignments: list[ShapeAssignment]) -> ChapterLayout | None:
    """Determine ChapterLayout from the count of content groups.

    Count unique group_index values among TITLE/CONTENT assignments.
    Returns None if no valid group structure found.
    """
    group_indices: set[int] = set()
    for a in assignments:
        if (
            a.include
            and a.group_index is not None
            and a.content_type in (ComponentContentType.TITLE, ComponentContentType.CONTENT)
        ):
            group_indices.add(a.group_index)

    count = len(group_indices)
    if count < 1 or count > 4:
        return None

    return ChapterLayout(count)


def validate_compatibility(
    assignments: list[ShapeAssignment],
    layout_type: ChapterLayout,
) -> tuple[bool, str]:
    """Validate that the assignment-derived style is compatible with ChapterContentPage.

    Rules:
    - TITLE location count must equal point_count
    - CONTENT location count must equal point_count
    - At least one TITLE and one CONTENT shape must exist
    - Group indices 0..point_count-1 must all be covered

    Returns (True, "") if valid, (False, reason) if invalid.
    """
    point_count = layout_type.value

    title_groups = set()
    content_groups = set()

    for a in assignments:
        if not a.include or a.group_index is None:
            continue
        if a.content_type == ComponentContentType.TITLE:
            title_groups.add(a.group_index)
        elif a.content_type == ComponentContentType.CONTENT:
            content_groups.add(a.group_index)

    if not title_groups:
        return False, "No TITLE shape found."
    if not content_groups:
        return False, "No CONTENT shape found."

    # TITLE groups must cover 0..point_count-1
    expected_groups = set(range(point_count))
    if title_groups != expected_groups:
        return False, (f"TITLE group indices {sorted(title_groups)} don't cover 0..{point_count - 1}.")
    if content_groups != expected_groups:
        return False, (f"CONTENT group indices {sorted(content_groups)} don't cover 0..{point_count - 1}.")

    return True, ""


# ---------------------------------------------------------------------------
# Step 5: StyleBuilder
# ---------------------------------------------------------------------------


class StyleBuilder:
    """Build a Style from shape assignments — assignment-driven, not heuristic."""

    TEXT_PLACEHOLDER = "Text"

    def build_style_from_assignments(
        self,
        slide: Slide,
        _layout_type: ChapterLayout,
        style_name: str,
        assignments: list[ShapeAssignment],
    ) -> Style:
        """Construct a Style from a slide and its shape assignments."""
        style = Style(style_name)

        # Build a map from shape_id to slide shape
        slide_shape_by_id: dict[int, Any] = {}
        for shape in slide.shapes:
            slide_shape_by_id[shape.shape_id] = shape

        # Collect included assignments sorted by zorder (slide position)
        included = [a for a in assignments if a.include]
        included.sort(key=lambda a: self._zorder_for(a, slide))

        # Group shapes by content_type + role, then merge same-XML shapes
        # across group indices into one CShape with multiple locations.
        role_groups: dict[str, list[tuple[ShapeAssignment, Any]]] = {}

        for a in included:
            ct_label = a.content_type.value if isinstance(a.content_type, ComponentContentType) else str(a.content_type)
            key = f"{ct_label}_group_{a.group_index}"
            shape_obj = slide_shape_by_id.get(a.shape_id)
            if shape_obj is None:
                continue
            role_groups.setdefault(key, []).append((a, shape_obj))

        # Now merge shapes that have the same content_type across group indices
        # by checking XML similarity
        merged_by_role: dict[str, list[tuple[ShapeAssignment, Any]]] = {}
        role_counter: dict[str, int] = {}

        for _key, items in role_groups.items():
            ct_str = (
                items[0][0].content_type.value
                if isinstance(items[0][0].content_type, ComponentContentType)
                else str(items[0][0].content_type)
            )

            # Try to merge with existing same-role shapes
            merged_key = None
            for existing_key, existing_items in merged_by_role.items():
                existing_ct = (
                    existing_items[0][0].content_type.value
                    if isinstance(existing_items[0][0].content_type, ComponentContentType)
                    else str(existing_items[0][0].content_type)
                )
                if existing_ct != ct_str:
                    continue

                # Check if all shapes are same-XML
                # For pictures/icons (xml=None), check if they're similar by type
                is_picture_shape = items[0][1].shape_type == MSO_SHAPE_TYPE.PICTURE
                existing_is_picture = existing_items[0][1].shape_type == MSO_SHAPE_TYPE.PICTURE

                if is_picture_shape and existing_is_picture:
                    # Pictures can be merged if they share the same content_type role
                    # and we just combine locations
                    merged_key = existing_key
                    break

                if not is_picture_shape and not existing_is_picture:
                    # Text/decoration shapes — check XML similarity
                    try:
                        xml1 = remove_custDataLst(items[0][1]._element.xml)
                        xml2 = remove_custDataLst(existing_items[0][1]._element.xml)
                        if ComponentsManager.are_same_shape(xml1, xml2):
                            merged_key = existing_key
                            break
                    except Exception:
                        pass

            if merged_key:
                merged_by_role[merged_key].extend(items)
            else:
                counter = role_counter.get(ct_str, 0)
                new_key = f"{ct_str}_{counter}"
                role_counter[ct_str] = counter + 1
                merged_by_role[new_key] = items

        # Build CShape objects from merged groups
        for shape_name, items in merged_by_role.items():
            ct = items[0][0].content_type
            if isinstance(ct, str):
                ct = ComponentContentType(ct)

            # Sort items by group_index for location ordering
            items.sort(key=lambda t: (t[0].group_index if t[0].group_index is not None else 99))

            first_shape_obj = items[0][1]
            is_picture = first_shape_obj.shape_type == MSO_SHAPE_TYPE.PICTURE

            xml_str: str | None = None
            if not is_picture:
                try:
                    raw_xml = first_shape_obj._element.xml
                    xml_str = remove_custDataLst(raw_xml)
                    xml_str = self._replace_shape_xml_names(xml_str)
                    text_placeholder = "" if ct == ComponentContentType.DECORATION else self.TEXT_PLACEHOLDER
                    xml_str = self._replace_shape_xml_text(xml_str, text_placeholder)
                except Exception as exc:
                    logger.warning("Failed to extract XML for shape {}: {}", shape_name, exc)
                    continue

            locations: list[Location] = []
            for _, shape_obj in items:
                loc = Location(
                    x=int(shape_obj.left),
                    y=int(shape_obj.top),
                    width=int(shape_obj.width),
                    height=int(shape_obj.height),
                )
                locations.append(loc)

            # Compute zorder from the first item's position in slide
            zorder = self._zorder_for(items[0][0], slide)

            cshape = CShape(
                xml=xml_str,
                zorder=zorder,
                content_type=ct,
                location=locations,
            )
            style.add_shape(shape_name, cshape)

        return style

    @staticmethod
    def _replace_shape_xml_text(xml_str: str, placeholder: str) -> str:
        """Replace imported text runs while preserving shape and run styling."""
        root = etree.fromstring(xml_str)
        for text_element in root.findall(".//a:t", namespaces=root.nsmap):
            if text_element.text:
                text_element.text = placeholder
        return etree.tostring(
            root,
            encoding="unicode",
            pretty_print=True,
            xml_declaration=False,
        )

    @staticmethod
    def _replace_shape_xml_names(xml_str: str) -> str:
        """Replace imported PowerPoint shape names with stable English identifiers."""
        root = etree.fromstring(xml_str)
        for index, c_nv_pr in enumerate(root.xpath(".//*[local-name()='cNvPr']"), start=1):
            fallback = f"shape_{c_nv_pr.get('id') or index}"
            c_nv_pr.set("name", _sanitize_ascii_name(c_nv_pr.get("name"), fallback=fallback))
        return etree.tostring(
            root,
            encoding="unicode",
            pretty_print=True,
            xml_declaration=False,
        )

    @staticmethod
    def _zorder_for(assignment: ShapeAssignment, slide: Slide) -> int:
        """Get zorder index for a shape from its position in slide.shapes."""
        for i, shape in enumerate(slide.shapes):
            if shape.shape_id == assignment.shape_id:
                return i
        return 0


# ---------------------------------------------------------------------------
# Step 7: ContentStyleImporter — main service
# ---------------------------------------------------------------------------


class ContentStyleImporter:
    """Import content styles from an uploaded PPTX presentation."""

    def __init__(
        self,
        components_manager: ComponentsManager,
        *,
        page_classifier: PageTypeClassifier | None = None,
        local_classifier: LocalShapeRoleClassifier | None = None,
        shape_role_agent: ShapeRoleAgent | None = None,
        style_builder: StyleBuilder | None = None,
        style_name_stem: str | None = None,
    ) -> None:
        self.cm = components_manager
        self.page_classifier = page_classifier or PageTypeClassifier()
        self.local_classifier = local_classifier or LocalShapeRoleClassifier()
        self.shape_role_agent = shape_role_agent or ShapeRoleAgent()
        self.style_builder = style_builder or StyleBuilder()
        self.style_name_stem = (
            _sanitize_ascii_name(style_name_stem, fallback="deck")
            if style_name_stem is not None
            else None
        )

    async def import_from_pptx(
        self,
        *,
        pptx_path: str | Path,
        user_id: uuid.UUID,
        llm_config_id: uuid.UUID | None = None,
        model: Model | Any | None = None,
        target_json_path: str | Path | None = None,
        min_page_confidence: float = 0.7,
        min_role_confidence: float = 0.7,
        dry_run: bool = True,
        overwrite_existing: bool = False,
        backup: bool = True,
        preview_dir: Path | None = None,
    ) -> ImportReport:
        """Full import pipeline per the Phase 2 spec."""
        path = Path(pptx_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        json_path = Path(target_json_path) if target_json_path else Path(settings.COMPONENTS_PATH)

        # Resolve LLM model once for reuse (same pattern as PageTypeClassifier)
        if model is None:
            request = BaseGenerationRequest(
                content="PPT content style import",
                user_id=user_id,
                llm_config_id=llm_config_id,
            )
            model = await get_llm_instance(request)

        # Fingerprint check
        if not overwrite_existing and check_fingerprint(path, json_path):
            return ImportReport(
                pptx_path=path,
                target_json_path=json_path,
                dry_run=dry_run,
                imported_count=0,
                skipped_count=0,
                failed_count=0,
                slides=[
                    ImportedSlideReport(
                        page_index=0,
                        page_type=PageType.UNKNOWN,
                        page_confidence=0.0,
                        status=ImportSlideStatus.SKIPPED,
                        layout=None,
                        style_name=None,
                        reason="PPTX already imported (fingerprint match).",
                    )
                ],
            )

        # Classify pages
        classifications = await self.page_classifier.classify_pages(
            pptx_path=str(path),
            user_id=user_id,
            llm_config_id=llm_config_id,
            model=model,
        )

        presentation = Presentation(str(path))
        slides: list[ImportedSlideReport] = []
        imported_count = 0
        skipped_count = 0
        failed_count = 0

        for classification in classifications:
            page_index = classification.page_index

            # Only process chapter_content pages
            if classification.page_type != PageType.CHAPTER_CONTENT:
                slides.append(
                    ImportedSlideReport(
                        page_index=page_index,
                        page_type=classification.page_type,
                        page_confidence=classification.confidence,
                        status=ImportSlideStatus.SKIPPED,
                        layout=None,
                        style_name=None,
                        reason=f"Only chapter_content pages are imported (got {classification.page_type.value}).",
                    )
                )
                skipped_count += 1
                continue

            # Confidence filter
            if classification.confidence < min_page_confidence:
                slides.append(
                    ImportedSlideReport(
                        page_index=page_index,
                        page_type=classification.page_type,
                        page_confidence=classification.confidence,
                        status=ImportSlideStatus.SKIPPED,
                        layout=None,
                        style_name=None,
                        reason=f"Page confidence {classification.confidence:.2f} below threshold {min_page_confidence:.2f}.",
                    )
                )
                skipped_count += 1
                continue

            # Process chapter_content page
            try:
                report = await self._process_content_slide(
                    presentation=presentation,
                    page_index=page_index,
                    classification=classification,
                    model=model,
                    _json_path=json_path,
                    pptx_path=path,
                    dry_run=dry_run,
                    overwrite_existing=overwrite_existing,
                    min_role_confidence=min_role_confidence,
                    preview_dir=preview_dir,
                )
                slides.append(report)
                if report.status == ImportSlideStatus.IMPORTED or report.status == ImportSlideStatus.DRY_RUN:
                    imported_count += 1
                elif report.status == ImportSlideStatus.SKIPPED:
                    skipped_count += 1
                else:
                    failed_count += 1
            except Exception as exc:
                logger.warning("Failed to process slide {}: {}", page_index, exc)
                slides.append(
                    ImportedSlideReport(
                        page_index=page_index,
                        page_type=classification.page_type,
                        page_confidence=classification.confidence,
                        status=ImportSlideStatus.FAILED,
                        layout=None,
                        style_name=None,
                        reason=str(exc),
                    )
                )
                failed_count += 1

        # Atomic save (if not dry_run and at least one import succeeded)
        if not dry_run and imported_count > 0:
            self.cm.save_to_json(json_path, backup=backup)
            store_fingerprint(path, json_path)

        return ImportReport(
            pptx_path=path,
            target_json_path=json_path,
            dry_run=dry_run,
            imported_count=imported_count,
            skipped_count=skipped_count,
            failed_count=failed_count,
            slides=slides,
        )

    @staticmethod
    def _build_dummy_heading(point_count: int) -> Any:
        """Build a minimal Heading tree with `point_count` children for validation."""
        from slidegen.services.document.markdown.elements import Heading

        root = Heading(level=2, text="Validation Test")
        for i in range(point_count):
            child = Heading(level=3, text=f"Point {i + 1}")
            root.append(child)
        return root

    async def _validate_render_roundtrip(
        self,
        *,
        style: Style,
        layout_type: ChapterLayout,
        preview_dir: Path | None = None,
    ) -> ValidationResult:
        """Force-render one slide with dummy content using the given style.

        Creates a temporary PPTX from a blank template, injects the style
        into ChapterContentPage.generate_slide, and verifies the output
        can be saved without error.
        """
        from pptx import Presentation as PPTXPresentation

        from slidegen.services.presentation.pages import ChapterContentPage

        point_count = layout_type.value
        content = self._build_dummy_heading(point_count)

        try:
            prs = PPTXPresentation()
            # Add two blank slides: one as layout source, one as target.
            prs.slides.add_slide(prs.slide_layouts[6])  # blank (layout source)
            prs.slides.add_slide(prs.slide_layouts[6])  # blank (placeholder)

            await ChapterContentPage.generate_slide(
                prs=prs,
                content=content,
                chapter_page_index=0,
                slide_index=1,
                style_override=style,
            )

            if len(prs.slides) < 3:
                return ValidationResult(ok=False, reason="No new slide was added during validation render.")

            if preview_dir is not None:
                preview_dir.mkdir(parents=True, exist_ok=True)
                preview_path = preview_dir / f"validation_{style.name}.pptx"
                prs.save(str(preview_path))
                return ValidationResult(ok=True, reason="Render OK.", preview_path=preview_path)
            else:
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=True) as tmp:
                    prs.save(tmp.name)
                return ValidationResult(ok=True, reason="Render OK.")

        except Exception as exc:
            logger.warning("Validation render failed for style '{}': {}", style.name, exc)
            return ValidationResult(ok=False, reason=str(exc))

    async def _process_content_slide(
        self,
        *,
        presentation: PresentationType,
        page_index: int,
        classification: PageClassification,
        model: Model | Any,
        _json_path: Path,
        pptx_path: Path,
        dry_run: bool,
        overwrite_existing: bool,
        min_role_confidence: float,
        preview_dir: Path | None = None,
    ) -> ImportedSlideReport:
        """Process a single chapter_content slide through the full pipeline."""
        slide = presentation.slides[page_index]
        summaries = self.page_classifier.summarize_slide(slide)

        # Step 1: local rules
        local_assignments = self.local_classifier.classify(summaries)

        # Step 2: LLM fallback for ambiguous shapes
        has_ambiguous = any(
            a.content_type == "skip"
            and a.include is False
            and a.reason != "Placeholder shapes are handled by the template."
            for a in local_assignments
        )

        if has_ambiguous:
            local_assignments = await self.shape_role_agent.assign_roles(
                model=model,
                summaries=summaries,
                local_assignments=local_assignments,
                page_classification=classification,
                min_role_confidence=min_role_confidence,
            )

        # Step 3: assign group indices
        local_assignments = self.local_classifier.assign_group_indices(
            local_assignments,
            summaries,
        )

        # Step 4: identify layout type
        layout_type = identify_layout_type(local_assignments)
        if layout_type is None:
            return ImportedSlideReport(
                page_index=page_index,
                page_type=classification.page_type,
                page_confidence=classification.confidence,
                status=ImportSlideStatus.SKIPPED,
                layout=None,
                style_name=None,
                reason="Could not identify a valid layout type (1-4 content groups).",
            )

        # Step 5: compatibility validation
        valid, reason = validate_compatibility(local_assignments, layout_type)
        if not valid:
            return ImportedSlideReport(
                page_index=page_index,
                page_type=classification.page_type,
                page_confidence=classification.confidence,
                status=ImportSlideStatus.SKIPPED,
                layout=layout_type,
                style_name=None,
                reason=f"Incompatible with ChapterContentPage: {reason}",
            )

        # Step 6: build style
        style_name = self._generate_style_name(pptx_path, page_index)

        # Handle name collision
        layout = self.cm.get_layout_type(layout_type)
        if layout and not overwrite_existing:
            style_name = self._unique_style_name(style_name, layout)

        # Ensure target layout exists in ComponentsManager
        if layout is None:
            layout = LayoutType(layout_type.str_value)
            if not dry_run:
                self.cm.layout_types[layout_type.str_value] = layout

        style = self.style_builder.build_style_from_assignments(
            slide=slide,
            _layout_type=layout_type,
            style_name=style_name,
            assignments=local_assignments,
        )

        duplicate_style = self._find_duplicate_style(style, layout)
        if duplicate_style is not None and (not overwrite_existing or duplicate_style.name != style_name):
            return ImportedSlideReport(
                page_index=page_index,
                page_type=classification.page_type,
                page_confidence=classification.confidence,
                status=ImportSlideStatus.SKIPPED,
                layout=layout_type,
                style_name=duplicate_style.name,
                reason=f"Duplicate style already exists in {layout_type.str_value}: {duplicate_style.name}.",
            )

        # --- Validation gate ---
        validation = await self._validate_render_roundtrip(
            style=style,
            layout_type=layout_type,
            preview_dir=preview_dir if dry_run else None,
        )
        if not validation.ok:
            return ImportedSlideReport(
                page_index=page_index,
                page_type=classification.page_type,
                page_confidence=classification.confidence,
                status=ImportSlideStatus.FAILED,
                layout=layout_type,
                style_name=style_name,
                reason=f"Validation render failed: {validation.reason}",
                validation=validation,
            )

        # Add to ComponentsManager (in-memory) only for real imports.
        if not dry_run:
            if overwrite_existing and layout.get_style(style_name):
                layout.styles.pop(style_name)
            layout.add_style(style)

        warnings: list[str] = []
        for a in local_assignments:
            if not a.include and a.reason != "Placeholder shapes are handled by the template.":
                warnings.append(f"Shape {a.shape_id} excluded: {a.reason}")

        status = ImportSlideStatus.DRY_RUN if dry_run else ImportSlideStatus.IMPORTED

        return ImportedSlideReport(
            page_index=page_index,
            page_type=classification.page_type,
            page_confidence=classification.confidence,
            status=status,
            layout=layout_type,
            style_name=style_name,
            reason=f"Detected {layout_type.value} reusable content groups.",
            warnings=warnings,
            validation=validation,
        )

    def _generate_style_name(self, pptx_path: str | Path, page_index: int) -> str:
        """Generate style name: upload_<ppt_stem>_<fingerprint>_p<page_number>."""
        path = Path(pptx_path)
        ppt_stem = self.style_name_stem or _sanitize_ascii_name(path.stem, fallback="deck")
        short_fingerprint = compute_pptx_fingerprint(path)[:8]
        page_number = page_index + 1
        return f"upload_{ppt_stem}_{short_fingerprint}_p{page_number}"

    @staticmethod
    def _unique_style_name(style_name: str, layout: LayoutType) -> str:
        """Return a non-conflicting style name for append-only imports."""
        if style_name not in layout.styles:
            return style_name

        suffix = 2
        while f"{style_name}_{suffix}" in layout.styles:
            suffix += 1
        return f"{style_name}_{suffix}"

    @classmethod
    def _find_duplicate_style(cls, style: Style, layout: LayoutType) -> Style | None:
        """Return an existing style with the same visual shape structure."""
        for existing_style in layout.style_list:
            if cls._are_same_style(style, existing_style):
                return existing_style
        return None

    @classmethod
    def _are_same_style(cls, left: Style, right: Style) -> bool:
        if len(left.shape_list) != len(right.shape_list):
            return False

        unmatched = list(right.shape_list)
        for left_shape in left.shape_list:
            match_index = next(
                (
                    index
                    for index, right_shape in enumerate(unmatched)
                    if cls._are_same_cshape(left_shape, right_shape)
                ),
                None,
            )
            if match_index is None:
                return False
            unmatched.pop(match_index)
        return True

    @staticmethod
    def _are_same_cshape(left: CShape, right: CShape) -> bool:
        if left.content_type != right.content_type:
            return False
        if left.zorder != right.zorder:
            return False
        if left.location != right.location:
            return False
        if left.xml is None or right.xml is None:
            return left.xml is None and right.xml is None
        return ComponentsManager.are_same_shape(left.xml, right.xml)
