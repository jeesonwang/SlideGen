from __future__ import annotations

from typing import TYPE_CHECKING, Any

import pptx
from pptx.presentation import Presentation

from slidegen.exceptions import MarkdownDocumentError
from slidegen.services.document.markdown import Heading, MarkdownDocument
from slidegen.services.presentation.design_tokens import extract_design_tokens_from_presentation
from slidegen.services.presentation.post_render_validator import PostRenderValidator
from slidegen.services.presentation.recipe_agent import RecipeAgent, resolve_recipe
from slidegen.services.presentation.semantic import BlockKind, BlockSpec, SlideKind, SlideSpec, build_content_slide_spec
from slidegen.services.presentation.slide_renderer import AssetProvider, SlideRenderer
from slidegen.services.slidegen.outline_structure import iter_chapter_slide_groups

if TYPE_CHECKING:
    from agno.models.base import Model


class MarkdownToPresentation:
    """Generate a PPT presentation from a markdown document."""

    def __init__(self, *, recipe_model: Model | None = None, asset_provider: AssetProvider | None = None) -> None:
        self.recipe_model = recipe_model
        self.asset_provider = asset_provider

    async def generate(
        self,
        template_prs: Presentation,
        markdown_document: MarkdownDocument,
    ) -> Presentation:
        """generate the complete PPT presentation"""
        main_heading = markdown_document.main
        if main_heading is None:
            raise MarkdownDocumentError("Markdown document must have a main heading")

        chapter_slide_groups = list(iter_chapter_slide_groups(main_heading))
        chapters: list[Heading] = [group.chapter for group in chapter_slide_groups]
        if not chapters:
            raise MarkdownDocumentError("Markdown document must have at least one level 2 heading")

        output_prs = pptx.Presentation()
        output_prs.slide_width = template_prs.slide_width
        output_prs.slide_height = template_prs.slide_height
        tokens = extract_design_tokens_from_presentation(template_prs, "general")
        renderer = SlideRenderer(tokens, asset_provider=self.asset_provider)
        agent = RecipeAgent(model=self.recipe_model) if self.recipe_model is not None else None

        async def render_spec(spec: SlideSpec) -> None:
            recipe = await resolve_recipe(spec, tokens, agent=agent, enable_agent=agent is not None)
            slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
            await renderer.render(slide, recipe, spec)

        await render_spec(self._cover_spec(main_heading))
        await render_spec(self._agenda_spec(chapters))
        for group in chapter_slide_groups:
            await render_spec(self._chapter_home_spec(group.chapter))
            for content_heading in group.slides:
                await render_spec(build_content_slide_spec(content_heading))
        await render_spec(self._closing_spec())

        issues = PostRenderValidator(mode="fail").validate(output_prs)
        if issues:
            messages = "; ".join(issue.message for issue in issues if issue.level == "error")
            if messages:
                raise MarkdownDocumentError(f"Recipe-rendered presentation failed geometry validation: {messages}")
        return output_prs

    def _cover_spec(self, main_heading: Heading) -> SlideSpec:
        return SlideSpec(
            kind=SlideKind.COVER,
            title=main_heading.element_text,
            source_level=1,
            blocks=(BlockSpec(kind=BlockKind.SUBTITLE, title="", text=""),),
        )

    def _agenda_spec(self, chapters: list[Heading]) -> SlideSpec:
        return SlideSpec(
            kind=SlideKind.AGENDA,
            title="目录",
            source_level=1,
            blocks=tuple(BlockSpec(kind=BlockKind.POINT, title=chapter.element_text, text="") for chapter in chapters),
        )

    def _chapter_home_spec(self, chapter: Heading) -> SlideSpec:
        return SlideSpec(
            kind=SlideKind.SECTION_COVER,
            title=chapter.element_text,
            source_level=2,
            blocks=(BlockSpec(kind=BlockKind.TITLE, title=chapter.element_text, text=""),),
        )

    def _closing_spec(self) -> SlideSpec:
        return SlideSpec(
            kind=SlideKind.CLOSING,
            title="谢谢",
            source_level=1,
            blocks=(BlockSpec(kind=BlockKind.TITLE, title="谢谢", text=""),),
        )
