from time import perf_counter

from loguru import logger
from pptx.presentation import Presentation

from slidegen.exceptions import MarkdownDocumentError
from slidegen.services.document.markdown import Heading, MarkdownDocument
from slidegen.services.presentation.pages import (
    CatalogPage,
    ChapterContentPage,
    ChapterHomePage,
    CoverPage,
    EndPage,
)
from slidegen.services.slidegen.outline_structure import iter_chapter_slide_groups


class PresentationOrchestrator:
    """
    Generate a PPT presentation from a markdown document

    Template presentation object. Template presentation must have at least 5 slides
    and must be in accordance with `CoverPage`, `CatalogPage`, `ChapterHomePage`, `ChapterContentPage`, `EndPage`.
    """

    def __init__(self) -> None:
        self.slide_index = 0
        self.chapter_index = 1

    async def generate(
        self,
        template_prs: Presentation,
        markdown_document: MarkdownDocument,
        cover_page_index: int = 0,
        catalog_page_index: int = 1,
    ) -> Presentation:
        """generate the complete PPT presentation"""
        started_at = perf_counter()

        headings = [h for h in markdown_document.descendants if hasattr(h, "level") and h.level == 1]

        if not headings:
            raise MarkdownDocumentError("Markdown document must have at least one level 1 heading")

        main_heading = markdown_document.main
        if main_heading is None:
            raise MarkdownDocumentError("Markdown document must have a main heading")
        logger.info("PPT conversion: generating cover slide for '{}'", main_heading.element_text)
        await CoverPage.generate_slide(template_prs, main_heading, cover_page_index=cover_page_index)

        chapter_slide_groups = list(iter_chapter_slide_groups(main_heading))
        chapters: list[Heading] = [group.chapter for group in chapter_slide_groups]
        total_content_slides = sum(len(group.slides) for group in chapter_slide_groups)

        if not chapters:
            raise MarkdownDocumentError("Markdown document must have at least one level 2 heading")
        logger.info(
            "PPT conversion: markdown parsed into {} chapters and {} content slides",
            len(chapters),
            total_content_slides,
        )

        logger.info("PPT conversion: generating catalog slides")
        catalog_last_index = await CatalogPage.generate_slide(
            template_prs, chapters, catalog_page_index=catalog_page_index
        )
        logger.info("PPT conversion: catalog generated through slide index {}", catalog_last_index)

        chapter_home_page_index = catalog_last_index + 1
        chapter_content_page_index = chapter_home_page_index + 1
        end_page_index = chapter_content_page_index + 1
        current_slide_index = end_page_index + 1

        generated_content_slides = 0
        for chapter_index, group in enumerate(chapter_slide_groups):
            logger.info(
                "PPT conversion: generating chapter {}/{} home slide: {}",
                chapter_index + 1,
                len(chapter_slide_groups),
                group.chapter.element_text,
            )
            await ChapterHomePage.generate_slide(
                template_prs,
                group.chapter,
                chapter_home_page_index=chapter_home_page_index,
                chapter_number=chapter_index + 1,
                slide_index=current_slide_index,
            )
            current_slide_index += 1

            for slide in group.slides:
                generated_content_slides += 1
                logger.info(
                    "PPT conversion: generating content slide {}/{}: {}",
                    generated_content_slides,
                    total_content_slides,
                    slide.element_text,
                )
                await ChapterContentPage.generate_slide(
                    template_prs,
                    slide,
                    chapter_page_index=chapter_content_page_index,
                    slide_index=current_slide_index,
                )
                current_slide_index += 1

        logger.info("PPT conversion: generating ending slide")
        await EndPage.generate_slide(template_prs, end_page_index=end_page_index, slide_index=current_slide_index)

        self._cleanup_template_slides(
            template_prs, [chapter_home_page_index, chapter_content_page_index, end_page_index]
        )
        logger.info(
            "PPT conversion: completed {} generated slides in {:.2f}s",
            len(template_prs.slides),
            perf_counter() - started_at,
        )

        return template_prs

    def _cleanup_template_slides(self, template_prs: Presentation, be_removed_slides_index: list[int]) -> None:
        # delete slides from back to front
        be_removed_slides_index.sort(reverse=True)
        for i in be_removed_slides_index:
            CoverPage.remove_slide(template_prs, i)
