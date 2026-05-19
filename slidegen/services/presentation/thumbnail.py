"""Thumbnail generation service for PPTX templates.

Uses LibreOffice to convert PPTX to PDF, then PyMuPDF to render pages as PNG.
Supports single-page and multi-page grid thumbnails with optional placeholder highlighting.
"""

import importlib.util
import shutil
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from loguru import logger
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.shapes.autoshape import Shape

from slidegen.core.config import settings

from slidegen.services.presentation.pdf_exporter import (
    LibreOfficeNotFoundError as PdfLibreOfficeNotFound,
)
from slidegen.services.presentation.pdf_exporter import PdfExportError, pdf_exporter

# Constants for thumbnail generation
THUMBNAIL_WIDTH = 300  # Default thumbnail width in pixels for grids
CONVERSION_DPI = 100  # DPI for PDF to image conversion
MAX_COLS = 6  # Maximum number of columns in grid
DEFAULT_COLS = 5  # Default number of columns in grid
JPEG_QUALITY = 95  # JPEG compression quality

# Grid layout constants
GRID_PADDING = 20  # Padding between thumbnails in pixels
BORDER_WIDTH = 2  # Border width around thumbnails
FONT_SIZE_RATIO = 0.12  # Font size as fraction of thumbnail width
LABEL_PADDING_RATIO = 0.4  # Label padding as fraction of font size

# Placeholder highlighting constants
PLACEHOLDER_OUTLINE_COLOR = (255, 0, 0, 255)  # Bright red, fully opaque
PLACEHOLDER_STROKE_WIDTH_RATIO = 150  # Divisor for proportional stroke width


class ThumbnailGenerationError(Exception):
    """Exception raised when thumbnail generation fails."""

    pass


class LibreOfficeNotFoundError(ThumbnailGenerationError):
    """Exception raised when LibreOffice is not installed."""

    pass


class ThumbnailGenerator:
    """Generates thumbnails for PPTX templates using LibreOffice and PyMuPDF.

    Features:
    - Single-page thumbnails for template preview
    - Multi-page grid thumbnails showing all slides
    - Optional placeholder highlighting for template development
    - Automatic caching with modification time validation
    - Dependency checking for LibreOffice and PyMuPDF

    Examples:
        # Generate simple thumbnail
        generator = ThumbnailGenerator()
        thumbnail = generator.generate_thumbnail("general")

        # Generate grid with all slides
        grid = generator.generate_grid_thumbnail("general", cols=5)

        # Generate with placeholder highlighting (for development)
        debug_thumb = generator.generate_thumbnail(
            "general",
            outline_placeholders=True,
            output_suffix="_debug"
        )
    """

    def __init__(
        self,
        templates_dir: Path | None = None,
        thumbnails_dir: Path | None = None,
        thumbnail_width: int | None = None,
    ):
        self.templates_dir = templates_dir or (settings.COMPONENTS_BASE_PATH / "templates")
        self.thumbnails_dir = thumbnails_dir or settings.THUMBNAILS_DIR
        self.thumbnail_width = thumbnail_width or settings.THUMBNAIL_WIDTH

        # Ensure thumbnails directory exists
        self.thumbnails_dir.mkdir(parents=True, exist_ok=True)

    def _get_template_path(self, template_name: str) -> Path:
        """Get the path to a template file, trying multiple naming conventions.

        Tries in order:
        1. {template_name}.pptx - Direct naming
        2. template_{template_name}.pptx - Legacy naming with prefix

        Args:
            template_name: Name of the template (e.g., "general" or "mytemplate").

        Returns:
            Path to the template file (may not exist yet).
            Returns the first existing match, or the direct naming path if none exist.

        Example:
            >>> # For template_name="general"
            >>> # Tries: "general.pptx", then "template_general.pptx"
            >>> path = generator._get_template_path("general")
        """
        # Try direct naming first (new style)
        direct_path = self.templates_dir / f"{template_name}.pptx"
        if direct_path.exists():
            return direct_path

        # Try legacy naming with "template_" prefix
        legacy_path = self.templates_dir / f"template_{template_name}.pptx"
        if legacy_path.exists():
            return legacy_path

        # Return direct path as default (caller will handle FileNotFoundError)
        return direct_path

    def check_dependencies(self) -> dict[str, bool]:
        """Check if all required dependencies are available.

        Returns:
            Dictionary with dependency names as keys and availability as values.

        Example:
            >>> deps = generator.check_dependencies()
            >>> if not all(deps.values()):
            ...     print(f"Missing dependencies: {[k for k, v in deps.items() if not v]}")
        """
        dependencies = {}

        # Check LibreOffice
        try:
            self._find_libreoffice()
            dependencies["libreoffice"] = True
        except LibreOfficeNotFoundError:
            dependencies["libreoffice"] = False

        # Check PyMuPDF (fitz)
        dependencies["pymupdf"] = importlib.util.find_spec("fitz") is not None

        # Check python-pptx
        dependencies["python-pptx"] = importlib.util.find_spec("pptx") is not None

        # Check Pillow
        dependencies["pillow"] = importlib.util.find_spec("PIL") is not None

        return dependencies

    def _find_libreoffice(self) -> str:
        """Find the LibreOffice executable path.

        Returns:
            Path to the LibreOffice executable.

        Raises:
            LibreOfficeNotFoundError: If LibreOffice is not found.
        """
        try:
            return pdf_exporter._find_libreoffice()
        except PdfLibreOfficeNotFound as e:
            raise LibreOfficeNotFoundError(str(e))

    def _convert_pptx_to_pdf(self, pptx_path: Path, output_dir: Path) -> Path:
        """Convert PPTX file to PDF using LibreOffice.

        Delegates to PdfExporter.

        Args:
            pptx_path: Path to the PPTX file.
            output_dir: Directory to save the PDF.

        Returns:
            Path to the generated PDF file.

        Raises:
            ThumbnailGenerationError: If conversion fails.
        """
        pdf_path = output_dir / f"{pptx_path.stem}.pdf"
        try:
            pdf_exporter.convert(str(pptx_path), str(pdf_path))
        except PdfExportError as e:
            raise ThumbnailGenerationError(str(e))
        return pdf_path

    def _render_pdf_page(
        self,
        pdf_path: Path,
        output_path: Path,
        width: int,
        page_num: int = 0,
    ) -> None:
        """Render a specific page of a PDF as PNG.

        Args:
            pdf_path: Path to the PDF file.
            output_path: Path to save the PNG thumbnail.
            width: Target width in pixels.
            page_num: Page number to render (0-indexed).

        Raises:
            ThumbnailGenerationError: If rendering fails.
        """
        try:
            doc = fitz.open(pdf_path)

            if len(doc) == 0:
                error_msg = f"PDF has no pages: {pdf_path}"
                logger.error(error_msg)
                raise ThumbnailGenerationError(error_msg)

            logger.debug(f"PDF contains {len(doc)} page(s)")

            if page_num >= len(doc):
                error_msg = f"Page {page_num} does not exist in PDF with {len(doc)} pages"
                logger.error(error_msg)
                raise ThumbnailGenerationError(error_msg)

            page = doc[page_num]

            # Calculate zoom factor to achieve target width
            zoom = width / page.rect.width
            matrix = fitz.Matrix(zoom, zoom)

            # Render page to pixmap
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)

            # Save as PNG
            pixmap.save(str(output_path))

            doc.close()
            logger.debug(f"Thumbnail saved to: {output_path}")

        except fitz.FileDataError as e:
            error_msg = f"Failed to open PDF {pdf_path}: {e}"
            logger.error(error_msg)
            raise ThumbnailGenerationError(error_msg)

    def _get_placeholder_regions(
        self, pptx_path: Path
    ) -> tuple[dict[int, list[dict[str, float]]], tuple[float, float]]:
        """Extract text placeholder regions from all slides.

        Args:
            pptx_path: Path to the PPTX file.

        Returns:
            Tuple of (placeholder_regions, slide_dimensions) where:
            - placeholder_regions: Dict mapping slide index to list of regions
            - slide_dimensions: Tuple of (width_inches, height_inches)

        Each region dict contains: left, top, width, height (all in inches)
        """
        try:
            prs = Presentation(str(pptx_path))
            placeholder_regions = {}

            # Get slide dimensions in inches (EMU to inches conversion)
            slide_width_inches = (prs.slide_width or 9144000) / 914400.0
            slide_height_inches = (prs.slide_height or 5143500) / 914400.0

            for slide_idx, slide in enumerate(prs.slides):
                regions: list[dict[str, float]] = []

                # Recursively collect all shapes with text
                def collect_text_shapes(
                    shape: Shape, regions_list: list[dict[str, float]], parent_left: float = 0, parent_top: float = 0
                ) -> None:
                    """Recursively collect shapes with text, calculating absolute positions."""
                    if hasattr(shape, "shapes"):  # GroupShape
                        group_left = shape.left if hasattr(shape, "left") else 0
                        group_top = shape.top if hasattr(shape, "top") else 0
                        abs_group_left = parent_left + group_left
                        abs_group_top = parent_top + group_top

                        for child in shape.shapes:
                            collect_text_shapes(child, regions_list, abs_group_left, abs_group_top)
                    else:
                        # Regular shape - check if it has text
                        if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
                            shape_left = shape.left if hasattr(shape, "left") else 0
                            shape_top = shape.top if hasattr(shape, "top") else 0

                            regions_list.append(
                                {
                                    "left": (parent_left + shape_left) / 914400.0,  # EMU to inches
                                    "top": (parent_top + shape_top) / 914400.0,
                                    "width": (shape.width if hasattr(shape, "width") else 0) / 914400.0,
                                    "height": (shape.height if hasattr(shape, "height") else 0) / 914400.0,
                                }
                            )

                for shape in slide.shapes:
                    collect_text_shapes(shape, regions)

                if regions:
                    placeholder_regions[slide_idx] = regions

            logger.debug(f"Found placeholders on {len(placeholder_regions)} slides")
            return placeholder_regions, (slide_width_inches, slide_height_inches)

        except Exception as e:
            logger.warning(f"Failed to extract placeholder regions: {e}")
            return {}, (10.0, 7.5)  # Return empty dict with default slide size

    def _apply_placeholder_outlines(
        self,
        image: Image.Image,
        regions: list[dict[str, float]],
        slide_dimensions: tuple[float, float],
    ) -> Image.Image:
        """Apply red outlines to placeholder regions on an image.

        Args:
            image: PIL Image to modify.
            regions: List of region dicts with left, top, width, height in inches.
            slide_dimensions: Tuple of (width_inches, height_inches).

        Returns:
            Modified PIL Image with outlines.
        """
        # Convert to RGBA for transparency support
        if image.mode != "RGBA":
            image = image.convert("RGBA")

        orig_w, orig_h = image.size
        slide_width_inches, slide_height_inches = slide_dimensions

        # Calculate scale factors
        x_scale = orig_w / slide_width_inches
        y_scale = orig_h / slide_height_inches

        # Create overlay for outlines
        overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
        draw = ImageDraw.Draw(overlay)

        # Draw outline for each region
        for region in regions:
            # Convert from inches to pixels
            px_left = int(region["left"] * x_scale)
            px_top = int(region["top"] * y_scale)
            px_width = int(region["width"] * x_scale)
            px_height = int(region["height"] * y_scale)

            # Calculate stroke width proportional to image size
            stroke_width = max(5, min(orig_w, orig_h) // PLACEHOLDER_STROKE_WIDTH_RATIO)

            # Draw rectangle outline
            draw.rectangle(
                [(px_left, px_top), (px_left + px_width, px_top + px_height)],
                outline=PLACEHOLDER_OUTLINE_COLOR,
                width=stroke_width,
            )

        # Composite overlay onto image
        image = Image.alpha_composite(image, overlay)
        return image.convert("RGB")

    def _render_pdf_first_page(self, pdf_path: Path, output_path: Path, width: int) -> None:
        """Render the first page of a PDF as PNG.

        Deprecated: Use _render_pdf_page instead.
        This method is kept for backward compatibility.
        """
        self._render_pdf_page(pdf_path, output_path, width, page_num=0)

    def is_thumbnail_valid(self, template_name: str, suffix: str = "") -> bool:
        """Check if a cached thumbnail is still valid.

        A thumbnail is valid if:
        1. The thumbnail file exists
        2. The thumbnail was created after the template was last modified

        Args:
            template_name: Name of the template (e.g., "general").
            suffix: Optional suffix for the thumbnail filename.

        Returns:
            True if the thumbnail is valid, False otherwise.
        """
        filename = f"{template_name}{suffix}.png"
        thumbnail_path = self.thumbnails_dir / filename
        template_path = self._get_template_path(template_name)

        if not thumbnail_path.exists():
            return False

        if not template_path.exists():
            logger.warning(f"Template file does not exist: {template_path}")
            return False

        # Compare modification times
        thumbnail_mtime = thumbnail_path.stat().st_mtime
        template_mtime = template_path.stat().st_mtime

        return thumbnail_mtime > template_mtime

    def generate_thumbnail(
        self,
        template_name: str,
        force: bool = False,
        outline_placeholders: bool = False,
        output_suffix: str = "",
    ) -> Path:
        """Generate a thumbnail for a template.

        Args:
            template_name: Name of the template (e.g., "general").
            force: If True, regenerate even if cached thumbnail is valid.
            outline_placeholders: If True, draw red outlines around text placeholders.
            output_suffix: Suffix to add to output filename (e.g., "_debug").

        Returns:
            Path to the thumbnail PNG file.

        Raises:
            FileNotFoundError: If the template does not exist.
            LibreOfficeNotFoundError: If LibreOffice is not installed.
            ThumbnailGenerationError: If generation fails.
        """
        filename = f"{template_name}{output_suffix}.png"
        thumbnail_path = self.thumbnails_dir / filename
        template_path = self._get_template_path(template_name)

        # Check if template exists
        if not template_path.exists():
            error_msg = f"Template not found: {template_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)

        # Check cache validity (unless force or outlining is enabled)
        if not force and not outline_placeholders and self.is_thumbnail_valid(template_name, output_suffix):
            logger.debug(f"Using cached thumbnail for: {template_name}{output_suffix}")
            return thumbnail_path

        logger.info(f"Generating thumbnail for template: {template_name}{output_suffix}")

        # Use temporary directory for intermediate files
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            # Convert PPTX to PDF
            pdf_path = self._convert_pptx_to_pdf(template_path, temp_path)

            # Render first page to temporary PNG
            temp_png = temp_path / "temp.png"
            self._render_pdf_page(pdf_path, temp_png, self.thumbnail_width, page_num=0)

            # Apply placeholder outlines if requested
            if outline_placeholders:
                logger.debug("Applying placeholder outlines")
                placeholder_regions, slide_dimensions = self._get_placeholder_regions(template_path)

                if 0 in placeholder_regions:
                    # Load image and apply outlines
                    opened_img = Image.open(temp_png)
                    processed_img = self._apply_placeholder_outlines(
                        opened_img, placeholder_regions[0], slide_dimensions
                    )
                    processed_img.save(str(thumbnail_path), "PNG")
                else:
                    # No placeholders on first slide, just copy
                    logger.debug("No placeholders found on first slide")
                    shutil.copy(temp_png, thumbnail_path)
            else:
                # Just copy the rendered image
                shutil.copy(temp_png, thumbnail_path)

        logger.info(f"Thumbnail generated: {thumbnail_path}")
        return thumbnail_path

    def generate_grid_thumbnail(
        self,
        template_name: str,
        cols: int = DEFAULT_COLS,
        force: bool = False,
        outline_placeholders: bool = False,
        output_suffix: str = "_grid",
    ) -> list[Path]:
        """Generate grid thumbnails showing all slides.

        Creates a grid layout of all slides with configurable columns.
        For presentations with many slides, multiple grid files are created.

        Args:
            template_name: Name of the template (e.g., "general").
            cols: Number of columns (max 6).
            force: If True, regenerate even if cached.
            outline_placeholders: If True, highlight text placeholders.
            output_suffix: Suffix for output files (default: "_grid").

        Returns:
            List of paths to generated grid image files.

        Raises:
            FileNotFoundError: If the template does not exist.
            LibreOfficeNotFoundError: If LibreOffice is not installed.
            ThumbnailGenerationError: If generation fails.

        Example:
            >>> grids = generator.generate_grid_thumbnail("general", cols=5)
            >>> # Creates: general_grid.png (or general_grid-1.png, general_grid-2.png, etc.)
        """
        template_path = self._get_template_path(template_name)

        # Check if template exists
        if not template_path.exists():
            error_msg = f"Template not found: {template_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)

        # Validate columns
        cols = min(cols, MAX_COLS)
        if cols < 1:
            cols = DEFAULT_COLS

        # Check cache validity (unless force or outlining is enabled)
        if not force and not outline_placeholders:
            # Check if any grid files exist and are valid
            grid_base = self.thumbnails_dir / f"{template_name}{output_suffix}.png"
            grid_numbered = self.thumbnails_dir / f"{template_name}{output_suffix}-1.png"

            if grid_base.exists() or grid_numbered.exists():
                # Get template modification time
                template_mtime = template_path.stat().st_mtime

                # Find all grid files
                cached_grids = list(self.thumbnails_dir.glob(f"{template_name}{output_suffix}*.png"))

                # Check if all cached grids are newer than template
                if cached_grids and all(g.stat().st_mtime > template_mtime for g in cached_grids):
                    logger.debug(f"Using cached grid thumbnails for: {template_name}")
                    return cached_grids

        logger.info(f"Generating grid thumbnail for: {template_name} with {cols} columns")

        # Use temporary directory
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            # Get placeholder regions if highlighting is enabled
            placeholder_regions = None
            slide_dimensions = None
            if outline_placeholders:
                logger.debug("Extracting placeholder regions")
                placeholder_regions, slide_dimensions = self._get_placeholder_regions(template_path)
                if placeholder_regions:
                    logger.debug(f"Found placeholders on {len(placeholder_regions)} slides")

            # Convert PPTX to PDF
            pdf_path = self._convert_pptx_to_pdf(template_path, temp_path)

            # Open PDF and get page count
            doc = fitz.open(pdf_path)
            num_pages = len(doc)
            logger.debug(f"PDF contains {num_pages} page(s)")

            if num_pages == 0:
                doc.close()
                raise ThumbnailGenerationError("PDF has no pages")

            # Render all pages to images
            slide_images = []
            for page_num in range(num_pages):
                page_path = temp_path / f"slide_{page_num:03d}.png"

                # Render page
                page = doc[page_num]
                zoom = THUMBNAIL_WIDTH / page.rect.width
                matrix = fitz.Matrix(zoom, zoom)
                pixmap = page.get_pixmap(matrix=matrix, alpha=False)
                pixmap.save(str(page_path))

                # Apply placeholder outlines if needed
                if (
                    outline_placeholders
                    and placeholder_regions
                    and page_num in placeholder_regions
                    and slide_dimensions is not None
                ):
                    opened_img = Image.open(page_path)
                    processed_img = self._apply_placeholder_outlines(
                        opened_img, placeholder_regions[page_num], slide_dimensions
                    )
                    processed_img.save(str(page_path), "PNG")

                slide_images.append(page_path)

            doc.close()

            # Create grids
            grid_files = self._create_grids(
                slide_images,
                cols,
                THUMBNAIL_WIDTH,
                template_name,
                output_suffix,
            )

        logger.info(f"Created {len(grid_files)} grid file(s) for {template_name}")
        return grid_files

    def _create_grids(
        self,
        image_paths: list[Path],
        cols: int,
        width: int,
        template_name: str,
        output_suffix: str,
    ) -> list[Path]:
        """Create thumbnail grid(s) from slide images.

        Args:
            image_paths: List of paths to slide images.
            cols: Number of columns.
            width: Thumbnail width in pixels.
            template_name: Template name for output filename.
            output_suffix: Suffix for output filename.

        Returns:
            List of paths to created grid files.
        """
        from PIL import ImageFont

        # Maximum images per grid: cols × (cols + 1) for better proportions
        max_images_per_grid = cols * (cols + 1)
        grid_files = []

        logger.debug(f"Creating grids with {cols} columns (max {max_images_per_grid} images per grid)")

        # Font size and label padding
        font_size = int(width * FONT_SIZE_RATIO)
        label_padding = int(font_size * LABEL_PADDING_RATIO)

        # Load font
        try:
            font = ImageFont.load_default(size=font_size)
        except Exception:
            font = ImageFont.load_default()

        # Split images into chunks
        for chunk_idx, start_idx in enumerate(range(0, len(image_paths), max_images_per_grid)):
            end_idx = min(start_idx + max_images_per_grid, len(image_paths))
            chunk_images = image_paths[start_idx:end_idx]

            # Create grid for this chunk
            grid = self._create_single_grid(
                chunk_images,
                cols,
                width,
                font,
                font_size,
                label_padding,
                start_idx,
            )

            # Generate output filename
            if len(image_paths) <= max_images_per_grid:
                # Single grid
                grid_filename = self.thumbnails_dir / f"{template_name}{output_suffix}.png"
            else:
                # Multiple grids
                grid_filename = self.thumbnails_dir / f"{template_name}{output_suffix}-{chunk_idx + 1}.png"

            # Save grid
            grid.save(str(grid_filename), "PNG")
            grid_files.append(grid_filename)

        return grid_files

    def _create_single_grid(
        self,
        image_paths: list[Path],
        cols: int,
        width: int,
        font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
        font_size: int,
        label_padding: int,
        start_slide_num: int = 0,
    ) -> Image.Image:
        """Create a single thumbnail grid from slide images.

        Args:
            image_paths: List of paths to slide images.
            cols: Number of columns.
            width: Thumbnail width in pixels.
            font: PIL font for labels.
            font_size: Font size in pixels.
            label_padding: Padding around labels.
            start_slide_num: Starting slide number for labels.

        Returns:
            PIL Image containing the grid.
        """
        # Get dimensions from first image
        with Image.open(image_paths[0]) as img:
            aspect = img.height / img.width
        height = int(width * aspect)

        # Calculate grid size
        rows = (len(image_paths) + cols - 1) // cols
        grid_w = cols * width + (cols + 1) * GRID_PADDING
        grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

        # Create grid
        grid = Image.new("RGB", (grid_w, grid_h), "white")
        draw = ImageDraw.Draw(grid)

        # Place thumbnails
        for i, img_path in enumerate(image_paths):
            row, col = i // cols, i % cols
            x = col * width + (col + 1) * GRID_PADDING
            y_base = row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING

            # Add label
            label = f"{start_slide_num + i + 1}"
            bbox = draw.textbbox((0, 0), label, font=font)
            text_w = bbox[2] - bbox[0]
            draw.text(
                (x + (width - text_w) // 2, y_base + label_padding),
                label,
                fill="black",
                font=font,
            )

            # Add thumbnail below label
            y_thumbnail = y_base + label_padding + font_size + label_padding

            with Image.open(img_path) as img:
                img.thumbnail((width, height), Image.Resampling.LANCZOS)
                w, h = img.size
                tx = x + (width - w) // 2
                ty = y_thumbnail + (height - h) // 2
                grid.paste(img, (tx, ty))

                # Add border
                if BORDER_WIDTH > 0:
                    draw.rectangle(
                        [
                            (tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                            (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1),
                        ],
                        outline="gray",
                        width=BORDER_WIDTH,
                    )

        return grid

    def get_thumbnail_path(self, template_name: str, suffix: str = "") -> Path | None:
        """Get the path to a cached thumbnail without generating.

        Args:
            template_name: Name of the template.
            suffix: Optional suffix for the thumbnail filename.

        Returns:
            Path to thumbnail if it exists, None otherwise.
        """
        filename = f"{template_name}{suffix}.png"
        thumbnail_path = self.thumbnails_dir / filename
        if thumbnail_path.exists():
            return thumbnail_path
        return None


# Singleton instance
thumbnail_generator = ThumbnailGenerator()
