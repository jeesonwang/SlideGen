from __future__ import annotations

from dataclasses import dataclass
from enum import Enum

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from slidegen.services.presentation.design_tokens import DesignTokens
from slidegen.services.presentation.semantic import BlockSpec, SlideSpec


class VisualObjectKind(str, Enum):
    SHAPE = "shape"
    TEXT = "text"
    LINE = "line"


class TextSource(str, Enum):
    SLIDE_TITLE = "slide_title"
    BLOCK_TITLE = "block_title"
    BLOCK_TEXT = "block_text"
    BLOCK_TITLE_TEXT = "block_title_text"
    INDEX = "index"


_SHAPE_MAP: dict[str, MSO_AUTO_SHAPE_TYPE] = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
}

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


@dataclass(frozen=True)
class VisualObject:
    object_id: str
    kind: VisualObjectKind
    x_frac: float
    y_frac: float
    w_frac: float
    h_frac: float
    z_layer: int = 10
    group_id: str | None = None
    shape: str = "rect"
    fill: str | None = None
    stroke: str | None = None
    text: str | None = None
    text_source: TextSource | None = None
    block_index: int | None = None
    font_size: int | None = None
    font_color: str | None = None
    font_face: str | None = None
    bold: bool = False
    align: str = "left"

    def to_absolute(self, tokens: DesignTokens) -> tuple[float, float, float, float]:
        return (
            round(self.x_frac * tokens.slide_width, 2),
            round(self.y_frac * tokens.slide_height, 2),
            round(self.w_frac * tokens.slide_width, 2),
            round(self.h_frac * tokens.slide_height, 2),
        )


@dataclass(frozen=True)
class VisualPlan:
    name: str
    objects: tuple[VisualObject, ...]

    def validate(self) -> None:
        if not self.objects:
            raise ValueError("VisualPlan must contain at least one object")
        seen_ids: set[str] = set()
        for obj in self.objects:
            if obj.object_id in seen_ids:
                raise ValueError(f"Duplicate VisualObject id: {obj.object_id}")
            seen_ids.add(obj.object_id)
            if obj.x_frac < 0 or obj.y_frac < 0:
                raise ValueError(f"VisualObject {obj.object_id} starts outside canvas")
            if obj.w_frac <= 0 or obj.h_frac <= 0:
                raise ValueError(f"VisualObject {obj.object_id} must have positive size")
            if obj.x_frac + obj.w_frac > 1.001:
                raise ValueError(f"VisualObject {obj.object_id} right edge out of canvas")
            if obj.y_frac + obj.h_frac > 1.001:
                raise ValueError(f"VisualObject {obj.object_id} bottom edge out of canvas")


class VisualPlanRenderer:
    def __init__(self, tokens: DesignTokens):
        self.tokens = tokens

    async def render(self, slide: Slide, plan: VisualPlan, spec: SlideSpec) -> None:
        plan.validate()
        for obj in sorted(plan.objects, key=lambda item: item.z_layer):
            if obj.kind == VisualObjectKind.SHAPE:
                self._render_shape(slide, obj)
            elif obj.kind == VisualObjectKind.TEXT:
                self._render_text(slide, obj, spec)
            elif obj.kind == VisualObjectKind.LINE:
                self._render_line(slide, obj)

    def _render_shape(self, slide: Slide, obj: VisualObject) -> None:
        left, top, width, height = obj.to_absolute(self.tokens)
        shape_type = _SHAPE_MAP.get(obj.shape, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        if obj.fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(obj.fill)
        else:
            shape.fill.background()
        if obj.stroke:
            shape.line.color.rgb = _rgb(obj.stroke)
        else:
            shape.line.fill.background()

    def _render_line(self, slide: Slide, obj: VisualObject) -> None:
        left, top, width, height = obj.to_absolute(self.tokens)
        line = slide.shapes.add_connector(
            1,
            Inches(left),
            Inches(top),
            Inches(left + width),
            Inches(top + height),
        )
        if obj.stroke:
            line.line.color.rgb = _rgb(obj.stroke)

    def _render_text(self, slide: Slide, obj: VisualObject, spec: SlideSpec) -> None:
        left, top, width, height = obj.to_absolute(self.tokens)
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = frame.paragraphs[0]
        para.text = self._resolve_text(obj, spec)
        para.alignment = _ALIGN_MAP.get(obj.align, PP_ALIGN.LEFT)
        para.font.name = obj.font_face or self.tokens.body_font
        para.font.size = Pt(obj.font_size or self.tokens.body_size)
        para.font.bold = obj.bold
        para.font.color.rgb = _rgb(obj.font_color or self.tokens.text_primary)

    def _resolve_text(self, obj: VisualObject, spec: SlideSpec) -> str:
        if obj.text is not None:
            return obj.text
        if obj.text_source == TextSource.SLIDE_TITLE:
            return spec.title
        block = _block_at(spec, obj.block_index)
        if block is None:
            return ""
        if obj.text_source == TextSource.BLOCK_TITLE:
            return block.title
        if obj.text_source == TextSource.BLOCK_TEXT:
            return block.text
        if obj.text_source == TextSource.BLOCK_TITLE_TEXT:
            return f"{block.title}\n{block.text}".strip()
        if obj.text_source == TextSource.INDEX:
            return f"{(obj.block_index or 0) + 1:02d}"
        return ""


def build_composite_points_visual_plan(spec: SlideSpec, tokens: DesignTokens) -> VisualPlan:
    blocks = spec.blocks[:4]
    objects: list[VisualObject] = [
        VisualObject(
            object_id="background",
            kind=VisualObjectKind.SHAPE,
            x_frac=0,
            y_frac=0,
            w_frac=1,
            h_frac=1,
            z_layer=0,
            fill=tokens.light_bg,
        ),
        VisualObject(
            object_id="title_accent",
            kind=VisualObjectKind.SHAPE,
            x_frac=0.06,
            y_frac=0.09,
            w_frac=0.012,
            h_frac=0.095,
            z_layer=5,
            shape="rounded_rect",
            fill=tokens.accent,
        ),
        VisualObject(
            object_id="slide_title",
            kind=VisualObjectKind.TEXT,
            x_frac=0.085,
            y_frac=0.072,
            w_frac=0.76,
            h_frac=0.12,
            z_layer=20,
            text_source=TextSource.SLIDE_TITLE,
            font_face=tokens.title_font,
            font_size=30,
            font_color=tokens.primary,
            bold=True,
        ),
    ]

    for index, block in enumerate(blocks):
        x, y, w, h = _slot_for(index, len(blocks))
        group_id = f"point_{index}"
        objects.extend(_point_objects(index, block, x, y, w, h, group_id, tokens))

    return VisualPlan(name="CompositePointsVisualPlan", objects=tuple(objects))


def _point_objects(
    index: int,
    block: BlockSpec,
    x: float,
    y: float,
    w: float,
    h: float,
    group_id: str,
    tokens: DesignTokens,
) -> list[VisualObject]:
    del block
    number_size = min(0.072, h * 0.38)
    return [
        VisualObject(
            object_id=f"{group_id}_shadow",
            group_id=group_id,
            kind=VisualObjectKind.SHAPE,
            x_frac=x + 0.008,
            y_frac=y + 0.012,
            w_frac=w,
            h_frac=h,
            z_layer=3,
            shape="rounded_rect",
            fill=tokens.light_bg_alt,
        ),
        VisualObject(
            object_id=f"{group_id}_card",
            group_id=group_id,
            kind=VisualObjectKind.SHAPE,
            x_frac=x,
            y_frac=y,
            w_frac=w,
            h_frac=h,
            z_layer=5,
            shape="rounded_rect",
            fill="#FFFFFF",
            stroke=tokens.light_bg_alt,
        ),
        VisualObject(
            object_id=f"{group_id}_bar",
            group_id=group_id,
            kind=VisualObjectKind.SHAPE,
            x_frac=x,
            y_frac=y,
            w_frac=0.012,
            h_frac=h,
            z_layer=6,
            shape="rounded_rect",
            fill=tokens.accent if index % 2 == 0 else tokens.primary,
        ),
        VisualObject(
            object_id=f"{group_id}_number_bg",
            group_id=group_id,
            kind=VisualObjectKind.SHAPE,
            x_frac=x + 0.032,
            y_frac=y + 0.06,
            w_frac=number_size,
            h_frac=number_size,
            z_layer=10,
            shape="ellipse",
            fill=tokens.primary,
        ),
        VisualObject(
            object_id=f"{group_id}_number",
            group_id=group_id,
            kind=VisualObjectKind.TEXT,
            x_frac=x + 0.032,
            y_frac=y + 0.06 + number_size * 0.22,
            w_frac=number_size,
            h_frac=number_size * 0.48,
            z_layer=20,
            text_source=TextSource.INDEX,
            block_index=index,
            font_size=13,
            font_color=tokens.text_on_dark,
            bold=True,
            align="center",
        ),
        VisualObject(
            object_id=f"{group_id}_title",
            group_id=group_id,
            kind=VisualObjectKind.TEXT,
            x_frac=x + 0.12,
            y_frac=y + 0.048,
            w_frac=w - 0.16,
            h_frac=0.07,
            z_layer=20,
            text_source=TextSource.BLOCK_TITLE,
            block_index=index,
            font_face=tokens.title_font,
            font_size=18,
            font_color=tokens.primary,
            bold=True,
        ),
        VisualObject(
            object_id=f"{group_id}_body",
            group_id=group_id,
            kind=VisualObjectKind.TEXT,
            x_frac=x + 0.12,
            y_frac=y + 0.13,
            w_frac=w - 0.16,
            h_frac=max(0.06, h - 0.17),
            z_layer=20,
            text_source=TextSource.BLOCK_TEXT,
            block_index=index,
            font_size=12,
            font_color=tokens.text_secondary,
        ),
    ]


def _slot_for(index: int, count: int) -> tuple[float, float, float, float]:
    if count <= 1:
        return (0.09, 0.27, 0.82, 0.46)
    if count == 2:
        return (0.08 + index * 0.43, 0.29, 0.38, 0.42)
    if count == 3:
        slots = ((0.08, 0.27, 0.38, 0.23), (0.54, 0.27, 0.38, 0.23), (0.31, 0.56, 0.38, 0.23))
        return slots[index]
    slots = (
        (0.08, 0.26, 0.38, 0.23),
        (0.54, 0.26, 0.38, 0.23),
        (0.08, 0.55, 0.38, 0.23),
        (0.54, 0.55, 0.38, 0.23),
    )
    return slots[index]


def _block_at(spec: SlideSpec, index: int | None) -> BlockSpec | None:
    if index is None:
        return None
    if 0 <= index < len(spec.blocks):
        return spec.blocks[index]
    return None


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.removeprefix("#"))
