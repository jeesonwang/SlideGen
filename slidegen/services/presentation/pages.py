import copy
import hashlib
import io
import os
import random
import tempfile
import unicodedata
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Any, cast

from loguru import logger
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Length, Pt

from slidegen.core.config import settings
from slidegen.exceptions import CatalogTemplateNotFoundError, PPTGenError, PPTTemplateError
from slidegen.schemas.image_prompt import ImagePrompt
from slidegen.schemas.theme import PresentationTheme
from slidegen.services.document.markdown import Heading
from slidegen.services.presentation.components import (
    ChapterLayout,
    ComponentContentType,
    Location,
    Style,
    components_manager,
)
from slidegen.services.presentation.icon_searcher import IconSearcher, icon_searcher
from slidegen.services.presentation.image_generator import ImageGenerator
from slidegen.utils.slide import (
    add_para_by_xml,
    add_shape_by_xml,
    convert_paragraph_xml,
    runs_merge,
)

EMU_PER_PT = 12700
TITLE_TEXTBOX_HORIZONTAL_PADDING_EMU = 228600
TITLE_TEXTBOX_RIGHT_MARGIN_EMU = 457200
CATALOG_DEFAULT_LEFT_EMU = 914400
CATALOG_DEFAULT_TOP_EMU = 1371600
CATALOG_DEFAULT_NUMBER_WIDTH_EMU = 457200
CATALOG_DEFAULT_ITEM_HEIGHT_EMU = 342900
CATALOG_DEFAULT_TEXT_GAP_EMU = 228600
CATALOG_DEFAULT_ITEM_GAP_EMU = 228600
CATALOG_DEFAULT_RIGHT_MARGIN_EMU = 914400
CATALOG_DEFAULT_NUMBER_FONT_SIZE_PT = 18
CATALOG_DEFAULT_TEXT_FONT_SIZE_PT = 18
CATALOG_DEFAULT_NUMBER_FONT_BOLD = True
CATALOG_DEFAULT_NUMBER_ALIGNMENT = PP_ALIGN.CENTER


class Page:
    """PPT pages base class"""

    @staticmethod
    def _set_text(shape: Shape, text: str) -> None:
        """Set the text of the shape, keep the original paragraph style."""
        assert shape.has_text_frame, "Shape must have a text frame"
        if shape.is_placeholder:
            shape.text = text
            return
        tf = shape.text_frame
        # if the shape has text, merge the runs and set the run text
        if shape.text:
            para = tf.paragraphs[0]
            run = runs_merge(para)
            if run:
                run.text = text
        # if the shape has no text, add a new paragraph and keep the original paragraph style
        else:
            para = tf.paragraphs[0]
            para_xml = convert_paragraph_xml(para._element.xml, text)
            shape = add_para_by_xml(shape, para_xml)

    @staticmethod
    def remove_slide(prs: Presentation, index: int) -> None:
        """Delete the slide at the given index"""
        rId = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rId)
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(list(xml_slides)[index])

    @staticmethod
    def remove_shapes(sp_tree: CT_GroupShape, shapes: list[Shape]) -> None:
        for shp in shapes:
            if shp:
                sp_tree.remove(shp.element)

    @staticmethod
    def bring_shape_to_front(shape: BaseShape) -> None:
        """Move a shape to the top of the slide z-order."""
        element = shape.element
        parent = element.getparent()
        if parent is None:
            return
        parent.remove(element)
        if hasattr(parent, "insert_element_before"):
            parent.insert_element_before(element, "p:extLst")
        else:
            parent.append(element)

    @staticmethod
    def _shape_alignment(shape: BaseShape | Shape) -> None:
        """Set the alignment of the shape. Uniformly justify the text in the shape"""
        if shape.has_text_frame:
            tf = shape.text_frame  # type: ignore
            tf.vertical_anchor = MSO_ANCHOR.TOP
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY

    @staticmethod
    def _set_text_style(shape: Shape, style: dict[str, Any]) -> None:
        """Set the input `Shape` text style"""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        paragraph = tf.paragraphs[0]
        for key, value in style.items():
            if value is None:
                continue
            setattr(paragraph.font, key, value)

    @staticmethod
    def _get_text_frame_font_size_pt(shape: BaseShape) -> float | None:
        if not shape.has_text_frame:
            return None
        for paragraph in cast(Shape, shape).text_frame.paragraphs:
            if paragraph.font.size is not None:
                return paragraph.font.size.pt
            for run in paragraph.runs:
                if run.font.size is not None:
                    return run.font.size.pt
        return None

    @staticmethod
    def _estimate_single_line_text_width(text: str, font_size_pt: float) -> int:
        width_units = 0.0
        for char in text:
            if char.isspace():
                width_units += 0.35
            elif unicodedata.east_asian_width(char) in {"F", "W"}:
                width_units += 1.0
            elif unicodedata.east_asian_width(char) == "A":
                width_units += 0.85
            elif char.isupper():
                width_units += 0.68
            else:
                width_units += 0.56
        return int(width_units * font_size_pt * EMU_PER_PT + TITLE_TEXTBOX_HORIZONTAL_PADDING_EMU)

    @staticmethod
    def _expand_title_text_box(slide: Slide, shape: BaseShape, text: str) -> None:
        """Expand a no-wrap title shape horizontally when the title is likely too long."""
        if not shape.has_text_frame or not text.strip():
            return

        font_size_pt = Page._get_text_frame_font_size_pt(shape)
        if font_size_pt is None:
            font_size_pt = max(18.0, min(44.0, shape.height / EMU_PER_PT / 1.2))

        required_width = Page._estimate_single_line_text_width(text, font_size_pt)
        if required_width <= shape.width:
            return

        prs = slide.part.package.presentation_part.presentation
        right_boundary = max(0, prs.slide_width - TITLE_TEXTBOX_RIGHT_MARGIN_EMU)
        max_width = max(0, right_boundary - shape.left)
        target_width = min(required_width, max_width)
        if target_width > shape.width:
            shape.width = target_width

    @staticmethod
    def _scale_shape_location(
        prs: Presentation,
        loc: Location,
    ) -> Location:
        metadata = getattr(components_manager, "metadata", {}) or {}
        source_width = metadata.get("slide_width")
        source_height = metadata.get("slide_height")
        target_width = prs.slide_width
        target_height = prs.slide_height

        if not source_width or not source_height or not target_width or not target_height:
            return loc

        x_scale = int(target_width) / int(source_width)
        y_scale = int(target_height) / int(source_height)
        scaled_x = round(loc.x * x_scale)
        scaled_y = round(loc.y * y_scale)
        scaled_width = round(loc.width * x_scale)
        scaled_height = round(loc.height * y_scale)

        if loc.width == loc.height:
            square_size = min(scaled_width, scaled_height)
            scaled_x += round((scaled_width - square_size) / 2)
            scaled_y += round((scaled_height - square_size) / 2)
            scaled_width = square_size
            scaled_height = square_size

        return Location(
            x=scaled_x,
            y=scaled_y,
            width=scaled_width,
            height=scaled_height,
        )

    @staticmethod
    def _find_or_inject_placeholder(
        slide: Slide,
        page_type: str,
        role: str,
        text: str,
        *,
        placeholder_types: tuple[int, ...],
        shape_name_prefix: str = "injected",
    ) -> BaseShape | None:
        """Find a matching placeholder on the slide, or inject one from shape.json.

        Args:
            slide: The slide to operate on.
            page_type: Key for page_placeholders lookup (e.g. "cover").
            role: Role key for page_placeholders lookup (e.g. "title").
            text: Text content to set on the shape.
            placeholder_types: Tuple of PP_PLACEHOLDER int values to match.
            shape_name_prefix: Prefix for the injected shape name.

        Returns:
            The Shape if found or injected and text was set, None otherwise.
        """
        for ph in slide.shapes.placeholders:
            if ph.placeholder_format.type in placeholder_types:
                Page._set_text(ph, text)
                if ph.has_text_frame:
                    ph.text_frame.word_wrap = False
                    Page._expand_title_text_box(slide, ph, text)
                return ph

        ph_data = components_manager.get_page_placeholder(page_type, role)
        if ph_data is not None:
            prs = slide.part.package.presentation_part.presentation
            loc = Page._scale_shape_location(prs, ph_data.location)
            injected = ShapeFactory.create_text_shape(
                slide,
                ph_data.xml,
                text,
                loc,
                shape_id=slide.shapes._next_shape_id,
                shape_name=f"{shape_name_prefix}_{role}",
            )
            if injected.has_text_frame:
                injected.text_frame.word_wrap = False
            Page._expand_title_text_box(slide, injected, text)
            return injected

        return None

    @staticmethod
    def move_slide(pres: Presentation, slide: Slide, index: int) -> None:
        """
        Move the slide to the new index

        Args:
            pres: Presentation object
            slide: Slide object to be moved
            index: New index of the slide
        """
        old_index = pres.slides.index(slide)
        xml_slides = pres.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(index, slides[old_index])

    @staticmethod
    def duplicate_slide(pres: Presentation, index: int) -> Slide:
        """Duplicate the slide at the given index

        Args:
            pres: Presentation object
            index: Index of the slide to be duplicated

        Returns:
            Copied slide
        """
        template = pres.slides[index]
        copied_slide = pres.slides.add_slide(template.slide_layout)
        # Delete the existing shapes that are part of the layout
        for shp in copied_slide.shapes:
            copied_slide.shapes.element.remove(shp.element)

        # Perform a deep copy of the shapes from the template
        for shp in template.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = io.BytesIO(shp.image.blob)
                copied_slide.shapes.add_picture(
                    image_file=img,
                    left=shp.left,
                    top=shp.top,
                    width=shp.width,
                    height=shp.height,
                )
            else:
                el = shp.element
                newel = copy.deepcopy(el)
                custDataLst = newel.xpath(".//p:custDataLst")
                for cd in custDataLst:
                    cd.getparent().remove(cd)
                copied_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

        return copied_slide

    @staticmethod
    async def generate_slide(*args: Any, **kwargs: Any) -> Any:
        raise NotImplementedError("Page.generate_slide is abstract; override in subclasses")


class ShapeFactory:
    """Factory helpers for creating shapes on slides."""

    @staticmethod
    def create_xml_shape(
        slide: Slide,
        xml: str,
        location: Location,
        *,
        shape_id: int | None = None,
        shape_name: str = "shape",
        text: str | None = None,
    ) -> BaseShape:
        """Create a shape from XML at the given location."""
        return add_shape_by_xml(
            slide=slide,
            shape_xml=xml,
            shape_id=shape_id if shape_id is not None else slide.shapes._next_shape_id,
            shape_name=shape_name,
            text_content=text if text is not None else "",
            location=location,
        )

    @staticmethod
    def create_text_shape(
        slide: Slide,
        xml: str,
        text: str,
        location: Location,
        *,
        shape_id: int | None = None,
        shape_name: str = "text",
    ) -> Shape:
        """Create a text shape from XML and set its text consistently."""
        shape = ShapeFactory.create_xml_shape(
            slide,
            xml,
            location,
            shape_id=shape_id,
            shape_name=shape_name,
            text=text,
        )
        text_shape = cast(Shape, shape)
        Page._set_text(text_shape, text)
        return text_shape

    @staticmethod
    def create_image_shape(slide: Slide, image_path: str, location: Location) -> BaseShape:
        """Create an image shape at the given location."""
        return slide.shapes.add_picture(
            image_path,
            Length(location.x),
            Length(location.y),
            Length(location.width),
            Length(location.height),
        )


class CoverPage(Page):
    """Presentation cover page"""

    @staticmethod
    async def generate_slide(prs: Presentation, content: Heading, *, cover_page_index: int = 0) -> None:
        """
        Generate the cover page

        Args:
            prs: Presentation object
            content: Heading object, the main heading of the markdown document(level 1)
            cover_page_index: index of the cover page
        """
        cover_page = prs.slides[cover_page_index]
        assert content.level == 1, "Cover page must have a level 1 heading"
        main_title = content.element_text
        if not main_title.strip():
            main_title = "Presentation Title"
        # TODO: add subtitle
        title_shape = Page._find_or_inject_placeholder(
            slide=cover_page,
            page_type="cover",
            role="title",
            text=main_title,
            placeholder_types=(PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE),
        )
        if not title_shape:
            raise PPTTemplateError(
                f"{CoverPage.__name__}: No title placeholder found in cover slide "
                f"and no fallback data in shape.json page_placeholders.cover"
            )


class CatalogLayout(Enum):
    """
    Catalog layout enum
    """

    VERTICAL = "vertical"
    HORIZONTAL = "horizontal"
    UNDEFINED = "undefined"


class CatalogItem:
    """
    Catalog item including number shape, text shape and background shape.
    """

    def __init__(
        self, number_shape: dict[str, Any], text_shape: dict[str, Any], background_shape: dict[str, Any] | None = None
    ):
        self.number_shape = number_shape
        self.text_shape = text_shape
        self.background_shape = background_shape

    def asdict(self) -> dict[str, Any]:
        return {
            "number_shape": self.number_shape,
            "text_shape": self.text_shape,
            "background_shape": self.background_shape,
        }


class CatalogList(list[CatalogItem]):
    """
    Catalog list including a list of `CatalogItem`.
    """

    def asdict(self) -> list[dict[str, Any]]:
        return [item.asdict() for item in self]


class CatalogPage(Page):
    """Presentation catalog page"""

    # vertical tolerance coefficient
    vertical_tolerance = 1.5

    @staticmethod
    def _calculate_distance(shape1: dict[str, Any], shape2: dict[str, Any]) -> float:
        """Calculate the distance between two shapes"""
        return ((shape1["left"] - shape2["left"]) ** 2 + (shape1["top"] - shape2["top"]) ** 2) ** 0.5

    @staticmethod
    def _shape_info(shape: BaseShape) -> dict[str, Any]:
        """Extract shape information into a dictionary for catalog processing.

        Args:
            shape: The shape to extract information from.

        Returns:
            Dictionary containing shape properties: text, position, dimensions, type, and reference.
        """
        return {
            "text": cast(Shape, shape).text.strip() if shape.has_text_frame else None,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "shape_type": shape.shape_type,
            "shape_id": shape.shape_id,
            "shape": shape,
        }

    # Roman numerals 1-20
    _ROMAN_NUMERALS: set[str] = {
        "I",
        "II",
        "III",
        "IV",
        "V",
        "VI",
        "VII",
        "VIII",
        "IX",
        "X",
        "XI",
        "XII",
        "XIII",
        "XIV",
        "XV",
        "XVI",
        "XVII",
        "XVIII",
        "XIX",
        "XX",
    }

    _ROMAN_TO_INT: dict[str, int] = {
        "I": 1,
        "II": 2,
        "III": 3,
        "IV": 4,
        "V": 5,
        "VI": 6,
        "VII": 7,
        "VIII": 8,
        "IX": 9,
        "X": 10,
        "XI": 11,
        "XII": 12,
        "XIII": 13,
        "XIV": 14,
        "XV": 15,
        "XVI": 16,
        "XVII": 17,
        "XVIII": 18,
        "XIX": 19,
        "XX": 20,
    }

    @staticmethod
    def _is_chapter_number(text: str) -> bool:
        """Check if the given text represents a chapter number.

        Supports three formats:
        - Pure digits (including leading zeros): "1", "01", "99"
        - "N." format: "1.", "01."
        - Roman numerals 1-20: "I", "IV", "VIII", "XX"
        """
        if not text or len(text) > 5:
            return False
        if text.isdigit():
            return True
        if text.endswith(".") and text[:-1].isdigit():
            return True
        if text.isalpha() and text.isupper():
            return text in CatalogPage._ROMAN_NUMERALS
        return False

    @staticmethod
    def _chapter_number_sort_key(shape_info: dict[str, Any]) -> int:
        """Return an integer sort key for a chapter number shape.

        Handles three formats:
        - "N." format: strips trailing dot, returns int
        - Pure digits: returns int
        - Roman numerals: returns mapped int value
        """
        text: str = shape_info["text"].strip()
        if text.endswith("."):
            text = text[:-1]
        if text.isdigit():
            return int(text)
        if text in CatalogPage._ROMAN_TO_INT:
            return CatalogPage._ROMAN_TO_INT[text]
        # Fallback: should not happen if shape passed _is_chapter_number()
        return 9999

    @staticmethod
    def _calculate_max_per_page(
        catalog_items: "CatalogList",
        layout_direction: "CatalogLayout",
        slide_height: int,
        slide_width: int,
    ) -> int:
        """Calculate max catalog items per page based on spacing and margins.

        Uses the average spacing between existing items and the distance from the
        first item to the slide edge to determine how many items fit on one page.
        Falls back to len(catalog_items) if spacing can't be determined (< 2 items
        or undefined layout).
        """
        if len(catalog_items) == 0 or layout_direction == CatalogLayout.UNDEFINED:
            return len(catalog_items)

        step = CatalogPage._calculate_catalog_step(catalog_items, layout_direction)
        if step <= 0:
            return len(catalog_items)

        first_item = catalog_items[0]
        if layout_direction == CatalogLayout.VERTICAL:
            usable = slide_height - min(first_item.number_shape["top"], first_item.text_shape["top"])
        else:  # HORIZONTAL
            usable = slide_width - min(first_item.number_shape["left"], first_item.text_shape["left"])

        return max(len(catalog_items), int(usable / step))

    @staticmethod
    def _calculate_catalog_step(
        catalog_items: "CatalogList",
        layout_direction: "CatalogLayout",
    ) -> int:
        """Calculate the spacing between catalog items along the layout direction.

        For multiple items, computes average spacing from existing positions.
        For a single item, estimates spacing from item dimensions plus default gap.

        Args:
            catalog_items: List of catalog items with position information.
            layout_direction: Direction in which items are arranged.

        Returns:
            Step size in EMU (English Metric Units), or 0 if calculation is not possible.
        """
        if len(catalog_items) >= 2:
            key = "top" if layout_direction == CatalogLayout.VERTICAL else "left"
            positions = [item.number_shape[key] for item in catalog_items]
            return int(
                sum(abs(positions[i + 1] - positions[i]) for i in range(len(positions) - 1))
                / (len(positions) - 1)
            )

        if len(catalog_items) != 1:
            return 0

        item = catalog_items[0]
        shapes = [item.number_shape, item.text_shape]
        if item.background_shape:
            shapes.append(item.background_shape)

        if layout_direction == CatalogLayout.VERTICAL:
            start_pos = min(shape["top"] for shape in shapes)
            end_pos = max(shape["top"] + shape["height"] for shape in shapes)
        elif layout_direction == CatalogLayout.HORIZONTAL:
            start_pos = min(shape["left"] for shape in shapes)
            end_pos = max(shape["left"] + shape["width"] for shape in shapes)
        else:
            return 0

        return max(
            CATALOG_DEFAULT_ITEM_HEIGHT_EMU + CATALOG_DEFAULT_ITEM_GAP_EMU,
            end_pos - start_pos + CATALOG_DEFAULT_ITEM_GAP_EMU,
        )

    @staticmethod
    def _infer_single_item_layout_direction(item: CatalogItem) -> CatalogLayout:
        """Infer layout direction from the spatial relationship between number and text shapes.

        The logic maps the relationship between number and text shapes to the overall
        catalog layout direction:
        - If number and text are horizontally separated (side-by-side), the catalog
          items are meant to stack VERTICALLY (one row per item).
        - If number and text are vertically separated (one above the other), the
          catalog items are arranged HORIZONTALLY (one column per item).

        Args:
            item: A single catalog item with number and text shape information.

        Returns:
            CatalogLayout.VERTICAL if items should stack vertically,
            CatalogLayout.HORIZONTAL if items should arrange horizontally.
        """
        number_center_x = item.number_shape["left"] + item.number_shape["width"] / 2
        number_center_y = item.number_shape["top"] + item.number_shape["height"] / 2
        text_center_x = item.text_shape["left"] + item.text_shape["width"] / 2
        text_center_y = item.text_shape["top"] + item.text_shape["height"] / 2

        horizontal_distance = abs(text_center_x - number_center_x)
        vertical_distance = abs(text_center_y - number_center_y)

        # Horizontal separation between number and text → vertical item stacking
        if horizontal_distance >= vertical_distance:
            return CatalogLayout.VERTICAL
        # Vertical separation between number and text → horizontal item arrangement
        return CatalogLayout.HORIZONTAL

    @staticmethod
    def _resolve_layout_direction(catalog_items: "CatalogList") -> CatalogLayout:
        """Determine the layout direction of catalog items.

        Uses different strategies based on the number of items available:
        - Multiple items: analyzes position patterns to determine direction.
        - Single item: infers direction from number-text spatial relationship.
        - No items: returns UNDEFINED.

        Args:
            catalog_items: List of catalog items to analyze.

        Returns:
            The determined layout direction (VERTICAL, HORIZONTAL, or UNDEFINED).
        """
        number_shapes = [item.number_shape for item in catalog_items]
        if len(number_shapes) >= 2:
            return CatalogPage._layout_direction(number_shapes)
        if len(catalog_items) == 1:
            return CatalogPage._infer_single_item_layout_direction(catalog_items[0])
        return CatalogLayout.UNDEFINED

    @staticmethod
    def _create_default_catalog_items(slide: Slide) -> CatalogList:
        prs = slide.part.package.presentation_part.presentation
        slide_width = int(prs.slide_width)
        text_left = CATALOG_DEFAULT_LEFT_EMU + CATALOG_DEFAULT_NUMBER_WIDTH_EMU + CATALOG_DEFAULT_TEXT_GAP_EMU
        text_width = max(
            CATALOG_DEFAULT_NUMBER_WIDTH_EMU,
            slide_width - text_left - CATALOG_DEFAULT_RIGHT_MARGIN_EMU,
        )

        number_shape = slide.shapes.add_textbox(
            Length(CATALOG_DEFAULT_LEFT_EMU),
            Length(CATALOG_DEFAULT_TOP_EMU),
            Length(CATALOG_DEFAULT_NUMBER_WIDTH_EMU),
            Length(CATALOG_DEFAULT_ITEM_HEIGHT_EMU),
        )
        number_shape.text = "01"
        number_paragraph = number_shape.text_frame.paragraphs[0]
        number_paragraph.alignment = CATALOG_DEFAULT_NUMBER_ALIGNMENT
        number_paragraph.font.size = Pt(CATALOG_DEFAULT_NUMBER_FONT_SIZE_PT)
        number_paragraph.font.bold = CATALOG_DEFAULT_NUMBER_FONT_BOLD

        text_shape = slide.shapes.add_textbox(
            Length(text_left),
            Length(CATALOG_DEFAULT_TOP_EMU),
            Length(text_width),
            Length(CATALOG_DEFAULT_ITEM_HEIGHT_EMU),
        )
        text_shape.text = "Catalog Item"
        text_paragraph = text_shape.text_frame.paragraphs[0]
        text_paragraph.font.size = Pt(CATALOG_DEFAULT_TEXT_FONT_SIZE_PT)

        return CatalogList(
            [
                CatalogItem(
                    CatalogPage._shape_info(number_shape),
                    CatalogPage._shape_info(text_shape),
                )
            ]
        )

    @staticmethod
    def _get_or_create_catalog_items(slide: Slide) -> CatalogList:
        """Retrieve catalog items from the slide template, or create default items as fallback.

        Args:
            slide: The catalog slide to process.

        Returns:
            List of catalog items found or created.

        Raises:
            PPTTemplateError: If a non-recoverable template error occurs.
        """
        try:
            return CatalogPage._get_catalog_items(slide)
        except CatalogTemplateNotFoundError:
            logger.info("Catalog slide has no template items; creating default catalog item fallback")
            return CatalogPage._create_default_catalog_items(slide)

    @staticmethod
    def _clone_shape_to_slide(
        sp_tree: Any,
        slide: Slide,
        src_shape: Any,
        dx: int,
        dy: int,
    ) -> Any:
        """Clone a shape element, reposition it, and insert into the slide.

        Args:
            sp_tree: The slide's shape tree (CT_GroupShape from python-pptx internals).
            slide: Target slide for the cloned shape.
            src_shape: Source shape to clone.
            dx: Horizontal offset in EMU.
            dy: Vertical offset in EMU.

        Returns:
            The new Shape object wrapper.

        Note:
            Type annotations use Any due to python-pptx internal types not being
            fully exposed in the public API.
        """
        new_el = copy.deepcopy(src_shape.element)
        # Remove custDataLst to avoid conflicts
        for cd in new_el.xpath(".//p:custDataLst"):
            cd.getparent().remove(cd)
        # Assign fresh non-visual ids so cloned shapes remain valid in the same slide.
        next_shape_id = slide.shapes._next_shape_id
        for offset, c_nv_pr in enumerate(new_el.xpath(".//p:cNvPr")):
            c_nv_pr.set("id", str(next_shape_id + offset))
        # Adjust position
        off_el = new_el.xpath(".//a:off")[0]
        new_left = int(off_el.get("x", "0")) + dx
        new_top = int(off_el.get("y", "0")) + dy
        off_el.set("x", str(new_left))
        off_el.set("y", str(new_top))

        sp_tree.insert_element_before(new_el, "p:extLst")

        # Find the newly created Shape wrapper by matching the element
        for s in slide.shapes:
            if s._element is new_el:
                return s
        raise PPTGenError("Failed to find cloned shape after insertion")

    @staticmethod
    def _layout_direction(number_shapes: list[dict[str, Any]]) -> CatalogLayout:
        """Judge the layout direction of the catalog page"""
        if len(number_shapes) < 2:
            raise PPTTemplateError("To judge the layout direction, catalog page must have at least two chapter numbers")

        sorted_numbers = sorted(number_shapes, key=lambda x: (x["left"], x["top"]))
        horizontal_diffs = []
        vertical_diffs = []

        for i in range(len(sorted_numbers) - 1):
            horizontal_diffs.append(abs(sorted_numbers[i + 1]["left"] - sorted_numbers[i]["left"]))
            vertical_diffs.append(abs(sorted_numbers[i + 1]["top"] - sorted_numbers[i]["top"]))
        avg_horizontal_diff = sum(horizontal_diffs) / len(horizontal_diffs) if horizontal_diffs else 0
        avg_vertical_diff = sum(vertical_diffs) / len(vertical_diffs) if vertical_diffs else 0

        if avg_horizontal_diff > avg_vertical_diff:
            return CatalogLayout.HORIZONTAL
        else:
            return CatalogLayout.VERTICAL

    @staticmethod
    def _get_catalog_items(slide: Slide) -> CatalogList:
        number_shapes: list[dict[str, Any]] = []
        text_shapes: list[dict[str, Any]] = []
        all_shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                # placeholder shapes are not included in the all_shapes list
                continue
            shape_info = CatalogPage._shape_info(shape)
            if shape.has_text_frame:
                text_shapes.append(shape_info)
            all_shapes.append(shape_info)
        for shape_info in text_shapes:
            # check if the text is a chapter number
            text = shape_info["text"].strip()
            if CatalogPage._is_chapter_number(text):
                number_shapes.append(shape_info)
        try:
            number_shapes.sort(key=CatalogPage._chapter_number_sort_key)
        except ValueError:
            raise PPTTemplateError("Chapter number must be a number")

        match len(number_shapes):
            case 0:
                raise CatalogTemplateNotFoundError("Catalog page must have at least one chapter numbers")
            case 1:
                layout_direction = CatalogLayout.UNDEFINED
            case _:
                layout_direction = CatalogPage._layout_direction(number_shapes)

        except_number_shapes = [shape for shape in text_shapes if shape not in number_shapes]
        catalog_list = CatalogList()
        # Find the closest text shape for each number shape
        for number_shape in number_shapes:
            min_distance = float("inf")
            closest_text_shape = None

            for text_shape in except_number_shapes:
                if layout_direction == CatalogLayout.HORIZONTAL:
                    # For horizontal layout, find the text shape below the number shape
                    if text_shape["top"] > number_shape["top"]:
                        distance = CatalogPage._calculate_distance(number_shape, text_shape)
                        horizontal_overlap = min(
                            number_shape["left"] + number_shape["width"],
                            text_shape["left"] + text_shape["width"],
                        ) - max(number_shape["left"], text_shape["left"])
                        if horizontal_overlap > 0 and distance < min_distance:
                            min_distance = distance
                            closest_text_shape = text_shape
                elif layout_direction == CatalogLayout.VERTICAL:
                    # For vertical layout, find the text shape to the right of the number shape
                    if text_shape["left"] > number_shape["left"]:
                        distance = CatalogPage._calculate_distance(number_shape, text_shape)
                        vertical_overlap = min(
                            number_shape["top"] + number_shape["height"],
                            text_shape["top"] + text_shape["height"],
                        ) - max(number_shape["top"], text_shape["top"])
                        if vertical_overlap > 0 and distance < min_distance:
                            min_distance = distance
                            closest_text_shape = text_shape
                else:
                    distance = CatalogPage._calculate_distance(number_shape, text_shape)
                    if distance < min_distance:
                        min_distance = distance
                        closest_text_shape = text_shape

            if closest_text_shape:
                catalog_list.append(CatalogItem(number_shape, closest_text_shape))
                try:
                    all_shapes.remove(number_shape)
                    if closest_text_shape is not None:
                        all_shapes.remove(closest_text_shape)
                except ValueError:
                    raise PPTTemplateError(
                        f"all shape: {all_shapes}\n current closest_text_shape: {closest_text_shape} \n current number_shape: {number_shape}"
                    )
        assert len(number_shapes) == len(catalog_list), (
            "The number of chapter numbers and chapter titles must be the same"
        )

        if len(all_shapes) >= len(number_shapes):
            # continue to calculate the distance to find the background shape
            for i, number_shape in enumerate(number_shapes):
                min_distance = float("inf")
                closest_background_shape = None
                for shape_info in all_shapes:
                    distance = CatalogPage._calculate_distance(number_shape, shape_info)
                    if distance < min_distance:
                        min_distance = distance
                        if min_distance < shape_info["height"] * CatalogPage.vertical_tolerance:
                            closest_background_shape = shape_info
                if closest_background_shape:
                    catalog_list[i].background_shape = closest_background_shape

        return catalog_list

    @staticmethod
    async def generate_slide(
        prs: Presentation,
        content: list[Heading],
        *,
        catalog_page_index: int = 1,
        begin_number: int = 1,
    ) -> int:
        """
        Generate the catalog page

        Args:
            prs: Presentation object
            content: list of Heading objects
            catalog_page_index: index of the catalog page
            begin_number: starting number of the catalog page
        Returns:
            index of the catalog page
        """
        if not content:
            raise PPTGenError("Catalog page must have content.")
        catalog_num = len(content)
        catalog_slide = prs.slides[catalog_page_index]
        catalog_items = CatalogPage._get_or_create_catalog_items(catalog_slide)

        # Determine layout direction for position calculations
        layout_direction = CatalogPage._resolve_layout_direction(catalog_items)

        if len(catalog_items) > catalog_num:
            sp_tree = catalog_slide.shapes._spTree
            # delete the excess shape pairs from the slide
            excess_items = catalog_items[catalog_num:]
            for item in excess_items:
                Page.remove_shapes(
                    sp_tree,
                    [
                        item.number_shape["shape"],
                        item.text_shape["shape"],
                        *([item.background_shape["shape"]] if item.background_shape else []),
                    ],
                )
            catalog_items = catalog_items[:catalog_num]  # type: ignore
        elif len(catalog_items) < catalog_num:
            # Template has fewer catalog slots than content — clone shapes to fill the page
            slide_height = prs.slide_height
            slide_width = prs.slide_width
            if slide_height is None or slide_width is None:
                raise PPTTemplateError("Presentation slide dimensions must be defined for catalog pagination")
            max_per_page = CatalogPage._calculate_max_per_page(
                catalog_items,
                layout_direction,
                int(slide_height),
                int(slide_width),
            )
            sp_tree = catalog_slide.shapes._spTree
            source_item = catalog_items[-1]
            target_count = min(max_per_page, catalog_num)

            # Calculate position step from existing items
            step = CatalogPage._calculate_catalog_step(catalog_items, layout_direction)

            n_existing = len(catalog_items)
            for clone_idx in range(1, target_count - n_existing + 1):
                if layout_direction == CatalogLayout.VERTICAL:
                    dx, dy = 0, step * clone_idx
                elif layout_direction == CatalogLayout.HORIZONTAL:
                    dx, dy = step * clone_idx, 0
                else:
                    dx, dy = 0, 0

                # Clone background shape first so it stays below the number and text.
                new_bg_info = None
                if source_item.background_shape:
                    new_bg_wrapper = CatalogPage._clone_shape_to_slide(
                        sp_tree,
                        catalog_slide,
                        source_item.background_shape["shape"],
                        dx,
                        dy,
                    )
                    new_bg_info = {
                        "text": None,
                        "left": source_item.background_shape["left"] + dx,
                        "top": source_item.background_shape["top"] + dy,
                        "width": source_item.background_shape["width"],
                        "height": source_item.background_shape["height"],
                        "shape_type": source_item.background_shape["shape_type"],
                        "shape_id": new_bg_wrapper.shape_id,
                        "shape": new_bg_wrapper,
                    }

                # Clone number shape
                new_number_shape_wrapper = CatalogPage._clone_shape_to_slide(
                    sp_tree,
                    catalog_slide,
                    source_item.number_shape["shape"],
                    dx,
                    dy,
                )
                new_number_info = {
                    "text": "",
                    "left": source_item.number_shape["left"] + dx,
                    "top": source_item.number_shape["top"] + dy,
                    "width": source_item.number_shape["width"],
                    "height": source_item.number_shape["height"],
                    "shape_type": source_item.number_shape["shape_type"],
                    "shape_id": new_number_shape_wrapper.shape_id,
                    "shape": new_number_shape_wrapper,
                }

                # Clone text shape
                new_text_shape_wrapper = CatalogPage._clone_shape_to_slide(
                    sp_tree,
                    catalog_slide,
                    source_item.text_shape["shape"],
                    dx,
                    dy,
                )
                new_text_info = {
                    "text": "",
                    "left": source_item.text_shape["left"] + dx,
                    "top": source_item.text_shape["top"] + dy,
                    "width": source_item.text_shape["width"],
                    "height": source_item.text_shape["height"],
                    "shape_type": source_item.text_shape["shape_type"],
                    "shape_id": new_text_shape_wrapper.shape_id,
                    "shape": new_text_shape_wrapper,
                }

                new_item = CatalogItem(new_number_info, new_text_info, new_bg_info)
                catalog_items.append(new_item)

        # Fill catalog items with content
        fill_count = min(len(catalog_items), catalog_num)
        for i in range(fill_count):
            cur_content = content[i].element_text
            cur_number = begin_number
            text_shape = catalog_items[i].text_shape["shape"]
            number_shape = catalog_items[i].number_shape["shape"]

            catalog_items[i].text_shape["text"] = cur_content
            CatalogPage._set_text(text_shape, cur_content)
            # The chapter number is formatted as "01"
            chapter_number = str(cur_number).zfill(2)
            catalog_items[i].number_shape["text"] = chapter_number
            CatalogPage._set_text(number_shape, chapter_number)
            begin_number += 1

        # If content still remains, create a new catalog page
        if fill_count < catalog_num:
            new_catalog_slide = CatalogPage.duplicate_slide(prs, catalog_page_index)
            catalog_page_index += 1
            CatalogPage.move_slide(prs, new_catalog_slide, catalog_page_index)
            # Recursively generate the new catalog page
            return await CatalogPage.generate_slide(
                prs,
                content[fill_count:],
                catalog_page_index=catalog_page_index,
                begin_number=begin_number,
            )

        return catalog_page_index


class ChapterHomePage(Page):
    """Chapter home page"""

    selected_style: int | None = None

    @staticmethod
    async def generate_slide(
        prs: Presentation,
        content: Heading,
        *,
        chapter_home_page_index: int = 2,
        chapter_number: int = 1,
        slide_index: int = 2,
    ) -> None:
        """
        Generate the chapter home page

        Args:
            prs: Presentation object
            content: Heading object
            chapter_home_page_index: index of the chapter home page
            chapter_number: current chapter number, begin from 1
            slide_index: index of the slide to be generated
        """
        assert content.level == 2, f"{ChapterHomePage.__name__}: Chapter home page must input a level 2 heading"
        template_slide = prs.slides[chapter_home_page_index]
        chapter_home_slide = prs.slides.add_slide(template_slide.slide_layout)

        title = content.element_text
        title_shape = Page._find_or_inject_placeholder(
            slide=chapter_home_slide,
            page_type="chapter_home",
            role="title",
            text=title,
            placeholder_types=(PP_PLACEHOLDER.TITLE,),
            shape_name_prefix="chapter_title",
        )
        if not title_shape:
            raise PPTTemplateError(
                f"{ChapterHomePage.__name__}: No title placeholder found in chapter home slide "
                f"and no fallback data in shape.json page_placeholders.chapter_home"
            )
        chapter_number_shape = None
        min_distance = float("inf")
        for shape in chapter_home_slide.shapes:
            if shape == title_shape:
                continue
            if shape.has_text_frame:
                if shape.top < title_shape.top:
                    distance = title_shape.top - shape.top
                    if distance < min_distance:
                        shape_text = shape.text.strip()
                        if (
                            (shape_text.startswith("0") and shape_text[1:].isdigit())
                            or shape_text.lower().startswith("part")
                            or (shape_text.endswith(".") and shape_text[:-1].isdigit())
                        ):
                            chapter_number_shape = shape
                            break
                        min_distance = distance
                        chapter_number_shape = shape
        if chapter_number_shape:
            chapter_index = ChapterHomePage.convert_chapter_number(chapter_number)
            ChapterHomePage._set_text(chapter_number_shape, chapter_index)
        ChapterHomePage.move_slide(prs, chapter_home_slide, slide_index)

    @staticmethod
    def convert_chapter_number(chapter_number: int) -> str:
        """Randomly select a style for the chapter number"""
        import inflect

        # If the style has been selected, use it
        if ChapterHomePage.selected_style is not None:
            style_type = ChapterHomePage.selected_style
        else:
            style_type = random.randint(1, 3)
            ChapterHomePage.selected_style = style_type

        p = inflect.engine()
        if style_type == 1:
            return str(chapter_number).zfill(2)  # 01, 02, 03, ...
        elif style_type == 2:
            return f"PART {str(chapter_number).zfill(2)}"  # PART 01, PART 02, PART 03, ...
        else:
            # PART ONE, PART TWO, PART THREE, ...
            return f"PART {p.number_to_words(chapter_number).upper()}"  # type: ignore


@dataclass(frozen=True)
class ChapterSlideAsset:
    shape_name: str
    location_index: int
    content_type: ComponentContentType
    path: str


@dataclass(frozen=True)
class ChapterSlideData:
    content_title: str
    section_titles: list[str]
    section_texts: list[str]
    style: Style
    assets: dict[tuple[str, int, ComponentContentType], ChapterSlideAsset]


class ChapterContentPage(Page):
    """
    Chapter content page

    Divide the chapter content slides into one-point, two-point, three-point, and four-point slides.
    """

    image_generator: ImageGenerator = ImageGenerator(
        settings.TEMP_DIRECTORY or os.path.join(os.getcwd(), "generated_images")
    )
    icon_searcher: IconSearcher = icon_searcher

    @staticmethod
    def _resolve_placeholder_image_path() -> str | None:
        placeholder_rel = os.path.join("components", "icons", "placeholder.png")
        placeholder_abs = os.path.join(Path(__file__).resolve().parents[3], placeholder_rel)
        if os.path.exists(placeholder_abs):
            return placeholder_abs
        if os.path.exists(placeholder_rel):
            return placeholder_rel
        return None

    @staticmethod
    def _normalize_icon_color(hex_color: str | None) -> str | None:
        if not hex_color:
            return None

        color = hex_color.strip()
        if color.startswith("#"):
            color = color[1:]
        if color.lower().startswith("0x"):
            color = color[2:]

        if len(color) != 6:
            logger.warning("Invalid icon theme color '{}': expected 6 hex digits", hex_color)
            return None

        try:
            int(color, 16)
        except ValueError:
            logger.warning("Invalid icon theme color '{}': not a hex color", hex_color)
            return None

        return color.upper()

    @staticmethod
    def _theme_icon_color(theme: PresentationTheme | None, icon_index: int) -> str | None:
        if theme is None:
            return None

        colors = [
            ChapterContentPage._normalize_icon_color(theme.colors.accent1),
            ChapterContentPage._normalize_icon_color(theme.colors.accent2),
        ]
        colors = [color for color in colors if color is not None]
        if not colors:
            return None
        return colors[icon_index % len(colors)]

    @staticmethod
    def _recolor_png_icon(
        icon_path: str,
        hex_color: str,
        *,
        output_dir: str | os.PathLike[str] | None = None,
    ) -> str:
        source_path = Path(icon_path)
        color = ChapterContentPage._normalize_icon_color(hex_color)
        if color is None:
            return icon_path

        if output_dir is None:
            output_root = Path(settings.TEMP_DIRECTORY or tempfile.gettempdir()) / "slidegen_recolored_icons"
        else:
            output_root = Path(output_dir)
        output_root.mkdir(parents=True, exist_ok=True)

        source_stat = source_path.stat()
        cache_key = f"{source_path.resolve()}:{source_stat.st_mtime_ns}:{source_stat.st_size}:{color}".encode()
        cache_hash = hashlib.sha1(cache_key).hexdigest()[:12]
        recolored_path = output_root / f"{source_path.stem}-{color}-{cache_hash}.png"
        if recolored_path.exists():
            return str(recolored_path)

        source_image = Image.open(source_path).convert("RGBA")
        rgb = tuple(int(color[index : index + 2], 16) for index in (0, 2, 4))
        recolored_image = Image.new("RGBA", source_image.size, (*rgb, 0))
        recolored_image.putalpha(source_image.getchannel("A"))
        recolored_image.save(recolored_path)
        return str(recolored_path)

    @staticmethod
    def _prepare_icon_for_theme(icon_path: str, theme: PresentationTheme | None, icon_index: int) -> str:
        color = ChapterContentPage._theme_icon_color(theme, icon_index)
        if color is None:
            return icon_path

        try:
            return ChapterContentPage._recolor_png_icon(icon_path, color)
        except Exception:
            logger.exception("{}: Failed to recolor icon '{}'", ChapterContentPage.__name__, icon_path)
            return icon_path

    @staticmethod
    def _get_slide_type(content: Heading) -> int:
        """
        Get the slide type of the chapter content page
        """
        return len(content)

    @staticmethod
    def _asset_key(
        shape_name: str,
        location_index: int,
        content_type: ComponentContentType,
    ) -> tuple[str, int, ComponentContentType]:
        return (shape_name, location_index, content_type)

    @staticmethod
    async def _resolve_picture_asset(
        shape_name: str,
        location_index: int,
        content_title: str,
        section_titles: list[str],
    ) -> ChapterSlideAsset:
        image_path = None
        try:
            prompt_text = section_titles[location_index] if location_index < len(section_titles) else content_title
            prompt = ImagePrompt(prompt=prompt_text, theme_prompt=None)

            logger.info(
                "{}: generating image asset for slide '{}' using prompt '{}...'",
                ChapterContentPage.__name__,
                content_title,
                prompt_text[:20],
            )
            image_result = await ChapterContentPage.image_generator.generate_image(prompt)
            if image_result.path and os.path.exists(image_result.path):
                image_path = image_result.path
        except Exception:
            logger.exception(f"{ChapterContentPage.__name__}: Image generation failed")

        if not image_path:
            image_path = ChapterContentPage._resolve_placeholder_image_path()
            if not image_path:
                raise PPTGenError(
                    f"{ChapterContentPage.__name__}: Unable to resolve image path for shape '{shape_name}'"
                )

        return ChapterSlideAsset(
            shape_name=shape_name,
            location_index=location_index,
            content_type=ComponentContentType.PICTURE,
            path=image_path,
        )

    @staticmethod
    async def _resolve_icon_asset(
        shape_name: str,
        location_index: int,
        content_title: str,
        section_titles: list[str],
        section_texts: list[str],
        theme: PresentationTheme | None,
        icon_index: int,
    ) -> ChapterSlideAsset:
        icon_path = None
        try:
            query = content_title
            if location_index < len(section_titles) and section_titles[location_index]:
                query = section_titles[location_index]
            elif location_index < len(section_texts) and section_texts[location_index]:
                query = section_texts[location_index]

            logger.info(
                "{}: searching icon for slide '{}' using query '{}'",
                ChapterContentPage.__name__,
                content_title,
                query,
            )
            results = await ChapterContentPage.icon_searcher.search_icons(query, k=1)
            if results:
                rel_path = results[0]
                abs_path = os.path.join(Path(__file__).resolve().parents[3], rel_path)
                icon_path = abs_path if os.path.exists(abs_path) else rel_path
        except Exception:
            logger.exception(f"{ChapterContentPage.__name__}: Icon search failed")

        if not icon_path:
            icon_path = ChapterContentPage._resolve_placeholder_image_path()
            if not icon_path:
                raise PPTGenError(
                    f"{ChapterContentPage.__name__}: Unable to resolve icon path for shape '{shape_name}'"
                )

        icon_path = ChapterContentPage._prepare_icon_for_theme(icon_path, theme, icon_index)
        return ChapterSlideAsset(
            shape_name=shape_name,
            location_index=location_index,
            content_type=ComponentContentType.ICON,
            path=icon_path,
        )

    @staticmethod
    async def _prepare_slide_data(
        content: Heading,
        *,
        style_override: Style | None = None,
        theme: PresentationTheme | None = None,
    ) -> ChapterSlideData:
        assert content.level in (2, 3), (
            f"{ChapterContentPage.__name__}: Chapter content page must have a level 2 or level 3 heading"
        )

        slide_type = ChapterContentPage._get_slide_type(content)
        if slide_type > 4:
            raise PPTGenError(f"{ChapterContentPage.__name__}: Invalid slide type: {slide_type}")

        section_titles = [child.element_text for child in content.children]
        section_texts = [child.text for child in content.children]
        chapter_layout = ChapterLayout(slide_type)
        if style_override is not None:
            style = style_override
        else:
            style = components_manager.get_random_style(chapter_layout)

        logger.debug(f"{ChapterContentPage.__name__}: {chapter_layout} {style.name if style else 'None'}")

        assets: dict[tuple[str, int, ComponentContentType], ChapterSlideAsset] = {}
        icon_index = 0
        sorted_shapes = sorted(style.shapes.items(), key=lambda x: x[1].zorder)
        for shape_name, shape in sorted_shapes:
            for idx, _loc in enumerate(shape.location):
                match shape.content_type:
                    case ComponentContentType.PICTURE:
                        asset = await ChapterContentPage._resolve_picture_asset(
                            shape_name,
                            idx,
                            content.element_text,
                            section_titles,
                        )
                        assets[ChapterContentPage._asset_key(shape_name, idx, ComponentContentType.PICTURE)] = asset
                    case ComponentContentType.ICON:
                        asset = await ChapterContentPage._resolve_icon_asset(
                            shape_name,
                            idx,
                            content.element_text,
                            section_titles,
                            section_texts,
                            theme,
                            icon_index,
                        )
                        assets[ChapterContentPage._asset_key(shape_name, idx, ComponentContentType.ICON)] = asset
                        icon_index += 1

        return ChapterSlideData(
            content_title=content.element_text,
            section_titles=section_titles,
            section_texts=section_texts,
            style=style,
            assets=assets,
        )

    @staticmethod
    def _get_prepared_asset(
        slide_data: ChapterSlideData,
        shape_name: str,
        location_index: int,
        content_type: ComponentContentType,
    ) -> ChapterSlideAsset:
        key = ChapterContentPage._asset_key(shape_name, location_index, content_type)
        try:
            return slide_data.assets[key]
        except KeyError as exc:
            raise PPTGenError(
                f"{ChapterContentPage.__name__}: Missing prepared asset for shape '{shape_name}' "
                f"at location {location_index}"
            ) from exc

    @staticmethod
    def _render_slide(prs: Presentation, new_slide: Slide, slide_data: ChapterSlideData) -> None:
        title_shape = Page._find_or_inject_placeholder(
            slide=new_slide,
            page_type="chapter_content",
            role="title",
            text=slide_data.content_title,
            placeholder_types=(PP_PLACEHOLDER.TITLE,),
            shape_name_prefix="content_title",
        )
        if not title_shape:
            logger.warning(
                f"{ChapterContentPage.__name__}: No title placeholder found and no fallback data "
                f"in shape.json page_placeholders.chapter_content; title will not be displayed"
            )

        index = 0
        sorted_shapes = sorted(slide_data.style.shapes.items(), key=lambda x: x[1].zorder)
        for shape_name, shape in sorted_shapes:
            locs = shape.location
            for idx, loc in enumerate(locs):
                added_shape: BaseShape
                scaled_loc = Page._scale_shape_location(prs, loc)
                match shape.content_type:
                    case ComponentContentType.CONTENT:
                        if len(slide_data.section_texts) != len(locs):
                            raise PPTGenError(
                                f"{ChapterContentPage.__name__}: "
                                f"Text content must be equal to the number of locations: "
                                f"{len(slide_data.section_texts)} != {len(locs)}"
                            )
                        assert shape.xml is not None
                        added_shape = ShapeFactory.create_text_shape(
                            new_slide,
                            shape.xml,
                            slide_data.section_texts[idx],
                            scaled_loc,
                            shape_id=index,
                            shape_name=shape_name,
                        )
                        Page._shape_alignment(added_shape)
                    case ComponentContentType.TITLE:
                        if len(slide_data.section_titles) != len(locs):
                            raise PPTGenError(
                                f"{ChapterContentPage.__name__}: "
                                f"Title must be equal to the number of locations: "
                                f"{len(slide_data.section_titles)} != {len(locs)}"
                            )
                        assert shape.xml is not None
                        added_shape = ShapeFactory.create_text_shape(
                            new_slide,
                            shape.xml,
                            slide_data.section_titles[idx],
                            scaled_loc,
                            shape_id=index,
                            shape_name=shape_name,
                        )
                        Page._shape_alignment(added_shape)
                    case ComponentContentType.PICTURE:
                        asset = ChapterContentPage._get_prepared_asset(
                            slide_data, shape_name, idx, ComponentContentType.PICTURE
                        )
                        added_shape = ShapeFactory.create_image_shape(new_slide, asset.path, scaled_loc)
                    case ComponentContentType.NUMBER:
                        assert shape.xml is not None
                        added_shape = ShapeFactory.create_text_shape(
                            new_slide,
                            shape.xml,
                            str(idx + 1).zfill(2),
                            scaled_loc,
                            shape_id=index,
                            shape_name=shape_name,
                        )
                    case ComponentContentType.ICON:
                        asset = ChapterContentPage._get_prepared_asset(
                            slide_data, shape_name, idx, ComponentContentType.ICON
                        )
                        added_shape = ShapeFactory.create_image_shape(new_slide, asset.path, scaled_loc)
                    case _:
                        added_shape = ShapeFactory.create_xml_shape(
                            new_slide,
                            shape.xml,  # type: ignore
                            scaled_loc,
                            shape_id=index,
                            shape_name=shape_name,
                        )
            index += 1

        if title_shape is not None:
            Page.bring_shape_to_front(title_shape)

    @staticmethod
    async def generate_slide(
        prs: Presentation,
        content: Heading,
        *,
        chapter_page_index: int = 3,
        slide_index: int = 3,
        style_override: Style | None = None,
        theme: PresentationTheme | None = None,
    ) -> None:
        """
        Generate the chapter content page

        Args:
            prs: Presentation object
            content: Heading object
            chapter_page_index: index of the template chapter content slide
            slide_index: index of the slide to be generated
            style_override: if set, use this Style instead of calling get_random_style()
            theme: if set, recolor inserted icons with accent1/accent2
        """
        slide_data = await ChapterContentPage._prepare_slide_data(
            content,
            style_override=style_override,
            theme=theme,
        )
        chapter_page = prs.slides[chapter_page_index]
        new_slide = prs.slides.add_slide(chapter_page.slide_layout)
        ChapterContentPage._render_slide(prs, new_slide, slide_data)
        ChapterContentPage.move_slide(prs, new_slide, slide_index)


class EndPage(Page):
    """End page"""

    @staticmethod
    async def generate_slide(
        prs: Presentation,
        content: Heading | None = None,
        *,
        end_page_index: int = 4,
        slide_index: int = 4,
    ) -> None:
        if content is None:
            content = Heading(text="Thank you!", level=2)
        template_slide = prs.slides[end_page_index]
        end_slide = prs.slides.add_slide(template_slide.slide_layout)
        title_shape = Page._find_or_inject_placeholder(
            slide=end_slide,
            page_type="end",
            role="title",
            text=content.element_text,
            placeholder_types=(PP_PLACEHOLDER.TITLE,),
            shape_name_prefix="end_title",
        )
        if not title_shape:
            raise PPTTemplateError(
                f"{EndPage.__name__}: No title placeholder found in end slide "
                f"and no fallback data in shape.json page_placeholders.end, "
                f"end slide index: {end_page_index}"
            )
        EndPage.move_slide(prs, end_slide, slide_index)
