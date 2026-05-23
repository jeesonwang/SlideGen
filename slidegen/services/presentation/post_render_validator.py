from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Literal

from pptx.util import Pt


@dataclass
class SlideGeometryIssue:
    level: Literal["error", "warning"]
    slide_index: int
    message: str
    shapes_involved: tuple[str, ...] = ()


class PostRenderValidator:
    def __init__(self, mode: Literal["off", "warn", "fail"] = "warn"):
        self.mode = mode

    def validate(self, prs: Any) -> list[SlideGeometryIssue]:
        if self.mode == "off":
            return []
        issues: list[SlideGeometryIssue] = []
        for i, slide in enumerate(prs.slides):
            issues.extend(self._check_out_of_bounds(slide, i, int(prs.slide_width), int(prs.slide_height)))
            issues.extend(self._check_readability(slide, i))
            issues.extend(self._check_overlaps(slide, i))
        return issues

    def _check_out_of_bounds(self, slide: Any, index: int, slide_w: int, slide_h: int) -> list[SlideGeometryIssue]:
        issues = []
        for shape in slide.shapes:
            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height
            if left < 0 or top < 0 or right > slide_w or bottom > slide_h:
                issues.append(
                    SlideGeometryIssue(
                        level="error" if self.mode == "fail" else "warning",
                        slide_index=index,
                        message=f"Shape '{shape.name}' extends beyond slide boundaries",
                        shapes_involved=(shape.name,),
                    )
                )
        return issues

    def _check_overlaps(self, slide: Any, index: int) -> list[SlideGeometryIssue]:
        issues = []
        text_shapes = [
            shape for shape in slide.shapes
            if shape.has_text_frame and getattr(shape, "text", "").strip()
        ]
        for i, first in enumerate(text_shapes):
            for second in text_shapes[i + 1:]:
                if self._overlap_area(first, second) <= 0:
                    continue
                issues.append(
                    SlideGeometryIssue(
                        level="error" if self.mode == "fail" else "warning",
                        slide_index=index,
                        message=f"Text shapes '{first.name}' and '{second.name}' overlap",
                        shapes_involved=(first.name, second.name),
                    )
                )
        return issues

    def _overlap_area(self, first: Any, second: Any) -> int:
        left = max(first.left, second.left)
        top = max(first.top, second.top)
        right = min(first.left + first.width, second.left + second.width)
        bottom = min(first.top + first.height, second.top + second.height)
        if right <= left or bottom <= top:
            return 0
        return int((right - left) * (bottom - top))

    def _check_readability(self, slide: Any, index: int) -> list[SlideGeometryIssue]:
        issues = []
        min_size = Pt(8)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.font.size and para.font.size < min_size:
                    issues.append(
                        SlideGeometryIssue(
                            level="warning",
                            slide_index=index,
                            message=f"Text in shape '{shape.name}' is smaller than 8pt",
                            shapes_involved=(shape.name,),
                        )
                    )
                for run in para.runs:
                    if run.font.size and run.font.size < min_size:
                        issues.append(
                            SlideGeometryIssue(
                                level="warning",
                                slide_index=index,
                                message=f"Text in shape '{shape.name}' is smaller than 8pt",
                                shapes_involved=(shape.name,),
                            )
                        )
        return issues
