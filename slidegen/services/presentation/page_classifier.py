import json
import uuid
from collections.abc import Callable
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Any, Literal

from agno.agent import Agent
from agno.models.base import Model
from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide
from pydantic import BaseModel, Field

from slidegen.schemas.gen_request import BaseGenerationRequest
from slidegen.services.slidegen.workflow import get_llm_instance


class PageType(str, Enum):
    COVER = "cover"
    CATALOG = "catalog"
    CHAPTER_HOME = "chapter_home"
    CHAPTER_CONTENT = "chapter_content"
    END = "end"
    UNKNOWN = "unknown"


@dataclass
class PageClassification:
    page_index: int
    page_type: PageType
    confidence: float
    reason: str
    method: Literal["rule", "llm", "fallback"]


@dataclass
class ShapeSummary:
    shape_id: int
    name: str
    shape_type: str
    text: str
    x: int
    y: int
    width: int
    height: int
    has_text_frame: bool
    is_placeholder: bool
    is_picture: bool
    font_size: float | None
    is_bold: bool | None


class AgentPageClassificationOutput(BaseModel):
    """Structured output returned by the agno page classifier agent."""

    page_type: PageType = Field(description="One of cover, catalog, chapter_home, chapter_content, end, unknown")
    confidence: float = Field(ge=0.0, le=1.0, description="Classification confidence from 0 to 1")
    reason: str = Field(min_length=1, description="Short reason for the classification")


AgentFactory = Callable[[Model | Any], Any]


PAGE_CLASSIFIER_INSTRUCTIONS = """
Classify a PowerPoint slide into exactly one page type:
- cover: presentation cover with a main title, subtitle, author, date, or topic.
- catalog: agenda, table of contents, roadmap, or chapter list.
- chapter_home: section divider or chapter opening page that highlights one chapter title or number.
- chapter_content: actual content slide with ideas, title/body groups, analysis, diagrams, data, or explanations.
- end: closing slide, thank-you slide, Q&A, contact, or final call-to-action.
- unknown: use when the slide cannot be classified reliably.

Return only the structured output requested by the output schema. Do not classify one_point/two_points layouts.
""".strip()

END_KEYWORDS = ("thank you", "thanks", "q&a", "qa", "questions", "谢谢", "感谢")


class PageTypeClassifier:
    """Classify PPT slides into the page types used by pages.py."""

    def __init__(
        self,
        *,
        max_text_chars: int = 300,
        min_llm_confidence: float = 0.6,
        agent_factory: AgentFactory | None = None,
    ) -> None:
        self.max_text_chars = max_text_chars
        self.min_llm_confidence = min_llm_confidence
        self.agent_factory = agent_factory or self._default_agent_factory

    def summarize_slide(self, slide: Slide) -> list[ShapeSummary]:
        summaries: list[ShapeSummary] = []
        for shape in slide.shapes:
            try:
                text = self._shape_text(shape)
                font_size, is_bold = self._shape_font(shape)
                summaries.append(
                    ShapeSummary(
                        shape_id=shape.shape_id,
                        name=shape.name,
                        shape_type=str(shape.shape_type),
                        text=text[: self.max_text_chars],
                        x=int(shape.left),
                        y=int(shape.top),
                        width=int(shape.width),
                        height=int(shape.height),
                        has_text_frame=bool(getattr(shape, "has_text_frame", False)),
                        is_placeholder=bool(getattr(shape, "is_placeholder", False)),
                        is_picture=shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
                        font_size=font_size,
                        is_bold=is_bold,
                    )
                )
            except Exception as exc:
                logger.warning("Failed to summarize shape on slide: {}", exc)
        return summaries

    def _shape_text(self, shape: Any) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return str(getattr(shape, "text", "") or "").strip()

    def _shape_font(self, shape: Any) -> tuple[float | None, bool | None]:
        if not getattr(shape, "has_text_frame", False):
            return None, None
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.font.size is not None or paragraph.font.bold is not None:
                return (
                    paragraph.font.size.pt if paragraph.font.size is not None else None,
                    paragraph.font.bold,
                )
            for run in paragraph.runs:
                if run.font.size is not None or run.font.bold is not None:
                    return (
                        run.font.size.pt if run.font.size is not None else None,
                        run.font.bold,
                    )
        return None, None

    def classify_by_rules(
        self,
        *,
        page_index: int,
        slide_count: int,
        summaries: list[ShapeSummary],
    ) -> PageClassification | None:
        text_shapes = [summary for summary in summaries if summary.text]
        combined_text = " ".join(summary.text.casefold() for summary in text_shapes)

        if not text_shapes and not any(summary.is_picture for summary in summaries):
            return PageClassification(
                page_index=page_index,
                page_type=PageType.UNKNOWN,
                confidence=1.0,
                reason="Slide has no text or picture content.",
                method="rule",
            )

        if page_index == slide_count - 1 and any(keyword in combined_text for keyword in END_KEYWORDS):
            return PageClassification(
                page_index=page_index,
                page_type=PageType.END,
                confidence=0.95,
                reason="Last slide contains a clear closing phrase.",
                method="rule",
            )

        if page_index == 0 and 1 <= len(text_shapes) <= 3:
            total_text = " ".join(summary.text for summary in text_shapes)
            if len(total_text) <= 180:
                return PageClassification(
                    page_index=page_index,
                    page_type=PageType.COVER,
                    confidence=0.9,
                    reason="First slide has only a small number of title-like text blocks.",
                    method="rule",
                )

        return None

    async def classify_with_agent(
        self,
        *,
        model: Model | Any,
        page_index: int,
        slide_count: int,
        summaries: list[ShapeSummary],
    ) -> PageClassification:
        prompt = self._build_prompt(page_index=page_index, slide_count=slide_count, summaries=summaries)
        agent = self.agent_factory(model)

        try:
            response = await agent.arun(prompt)
            output = self._coerce_agent_output(getattr(response, "content", response))
            if output.confidence < self.min_llm_confidence:
                return self._fallback(
                    page_index=page_index,
                    reason=(
                        f"LLM confidence {output.confidence:.2f} below threshold "
                        f"{self.min_llm_confidence:.2f}: {output.reason}"
                    ),
                )
            return PageClassification(
                page_index=page_index,
                page_type=output.page_type,
                confidence=output.confidence,
                reason=output.reason,
                method="llm",
            )
        except Exception as exc:
            logger.warning("Page classifier Agent failed for slide {}: {}", page_index, exc)
            return self._fallback(page_index=page_index, reason=str(exc))

    def _build_prompt(
        self,
        *,
        page_index: int,
        slide_count: int,
        summaries: list[ShapeSummary],
    ) -> str:
        payload = {
            "page_index": page_index,
            "slide_count": slide_count,
            "shape_count": len(summaries),
            "text_shape_count": sum(1 for summary in summaries if summary.text),
            "picture_count": sum(1 for summary in summaries if summary.is_picture),
            "shapes": [
                {
                    "shape_id": summary.shape_id,
                    "name": summary.name,
                    "shape_type": summary.shape_type,
                    "text": summary.text,
                    "x": summary.x,
                    "y": summary.y,
                    "width": summary.width,
                    "height": summary.height,
                    "has_text_frame": summary.has_text_frame,
                    "is_placeholder": summary.is_placeholder,
                    "is_picture": summary.is_picture,
                    "font_size": summary.font_size,
                    "is_bold": summary.is_bold,
                }
                for summary in summaries
            ],
        }
        return (
            "Classify this PowerPoint slide using the provided structured slide summary. "
            "Choose one page_type from the schema.\n\n"
            f"{json.dumps(payload, ensure_ascii=False)}"
        )

    def _coerce_agent_output(self, content: Any) -> AgentPageClassificationOutput:
        if isinstance(content, AgentPageClassificationOutput):
            return content
        if isinstance(content, dict):
            return AgentPageClassificationOutput.model_validate(content)
        if isinstance(content, str):
            return AgentPageClassificationOutput.model_validate_json(content)
        return AgentPageClassificationOutput.model_validate(content)

    def _fallback(self, *, page_index: int, reason: str) -> PageClassification:
        return PageClassification(
            page_index=page_index,
            page_type=PageType.UNKNOWN,
            confidence=0.0,
            reason=reason or "Unable to classify slide.",
            method="fallback",
        )

    def _default_agent_factory(self, model: Model | Any) -> Agent:
        return Agent(
            name="Page type classifier",
            instructions=[PAGE_CLASSIFIER_INSTRUCTIONS],
            model=model,
            structured_outputs=True,
            output_model=AgentPageClassificationOutput,
        )

    async def classify_pages(
        self,
        *,
        pptx_path: str | Path,
        user_id: uuid.UUID,
        llm_config_id: uuid.UUID | None = None,
        model: Model | Any | None = None,
    ) -> list[PageClassification]:
        path = Path(pptx_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        if model is None:
            request = BaseGenerationRequest(
                content="PPT page classification",
                user_id=user_id,
                llm_config_id=llm_config_id,
            )
            model = await get_llm_instance(request)

        presentation = Presentation(str(path))
        slide_count = len(presentation.slides)
        results: list[PageClassification] = []

        for page_index, slide in enumerate(presentation.slides):
            try:
                summaries = self.summarize_slide(slide)
                rule_result = self.classify_by_rules(
                    page_index=page_index,
                    slide_count=slide_count,
                    summaries=summaries,
                )
                if rule_result is not None:
                    results.append(rule_result)
                    continue

                results.append(
                    await self.classify_with_agent(
                        model=model,
                        page_index=page_index,
                        slide_count=slide_count,
                        summaries=summaries,
                    )
                )
            except Exception as exc:
                logger.warning("Failed to classify slide {}: {}", page_index, exc)
                results.append(self._fallback(page_index=page_index, reason=str(exc)))

        return results
