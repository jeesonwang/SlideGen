from __future__ import annotations

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.services.document.markdown import Heading


class NativePage:
    @staticmethod
    def _blank_slide(prs: Presentation, slide_index: int):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        old_index = len(slides) - 1
        xml_slides.remove(slides[old_index])
        xml_slides.insert(slide_index, slides[old_index])
        return slide

    @staticmethod
    def _add_text(slide, text: str, left: float, top: float, width: float, height: float, font_size: int) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(font_size)


class NativeCoverPage(NativePage):
    @staticmethod
    async def generate_slide(prs: Presentation, content: Heading, *, slide_index: int) -> None:
        slide = NativeCoverPage._blank_slide(prs, slide_index)
        NativeCoverPage._add_text(slide, content.element_text or "Presentation Title", 0.8, 2.4, 8.6, 1.0, 36)


class NativeCatalogPage(NativePage):
    @staticmethod
    async def generate_slide(prs: Presentation, chapters: list[Heading], *, slide_index: int) -> int:
        slide = NativeCatalogPage._blank_slide(prs, slide_index)
        NativeCatalogPage._add_text(slide, "Agenda", 0.8, 0.5, 8.6, 0.6, 28)
        for index, chapter in enumerate(chapters, start=1):
            NativeCatalogPage._add_text(slide, f"{index:02d}. {chapter.element_text}", 1.0, 1.1 + index * 0.45, 8.0, 0.35, 18)
        return slide_index


class NativeChapterHomePage(NativePage):
    @staticmethod
    async def generate_slide(prs: Presentation, content: Heading, *, chapter_number: int, slide_index: int) -> None:
        slide = NativeChapterHomePage._blank_slide(prs, slide_index)
        NativeChapterHomePage._add_text(slide, f"PART {chapter_number:02d}", 0.8, 1.8, 8.0, 0.5, 20)
        NativeChapterHomePage._add_text(slide, content.element_text, 0.8, 2.4, 8.0, 0.8, 34)


class NativeChapterContentPage(NativePage):
    @staticmethod
    async def generate_slide(prs: Presentation, content: Heading, *, slide_index: int) -> None:
        slide = NativeChapterContentPage._blank_slide(prs, slide_index)
        NativeChapterContentPage._add_text(slide, content.element_text, 0.7, 0.4, 8.8, 0.5, 26)
        lines = []
        for child in content.children:
            text = child.element_text.strip()
            if text:
                lines.append(text)
        body = "\n".join(lines) if lines else content.text.strip()
        NativeChapterContentPage._add_text(slide, body, 0.9, 1.3, 8.2, 4.6, 18)


class NativeEndPage(NativePage):
    @staticmethod
    async def generate_slide(prs: Presentation, *, slide_index: int) -> None:
        slide = NativeEndPage._blank_slide(prs, slide_index)
        NativeEndPage._add_text(slide, "Thank you!", 0.8, 2.5, 8.6, 0.8, 36)
