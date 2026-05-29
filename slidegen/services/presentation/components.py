import json
from dataclasses import dataclass
from enum import Enum, IntEnum
from pathlib import Path
from typing import Any, cast

from loguru import logger
from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide

from slidegen.core import settings


def remove_custDataLst(xml_str: str) -> str:
    """
    Remove the <p:custDataLst> part in the XML and return the processed XML string.

    Args:
        xml_str (str): The input XML string, containing the <p:custDataLst> part.

    Returns:
        str: The processed XML string, without the <p:custDataLst> part.
    """
    root = etree.fromstring(xml_str)
    ns = root.nsmap

    cust_data_list = root.find(".//p:custDataLst", namespaces=ns)

    if cust_data_list is not None:
        parent = cust_data_list.getparent()
        if parent is not None:
            parent.remove(cust_data_list)

    return etree.tostring(
        root,
        encoding="unicode",
        pretty_print=True,
        xml_declaration=False,
    )


class ComponentContentType(Enum):
    CONTENT = "content"
    PICTURE = "picture"
    NUMBER = "number"
    TITLE = "title"
    ICON = "icon"
    DECORATION = "decoration"

    def __eq__(self, other: str | object) -> bool:
        if isinstance(other, str):
            return self.value == other
        return super().__eq__(other)


class ChapterLayout(IntEnum):
    """Chapter layout type"""

    str_value: str

    ONE_POINT = (1, "one_point")
    TWO_POINTS = (2, "two_points")
    THREE_POINTS = (3, "three_points")
    FOUR_POINTS = (4, "four_points")

    def __new__(cls, number: int, str_value: str = "") -> "ChapterLayout":
        obj = int.__new__(cls, number)
        obj._value_ = number
        obj.str_value = str_value
        return obj

    def __str__(self) -> str:
        return f"{self.name} ({self.value})"


@dataclass
class Location:
    """Information about the location of the shape"""

    x: int
    y: int
    width: int
    height: int


@dataclass
class CShape:
    """Basic shape element in PPT"""

    xml: str | None
    zorder: int
    content_type: ComponentContentType
    location: list[Location]

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> "CShape":
        """Create a Shape object from dictionary data"""
        location_list = []
        for loc in data.get("location", []):
            location_list.append(
                Location(x=loc.get("x", 0), y=loc.get("y", 0), width=loc.get("width", 0), height=loc.get("height", 0))
            )

        ct = data.get("content_type")
        if ct is None:
            content_type = ComponentContentType.DECORATION
        elif isinstance(ct, ComponentContentType):
            content_type = ct
        else:
            content_type = ComponentContentType(ct)

        return cls(
            xml=data.get("xml"),
            zorder=data.get("zorder", 0),
            content_type=content_type,
            location=location_list,
        )

    def to_dict(self) -> dict[str, Any]:
        location_list = []
        for loc in self.location:
            location_list.append({"x": loc.x, "y": loc.y, "width": loc.width, "height": loc.height})

        return {
            "xml": self.xml,
            "zorder": self.zorder,
            "content_type": self.content_type.value,
            "location": location_list,
        }


@dataclass
class PagePlaceholder:
    """Stored placeholder info extracted from a template page slide."""

    xml: str
    placeholder_type: str
    location: Location

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> "PagePlaceholder":
        loc = data.get("location", {})
        location = Location(
            x=loc.get("x", 0),
            y=loc.get("y", 0),
            width=loc.get("width", 0),
            height=loc.get("height", 0),
        )
        return cls(
            xml=data["xml"],
            placeholder_type=data["placeholder_type"],
            location=location,
        )

    def to_dict(self) -> dict[str, Any]:
        return {
            "xml": self.xml,
            "placeholder_type": self.placeholder_type,
            "location": {
                "x": self.location.x,
                "y": self.location.y,
                "width": self.location.width,
                "height": self.location.height,
            },
        }


class Style:
    """Represent a style of shapes, including multiple shapes"""

    def __init__(self, name: str, shapes_data: dict[str, Any] | None = None):
        self.name = name
        self.shapes: dict[str, CShape] = {}

        if shapes_data:
            self.load_from_dict(shapes_data)

    def load_from_dict(self, shapes_data: dict[str, Any]) -> None:
        for shape_name, shape_data in shapes_data.items():
            self.shapes[shape_name] = CShape.from_dict(shape_data)

    def to_dict(self) -> dict[str, Any]:
        result = {}
        for shape_name, shape in self.shapes.items():
            result[shape_name] = shape.to_dict()
        return result

    @property
    def shape_list(self) -> list[CShape]:
        """Get a list of shapes"""
        return list(self.shapes.values())

    @property
    def shape_names(self) -> list[str]:
        """Get a list of shape names"""
        return list(self.shapes.keys())

    def get_shape(self, shape_name: str) -> CShape | None:
        """Get a shape by its name"""
        return self.shapes.get(shape_name)

    def add_shape(self, shape_name: str, shape: CShape) -> None:
        """Add a shape to the style"""
        self.shapes[shape_name] = shape

    def __len__(self) -> int:
        """Get the number of shapes"""
        return len(self.shapes)


class LayoutType:
    """Represent a layout type (like two_points or three_points)."""

    def __init__(self, name: str, data: dict[str, Any] | None = None):
        self.name = name
        self.styles: dict[str, Style] = {}
        self.style_order: Style | None = None

        if data:
            self.load_from_dict(data)

    def load_from_dict(self, data: dict[str, Any]) -> None:
        """Load layout type data from a dictionary"""
        for style_name, style_data in data.items():
            if style_name == "style_ordered":
                self.style_order = Style(style_name, style_data)
            self.styles[style_name] = Style(style_name, style_data)

    def to_dict(self) -> dict[str, Any]:
        result = {}
        for style_name, style in self.styles.items():
            result[style_name] = style.to_dict()

        return result

    def add_style(self, style: Style) -> None:
        """Add a style to the layout type"""
        self.styles[style.name] = style

    def get_style(self, style_name: str) -> Style | None:
        return self.styles.get(style_name)

    def __len__(self) -> int:
        return len(self.styles)

    @property
    def style_names(self) -> list[str]:
        return list(self.styles.keys())

    @property
    def style_list(self) -> list[Style]:
        return list(self.styles.values())


class ComponentsManager:
    """Manage all components, styles, and layouts"""

    def __init__(self, json_path: str | Path | None = None):
        self.metadata: dict[str, Any] = {}
        self.layout_types: dict[str, LayoutType] = {}
        self.page_placeholders: dict[str, dict[str, PagePlaceholder]] = {}

        if json_path:
            self.load_from_json(json_path)
            self.json_path = json_path

    def _load_from_dict(self, data: dict[str, Any]) -> None:
        self.metadata = {}
        self.layout_types = {}
        self.page_placeholders = {}
        for key, value in data.items():
            if key == "metadata":
                self.metadata = dict(value) if isinstance(value, dict) else {}
            elif key == "page_placeholders":
                for page_type, ph_data in value.items():
                    self.page_placeholders[page_type] = {}
                    for role, role_data in ph_data.items():
                        self.page_placeholders[page_type][role] = PagePlaceholder.from_dict(role_data)
            else:
                self.layout_types[key] = LayoutType(key, value)

    def load_from_json(self, json_path: str | Path) -> None:
        with open(json_path, encoding="utf-8") as f:
            data = json.load(f)
        self._load_from_dict(data)

    def save_to_json(self, json_path: str | Path, backup: bool = True) -> None:
        data: dict[str, Any] = {}
        if self.metadata:
            data["metadata"] = self.metadata
        for layout_name, layout in self.layout_types.items():
            data[layout_name] = layout.to_dict()

        if self.page_placeholders:
            data["page_placeholders"] = {}
            for page_type, roles in self.page_placeholders.items():
                data["page_placeholders"][page_type] = {}
                for role, ph in roles.items():
                    data["page_placeholders"][page_type][role] = ph.to_dict()

        from slidegen.services.presentation.component_importer import atomic_save_json

        atomic_save_json(data, json_path, backup=backup)

    def get_layout_type(self, layout_name: ChapterLayout | str) -> LayoutType | None:
        """Get a layout type by its name

        example:
        ```
        shapes_manager = ShapesManager(json_path)
        layout = shapes_manager.get_layout_type(ChapterLayout.TWO_POINTS)
        ```
        """
        layout_key = layout_name.str_value if isinstance(layout_name, ChapterLayout) else layout_name
        return self.layout_types.get(layout_key)

    def get_random_style(self, layout_name: ChapterLayout) -> Style:
        """Get a random style from a specified layout type

        Args:
            layout_name (str): The name of the layout type to get a random style from.

        Returns:
            Optional[Style]: A random style from the specified layout type, or None if the layout type does not exist or has no styles.

        example:
        ```
        shapes_manager = ShapesManager(json_path)
        style = shapes_manager.get_random_style(ChapterLayout.TWO_POINTS)
        ```
        """
        import random

        layout = self.get_layout_type(layout_name)
        if not layout or not layout.styles:
            raise ValueError(
                f"Layout type '{layout_name}' not found, available layout types: {self.layout_types_names}"
            )

        style_name = random.choice(list(layout.styles.keys()))
        return layout.styles[style_name]

    @staticmethod
    def are_same_shape(xml1: str, xml2: str) -> bool:
        """
        Compare two PPT shapes to determine if they are of the same type.

        Args:
            xml1 (str): The XML representation of the first shape
            xml2 (str): The XML representation of the second shape

        Returns:
            bool: True if the shapes are of the same type, False otherwise
        """
        try:
            root1 = etree.fromstring(xml1)
            root2 = etree.fromstring(xml2)
            nsmap = root1.nsmap

            xfrm_element1 = root1.find(".//a:xfrm", namespaces=nsmap)
            xfrm_element2 = root2.find(".//a:xfrm", namespaces=nsmap)

            if xfrm_element1 is not None and xfrm_element1.getparent() is not None:
                xfrm_element1.getparent().remove(xfrm_element1)

            if xfrm_element2 is not None and xfrm_element2.getparent() is not None:
                xfrm_element2.getparent().remove(xfrm_element2)

            cnvpr1 = root1.find(".//p:cNvPr", namespaces=nsmap)
            cnvpr2 = root2.find(".//p:cNvPr", namespaces=nsmap)

            if cnvpr1 is not None and cnvpr2 is not None:
                cnvpr1.set("id", "1")
                cnvpr1.set("name", "temp")
                cnvpr2.set("id", "1")
                cnvpr2.set("name", "temp")

            t_elements1 = root1.findall(".//a:t", namespaces=nsmap)
            t_elements2 = root2.findall(".//a:t", namespaces=nsmap)

            for t_elem in t_elements1:
                t_elem.text = "placeholder_text"

            for t_elem in t_elements2:
                t_elem.text = "placeholder_text"

            xml_str1 = etree.tostring(root1, encoding="unicode")
            xml_str2 = etree.tostring(root2, encoding="unicode")

            return xml_str1 == xml_str2

        except Exception as e:
            logger.exception(f"Compare shapes error: {e}")
            return False

    def add_style_from_slide(self, slide: Slide, layout_type: ChapterLayout, style_name: str) -> None:
        """Add a style to the manager from a slide

        Args:
            slide (Slide): The slide to add the style from.
            layout_type (str): The layout type to add the style to, like "ChapterLayout.TWO_POINTS".
            style_name (str): The name of the style to add, like "style_1".
        """
        layout = self.get_layout_type(layout_type)
        if not layout:
            raise ValueError(
                f"Layout type '{layout_type}' not found, available layout types: {self.layout_types_names}"
            )

        if style_name in layout.styles:
            raise ValueError(
                f"{self.__class__.__name__}: Style '{style_name}' already exists in layout type '{layout_type}'"
            )

        new_style = Style(style_name)

        shapes = slide.shapes
        shape_data_dict: dict[str, dict[str, Any]] = {}
        area = 0
        for i, shape in enumerate(shapes):
            if shape.is_placeholder:
                continue
            location = Location(x=shape.left, y=shape.top, width=shape.width, height=shape.height)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if self.is_icon(location):
                    content_type = ComponentContentType.ICON
                else:
                    content_type = ComponentContentType.PICTURE
                shape_data: dict[str, Any] = {
                    "xml": None,
                    "zorder": i,
                    "content_type": content_type,
                    "location": [
                        {"x": location.x, "y": location.y, "width": location.width, "height": location.height}
                    ],
                }
                shape_name = f"{shape.name}_{i}"
                shape_data_dict[shape_name] = shape_data
            elif shape.has_text_frame:
                xml_str = shape._element.xml
                xml_str = remove_custDataLst(xml_str)
                current_area = shape.height * shape.width
                if current_area > area:
                    area = current_area
                shape_text = cast(Shape, shape).text
                content_type = ComponentContentType.DECORATION
                if shape_text:
                    content_type = ComponentContentType.CONTENT
                shape_data = {
                    "xml": xml_str,
                    "zorder": i,
                    "content_type": content_type,
                    "location": [
                        {"x": location.x, "y": location.y, "width": location.width, "height": location.height}
                    ],
                }
                shape_name = f"{shape.name}_{i}"
                shape_data_dict[shape_name] = shape_data
            else:
                try:
                    xml_str = shape._element.xml
                    xml_str = remove_custDataLst(xml_str)

                    shape_data = {
                        "xml": xml_str,
                        "zorder": i,
                        "content_type": ComponentContentType.DECORATION,
                        "location": [
                            {"x": location.x, "y": location.y, "width": location.width, "height": location.height}
                        ],
                    }
                    shape_name = f"{shape.name}_{i}"
                    shape_data_dict[shape_name] = shape_data
                except Exception as e:
                    logger.error(f"Error extracting XML from shape {i}: {e}")
                    continue

        # Merge similar shapes (use a copy to avoid modifying the dictionary during iteration)
        shapes_to_remove = []
        for shape_name, shape_data in shape_data_dict.items():
            if shape_data["xml"] is None:
                continue

            if shape_data["content_type"] == ComponentContentType.CONTENT and isinstance(shape_data["location"], list):
                current_area = int(shape_data["location"][0]["height"]) * int(shape_data["location"][0]["width"])
                if current_area < (area - 10000):
                    if self.get_text_from_xml(shape_data["xml"]).isdigit():
                        shape_data["content_type"] = ComponentContentType.NUMBER
                    else:
                        shape_data["content_type"] = ComponentContentType.TITLE

            for shape_name_other, shape_data_other in list(shape_data_dict.items()):
                if shape_name == shape_name_other or shape_data_other["xml"] is None:
                    continue
                if self.are_same_shape(shape_data["xml"], shape_data_other["xml"]):
                    shape_data["location"].extend(shape_data_other["location"])
                    shape_data_other["xml"] = None
                    shapes_to_remove.append(shape_name_other)

        for shape_name in shapes_to_remove:
            if shape_name in shape_data_dict:
                shape_data_dict.pop(shape_name)

        for shape_name, shape_data in shape_data_dict.items():
            new_style.shapes[shape_name] = CShape.from_dict(shape_data)

        layout.add_style(new_style)
        logger.info(f"Added style '{style_name}' to layout type '{layout_type}' with {len(new_style.shapes)} shapes")

    def extract_page_placeholder(self, slide: Slide, page_type: str) -> None:
        """Extract placeholder shapes from a template slide for a given page type.

        Args:
            slide: The template slide to scan.
            page_type: Key name like "cover", "chapter_home", "chapter_content", "end".
        """
        if page_type not in self.page_placeholders:
            self.page_placeholders[page_type] = {}

        # Map placeholder types to standardized role keys for lookup
        _TYPE_TO_ROLE = {"TITLE": "title", "CENTER_TITLE": "title", "SUBTITLE": "subtitle"}
        for ph in slide.shapes.placeholders:
            type_name = ph.placeholder_format.type.name
            role_key = _TYPE_TO_ROLE.get(type_name)
            if role_key is None:
                continue

            xml_str = ph._element.xml
            xml_str = remove_custDataLst(xml_str)

            location = Location(x=ph.left, y=ph.top, width=ph.width, height=ph.height)
            self.page_placeholders[page_type][role_key] = PagePlaceholder(
                xml=xml_str,
                placeholder_type=type_name,
                location=location,
            )
        logger.info(f"Extracted {len(self.page_placeholders[page_type])} placeholder(s) for page type '{page_type}'")

    def get_page_placeholder(self, page_type: str, role: str) -> PagePlaceholder | None:
        """Retrieve stored placeholder data for a page type and role.

        Args:
            page_type: e.g. "cover", "chapter_home", "chapter_content", "end".
            role: e.g. "title", "center_title".

        Returns:
            PagePlaceholder if found, None otherwise.
        """
        page_data = self.page_placeholders.get(page_type)
        if page_data is None:
            return None
        return page_data.get(role)

    @staticmethod
    def is_icon(shape_location: Location) -> bool:
        """Distinguishing Icons and Images by Area"""
        return shape_location.width * shape_location.height < 3900000 * 3900000

    @property
    def layout_types_names(self) -> list[str]:
        return list(self.layout_types.keys())

    @staticmethod
    def get_text_from_xml(xml: str) -> str:
        root = etree.fromstring(xml)
        nsmap = root.nsmap
        t_elements = root.findall(".//a:t", namespaces=nsmap)
        return "".join([t.text for t in t_elements])

    def reload(self, json_path: str) -> None:
        with open(json_path, encoding="utf-8") as f:
            data = json.load(f)
        self._load_from_dict(data)
        logger.info(f"Reloaded components from {json_path}")


components_manager = ComponentsManager(settings.COMPONENTS_PATH)
