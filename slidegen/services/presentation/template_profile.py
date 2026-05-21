from __future__ import annotations

import re
from dataclasses import dataclass
from enum import Enum
from typing import Any

from pptx.presentation import Presentation
from pptx.slide import Slide

from slidegen.exceptions import PPTTemplateError

READY_THRESHOLD = 0.45
MAX_CANDIDATES_PER_ROLE = 8


class TemplateRole(str, Enum):
    COVER = "cover"
    CATALOG = "catalog"
    CHAPTER = "chapter"
    CONTENT = "content"
    END = "end"


@dataclass(frozen=True)
class TemplateRoleAssignment:
    role: TemplateRole
    slide_index: int
    confidence: float
    reason: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "role": self.role.value,
            "slide_index": self.slide_index,
            "confidence": round(self.confidence, 3),
            "reason": self.reason,
        }


@dataclass(frozen=True)
class TemplateProfile:
    slide_count: int
    assignments: tuple[TemplateRoleAssignment, ...]
    warnings: list[str]
    status: str
    missing_roles: list[str]

    def role_index(self, role: TemplateRole) -> int | None:
        for assignment in self.assignments:
            if assignment.role is role:
                return assignment.slide_index
        return None

    def has_role(self, role: TemplateRole) -> bool:
        return self.role_index(role) is not None

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_count": self.slide_count,
            "status": self.status,
            "assignments": [assignment.to_dict() for assignment in self.assignments],
            "warnings": self.warnings,
            "missing_roles": self.missing_roles,
        }


@dataclass(frozen=True)
class SlideFeatures:
    """Per-slide features extracted for role scoring."""

    index: int
    text: str
    line_count: int
    shape_count: int
    text_shape_count: int
    placeholder_count: int
    title_placeholder_count: int
    subtitle_placeholder_count: int
    body_placeholder_count: int
    has_numbered_lines: bool

    @property
    def normalized_text(self) -> str:
        return re.sub(r"\s+", " ", self.text.casefold()).strip()

    @property
    def text_length(self) -> int:
        return len(self.normalized_text)


KEYWORDS: dict[TemplateRole, tuple[str, ...]] = {
    TemplateRole.COVER: ("cover", "title", "overview", "review", "报告", "汇报", "封面"),
    TemplateRole.CATALOG: ("catalog", "contents", "agenda", "outline", "目录", "大纲", "议程"),
    TemplateRole.CHAPTER: ("chapter", "section", "part", "章节", "篇章", "第一章", "第二章"),
    TemplateRole.CONTENT: (
        "analysis",
        "summary",
        "detail",
        "data",
        "result",
        "market",
        "strategy",
        "growth",
        "分析",
        "策略",
        "增长",
    ),
    TemplateRole.END: ("thank", "thanks", "q&a", "question", "questions", "谢谢", "感谢", "答疑"),
}


def _has_keyword_match(text: str, keywords: tuple[str, ...]) -> bool:
    """Match English keywords with token boundaries and Chinese keywords by substring."""
    for keyword in keywords:
        if re.search(r"[一-鿿]", keyword):
            if keyword in text:
                return True
        else:
            pattern = rf"(?<![a-z0-9_]){re.escape(keyword.casefold())}(?![a-z0-9_])"
            if re.search(pattern, text.casefold()):
                return True
    return False


LEGACY_POSITION_HINTS: dict[TemplateRole, int] = {
    TemplateRole.COVER: 0,
    TemplateRole.CATALOG: 1,
    TemplateRole.CHAPTER: 2,
    TemplateRole.CONTENT: 3,
    TemplateRole.END: 4,
}


def profile_presentation_template(presentation: Presentation) -> TemplateProfile:
    slide_count = len(presentation.slides)
    if slide_count == 0:
        raise PPTTemplateError("PPT template must contain at least one slide")

    features = [_extract_slide_features(index, slide) for index, slide in enumerate(presentation.slides)]
    assignments = _assign_roles(features, slide_count)
    assigned_roles = {assignment.role for assignment in assignments}
    missing_roles = [role.value for role in TemplateRole if role not in assigned_roles]

    warnings: list[str] = [f"{role} role not detected; native fallback will be used" for role in missing_roles]
    low_confidence = [assignment for assignment in assignments if assignment.confidence < READY_THRESHOLD]
    warnings.extend(
        f"{assignment.role.value} role confidence is {assignment.confidence:.2f}; native fallback may be used"
        for assignment in low_confidence
    )

    status = "ready" if not missing_roles and not low_confidence else "review_required"
    return TemplateProfile(
        slide_count=slide_count,
        assignments=tuple(assignments),
        warnings=warnings,
        status=status,
        missing_roles=missing_roles,
    )


def _extract_slide_features(index: int, slide: Slide) -> SlideFeatures:
    texts: list[str] = []
    text_shape_count = 0
    placeholder_count = 0
    title_placeholder_count = 0
    subtitle_placeholder_count = 0
    body_placeholder_count = 0

    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            placeholder_count += 1
            placeholder_type = str(shape.placeholder_format.type).casefold()
            if "subtitle" in placeholder_type:
                subtitle_placeholder_count += 1
            elif "title" in placeholder_type:
                title_placeholder_count += 1
            elif any(name in placeholder_type for name in ("body", "object", "content")):
                body_placeholder_count += 1
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if not text:
            continue
        text_shape_count += 1
        texts.append(text)

    lines = [line.strip() for text in texts for line in text.splitlines() if line.strip()]
    has_numbered_lines = any(re.match(r"^(\d+[.)]|[一二三四五六七八九十]+[、.])", line) for line in lines)
    return SlideFeatures(
        index=index,
        text="\n".join(lines),
        line_count=len(lines),
        shape_count=len(slide.shapes),
        text_shape_count=text_shape_count,
        placeholder_count=placeholder_count,
        title_placeholder_count=title_placeholder_count,
        subtitle_placeholder_count=subtitle_placeholder_count,
        body_placeholder_count=body_placeholder_count,
        has_numbered_lines=has_numbered_lines,
    )


def _assign_roles(features: list[SlideFeatures], slide_count: int) -> list[TemplateRoleAssignment]:
    """Find a bounded maximum-weight role-to-slide assignment without role-order greediness."""
    roles = list(TemplateRole)

    scores: dict[tuple[TemplateRole, int], TemplateRoleAssignment] = {}
    for role in roles:
        for f in features:
            scores[(role, f.index)] = _score_slide_for_role(f, role, slide_count)

    candidates_by_role: dict[TemplateRole, list[TemplateRoleAssignment]] = {
        role: sorted(
            (scores[(role, f.index)] for f in features),
            key=lambda assignment: (assignment.confidence, -assignment.slide_index),
            reverse=True,
        )[:MAX_CANDIDATES_PER_ROLE]
        for role in roles
    }

    best_total = -1.0
    best_selected: list[TemplateRoleAssignment] = []

    def search(
        role_index: int,
        used_slide_indexes: set[int],
        selected: list[TemplateRoleAssignment],
        total: float,
    ) -> None:
        nonlocal best_total, best_selected
        if role_index == len(roles):
            if total > best_total:
                best_total = total
                best_selected = selected.copy()
            return

        role = roles[role_index]
        search(role_index + 1, used_slide_indexes, selected, total)
        for assignment in candidates_by_role[role]:
            if assignment.slide_index in used_slide_indexes:
                continue
            used_slide_indexes.add(assignment.slide_index)
            selected.append(assignment)
            search(role_index + 1, used_slide_indexes, selected, total + assignment.confidence)
            selected.pop()
            used_slide_indexes.remove(assignment.slide_index)

    search(0, set(), [], 0.0)
    return sorted(best_selected, key=lambda a: a.slide_index)


def _score_slide_for_role(
    feature: SlideFeatures,
    role: TemplateRole,
    slide_count: int,
) -> TemplateRoleAssignment:
    text = feature.normalized_text
    score = 0.0
    reasons: list[str] = []

    if _has_keyword_match(text, KEYWORDS[role]):
        score += 0.30
        reasons.append("keyword match")

    if slide_count >= 5 and feature.index == LEGACY_POSITION_HINTS[role]:
        score += 0.20
        reasons.append("legacy role position")

    if role is TemplateRole.COVER:
        if feature.index == 0:
            score += 0.25
            reasons.append("first slide")
        if feature.index == slide_count - 1 and slide_count == 1:
            score += 0.25
            reasons.append("sole slide")
        if feature.title_placeholder_count and feature.subtitle_placeholder_count:
            score += 0.25
            reasons.append("title and subtitle placeholders")
        if feature.text_length <= 220 and feature.line_count <= 4:
            score += 0.15
            reasons.append("short title-like text")
    elif role is TemplateRole.CATALOG:
        if feature.index <= min(2, slide_count - 1):
            score += 0.15
            reasons.append("early slide")
        if feature.has_numbered_lines or feature.line_count >= 3:
            score += 0.25
            reasons.append("list-like structure")
        if _has_keyword_match(text, ("table of contents",)):
            score += 0.25
            reasons.append("table-of-contents phrase")
    elif role is TemplateRole.CHAPTER:
        if 0 < feature.index < slide_count - 1:
            score += 0.10
            reasons.append("middle slide")
        if feature.text_length <= 160 and feature.line_count <= 3:
            score += 0.25
            reasons.append("section-divider density")
        if re.search(r"(?<![a-z0-9_])(chapter|section|part)\s+\d+(?![a-z0-9_])", text) or re.search(
            r"第[一二三四五六七八九十\d]+章", text
        ):
            score += 0.30
            reasons.append("chapter pattern")
    elif role is TemplateRole.CONTENT:
        if 0 < feature.index < slide_count - 1:
            score += 0.15
            reasons.append("non-edge slide")
        if feature.line_count >= 3:
            score += 0.30
            reasons.append("body text lines")
        if feature.text_length >= 80:
            score += 0.20
            reasons.append("body text density")
        if feature.shape_count >= 2 or feature.text_shape_count >= 1:
            score += 0.10
            reasons.append("content shape density")
        if feature.body_placeholder_count or feature.title_placeholder_count:
            score += 0.10
            reasons.append("content placeholder")
    elif role is TemplateRole.END:
        if feature.index == slide_count - 1:
            score += 0.25
            reasons.append("last slide")
        if feature.text_length <= 180:
            score += 0.15
            reasons.append("short closing text")

    return TemplateRoleAssignment(
        role=role,
        slide_index=feature.index,
        confidence=round(min(score, 1.0), 2),
        reason=", ".join(reasons) or "weak heuristic match",
    )
