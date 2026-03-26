from typing import Any

from pptx import Presentation

from .base import BaseDocumentReader, ContentType, DocumentReadResult


class PptxReader(BaseDocumentReader):
    """Reader for PowerPoint files (.pptx)."""

    @classmethod
    def get_supported_content_types(self) -> list[ContentType]:
        return [ContentType.PPTX]

    def convert(self, local_path: str, **kwargs: Any) -> None | DocumentReadResult:
        extension = kwargs.get("file_extension", "").lower()
        supported_extensions = [ct.value for ct in PptxReader.get_supported_content_types()]
        if extension not in supported_extensions:
            return None

        presentation = Presentation(local_path)
        slide_chunks: list[str] = []
        title: str | None = None

        for slide_index, slide in enumerate(presentation.slides, start=1):
            shape_texts: list[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if text:
                    shape_texts.append(text)
                    if title is None:
                        title = text.splitlines()[0].strip()

            if shape_texts:
                slide_chunks.append(f"Slide {slide_index}\n" + "\n".join(shape_texts))

        return DocumentReadResult(
            title=title,
            text_content="\n\n".join(slide_chunks).strip(),
        )
